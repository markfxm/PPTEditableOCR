from __future__ import annotations

import os
import re
import shutil
import subprocess
import sys
import json
import hashlib
import tempfile
import time
import traceback
import importlib
import importlib.util
import types
from dataclasses import dataclass, field, replace
from pathlib import Path
from statistics import median
from typing import Callable

BASE = Path(__file__).resolve().parent.parent
DEPS_DIR = BASE / ".py310deps"
IOPAINT_DEPS_DIR = BASE / ".py310iopaint"
ORIGINAL_SYS_PATH = list(sys.path)
for deps in [IOPAINT_DEPS_DIR, DEPS_DIR]:
    if deps.exists() and str(deps) not in sys.path:
        sys.path.insert(0, str(deps))


def _resolved_sys_path(path: str) -> Path:
    return Path(path or ".").resolve()


def prefer_system_cuda_torch(
    original_sys_path: list[str] | None = None,
    bundled_paths: set[Path] | None = None,
    import_module=importlib.import_module,
    frozen_app: bool | None = None,
) -> bool:
    if "torch" in sys.modules:
        return bool(getattr(getattr(sys.modules["torch"], "cuda", None), "is_available", lambda: False)())

    if frozen_app is None:
        frozen_app = bool(getattr(sys, "frozen", False))
    original_sys_path = list(original_sys_path or ORIGINAL_SYS_PATH)
    bundled_paths = {path.resolve() for path in (bundled_paths or {IOPAINT_DEPS_DIR, DEPS_DIR})}
    current_sys_path = list(sys.path)
    before_modules = set(sys.modules)
    keep_imported_torch = False
    try:
        sys.path[:] = [
            path
            for path in original_sys_path
            if _resolved_sys_path(path) not in bundled_paths
        ]
        torch_module = import_module("torch")
        cuda = getattr(torch_module, "cuda", None)
        is_available = getattr(cuda, "is_available", None)
        if callable(is_available) and is_available():
            import_module("torchvision")
            keep_imported_torch = True
            return True
        if frozen_app:
            keep_imported_torch = True
    except Exception:
        new_torch_modules = [
            name
            for name in list(sys.modules)
            if name not in before_modules and (name == "torch" or name.startswith("torch."))
        ]
        if frozen_app and new_torch_modules:
            keep_imported_torch = True
        pass
    finally:
        sys.path[:] = current_sys_path

    if not keep_imported_torch:
        for name in list(sys.modules):
            if name == "torch" or name.startswith("torch."):
                if name not in before_modules:
                    sys.modules.pop(name, None)
    return False


def prepare_iopaint_torch(progress: ProgressCB = None) -> bool:
    selected = prefer_system_cuda_torch()
    if selected:
        _log(progress, "已优先使用系统 CUDA Torch")
    return selected


import cv2
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from .quality_pipeline import (
    OpenAIImageRepairBackend,
    PageQualityResult,
    QualityMode,
    QualityPipeline,
    QualityStatus,
    build_background_erase_mask,
    decontaminate_asset_rgba,
    expand_text_erase_mask,
    refine_asset_alpha,
)


SKIP_TEXTS = {"NotebookLM"}
ProgressCB = Callable[[str], None] | None
PROGRESS_PREFIX = "__PPTTOEDIT_PROGRESS__|"
PAGE_READY_PREFIX = "__PPTTOEDIT_PAGE_READY__|"
CACHE_VERSION = 5
DEFAULT_PPT_WIDTH = Inches(13.333333)
OCR_BACKEND_LOCAL = "local"
OCR_BACKEND_REMOTE = "remote"
REMOTE_OCR_TOKEN_LENGTH = 40
EXPORT_IMAGE_JPEG_QUALITY = 85
_PADDLEOCR_SYMBOLS = None


def load_paddleocr_symbols():
    global _PADDLEOCR_SYMBOLS
    if _PADDLEOCR_SYMBOLS is not None:
        return _PADDLEOCR_SYMBOLS
    from paddleocr import PaddleOCR

    try:
        from paddleocr import PaddleOCRClient
    except ImportError:
        PaddleOCRClient = None
    try:
        from paddleocr import Model
    except ImportError:
        Model = None
    try:
        from paddleocr import OCROptions
    except ImportError:
        OCROptions = None
    _PADDLEOCR_SYMBOLS = PaddleOCR, PaddleOCRClient, Model, OCROptions
    return _PADDLEOCR_SYMBOLS


def preferred_iopaint_device(device_enum, torch_module=None):
    if torch_module is None:
        try:
            import torch as torch_module
        except Exception:
            torch_module = None
    cuda_available = False
    if torch_module is not None:
        cuda = getattr(torch_module, "cuda", None)
        is_available = getattr(cuda, "is_available", None)
        if callable(is_available):
            try:
                cuda_available = bool(is_available())
            except Exception:
                cuda_available = False
    if cuda_available and hasattr(device_enum, "cuda"):
        return device_enum.cuda
    return device_enum.cpu


def preferred_paddleocr_device(paddle_module=None) -> str:
    if paddle_module is None:
        try:
            import paddle as paddle_module
        except Exception:
            paddle_module = None
    if paddle_module is not None:
        device = getattr(paddle_module, "device", None)
        is_compiled_with_cuda = getattr(device, "is_compiled_with_cuda", None)
        if callable(is_compiled_with_cuda):
            try:
                if is_compiled_with_cuda():
                    return "gpu"
            except Exception:
                pass
    return "cpu"


def bundle_root() -> Path:
    if getattr(sys, "frozen", False) and hasattr(sys, "_MEIPASS"):
        return Path(sys._MEIPASS)
    return BASE


def bundled_models_root() -> Path | None:
    candidate = bundle_root() / "models"
    return candidate if candidate.exists() else None


def bundled_ocr_model_dir(model_name: str) -> Path | None:
    root = bundled_models_root()
    if not root:
        return None
    candidate = root / "paddlex" / "official_models" / model_name
    return candidate if candidate.exists() else None


def bundled_iopaint_model_root() -> Path | None:
    root = bundled_models_root()
    return root if root and (root / "torch" / "hub" / "checkpoints" / "big-lama.pt").exists() else None


def load_iopaint_lama_class():
    module = sys.modules.get("iopaint.model.lama")
    if module is not None and hasattr(module, "LaMa"):
        return module.LaMa

    iopaint_spec = importlib.util.find_spec("iopaint")
    package_locations = list(iopaint_spec.submodule_search_locations or []) if iopaint_spec else []
    if not package_locations:
        raise RuntimeError("未找到 IOPaint 运行依赖。")

    model_dir = Path(package_locations[0]) / "model"
    lama_path = model_dir / "lama.py"
    if not lama_path.exists():
        raise RuntimeError(f"未找到 IOPaint LaMa 模块：{lama_path}")

    if "iopaint.model" not in sys.modules:
        package = types.ModuleType("iopaint.model")
        package.__path__ = [str(model_dir)]  # type: ignore[attr-defined]
        package.__package__ = "iopaint"
        sys.modules["iopaint.model"] = package

    spec = importlib.util.spec_from_file_location("iopaint.model.lama", lama_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法加载 IOPaint LaMa 模块：{lama_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules["iopaint.model.lama"] = module
    spec.loader.exec_module(module)
    return module.LaMa


@dataclass
class OCRBox:
    text: str
    score: float
    bbox: tuple[int, int, int, int]
    erase_rect: tuple[int, int, int, int]
    enabled: bool = True
    manual: bool = False
    edited: bool = False
    rotation: int = 0
    line_height: int | None = None
    text_regions: tuple[tuple[tuple[int, int], ...], ...] = ()
    mask_mode: str = "pending"
    mask_reason: str | None = None

    def set_erase_rect(self, rect: tuple[int, int, int, int]):
        left, top, right, bottom = rect
        self.erase_rect = (
            int(left),
            int(top),
            int(max(left + 1, right)),
            int(max(top + 1, bottom)),
        )

    def set_bbox_from_rect(self, rect: tuple[int, int, int, int]):
        left, top, right, bottom = rect
        self.bbox = (
            int(left),
            int(top),
            int(max(1, right - left)),
            int(max(1, bottom - top)),
        )

    def reset_from_bbox(self, pad_x: int, pad_y: int, image_width: int, image_height: int):
        x, y, w, h = self.bbox
        left = max(0, x - pad_x)
        top = max(0, y - pad_y)
        right = min(image_width - 1, x + w + pad_x)
        bottom = min(image_height - 1, y + h + pad_y)
        self.erase_rect = (left, top, right, bottom)


@dataclass
class VisualAsset:
    asset_id: str
    bbox: tuple[int, int, int, int]
    enabled: bool = True
    source: str = "opencv"
    status: str = "candidate"
    layer: str = "below_text"
    image_path: Path | None = None
    mask_path: Path | None = None
    segmentation_mode: str = "opencv"
    confidence: float | None = None
    confirmed: bool = False
    model_id: str | None = None
    mask_version: int = 0
    segmentation_warning: str | None = None
    sam_selected_index: int | None = None
    sam_points: tuple[tuple[float, float, int], ...] = ()


@dataclass
class PPTSlide:
    index: int
    image_name: str
    image_path: Path
    image_width: int
    image_height: int
    boxes: list[OCRBox] = field(default_factory=list)
    watermark_rect: tuple[int, int, int, int] | None = None
    remove_watermark: bool = True
    ocr_status: str = "pending"
    visual_assets: list[VisualAsset] = field(default_factory=list)

    def reset_boxes(self, pad_x: int, pad_y: int):
        for box in self.boxes:
            box.reset_from_bbox(pad_x, pad_y, self.image_width, self.image_height)


class OCRPageProcessError(RuntimeError):
    def __init__(
        self,
        slide_index: int,
        message: str,
        returncode: int | None = None,
        output: str = "",
    ):
        super().__init__(message)
        self.slide_index = slide_index
        self.returncode = returncode
        self.output = output


@dataclass
class PPTProject:
    source_pptx: Path
    work_dir: Path
    images_dir: Path
    masks_dir: Path
    cleaned_dir: Path
    slides: list[PPTSlide]
    slide_width: int
    slide_height: int
    assets_dir: Path | None = None


def default_cache_path(source_pptx: Path) -> Path:
    return source_pptx.with_name(f"{source_pptx.stem}.ppttoedit.json")


def app_cache_path(source_pptx: Path) -> Path:
    cache_root = Path(os.environ.get("LOCALAPPDATA") or Path.home()) / "PPTEditableOCR" / "ocr_caches"
    digest = hashlib.sha1(str(source_pptx).encode("utf-8", errors="ignore")).hexdigest()[:16]
    return cache_root / f"{source_pptx.stem}-{digest}.ppttoedit.json"


def cache_path_candidates(source_pptx: Path) -> list[Path]:
    sidecar = default_cache_path(source_pptx)
    fallback = app_cache_path(source_pptx)
    return [sidecar] if sidecar == fallback else [sidecar, fallback]


def _rect_to_list(rect: tuple[int, int, int, int]) -> list[int]:
    return [int(value) for value in rect]


def _rect_from_data(value, fallback: tuple[int, int, int, int]) -> tuple[int, int, int, int]:
    if not isinstance(value, (list, tuple)) or len(value) != 4:
        return fallback
    try:
        return tuple(int(item) for item in value)
    except (TypeError, ValueError):
        return fallback


def _optional_int(value) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _regions_to_data(regions: tuple[tuple[tuple[int, int], ...], ...]) -> list[list[list[int]]]:
    return [[[int(x), int(y)] for x, y in region] for region in regions]


def _regions_from_data(value) -> tuple[tuple[tuple[int, int], ...], ...]:
    if not isinstance(value, (list, tuple)):
        return ()
    regions = []
    for region in value:
        if not isinstance(region, (list, tuple)) or len(region) < 3:
            continue
        points = []
        for point in region:
            if not isinstance(point, (list, tuple)) or len(point) < 2:
                points = []
                break
            try:
                points.append((int(point[0]), int(point[1])))
            except (TypeError, ValueError):
                points = []
                break
        if len(points) >= 3:
            regions.append(tuple(points))
    return tuple(regions)


def visual_asset_to_data(asset: VisualAsset) -> dict:
    return {
        "asset_id": asset.asset_id,
        "bbox": _rect_to_list(asset.bbox),
        "enabled": bool(asset.enabled),
        "source": asset.source,
        "status": asset.status,
        "layer": asset.layer,
        "image_path": str(asset.image_path) if asset.image_path else None,
        "mask_path": str(asset.mask_path) if asset.mask_path else None,
        "segmentation_mode": asset.segmentation_mode,
        "confidence": float(asset.confidence) if asset.confidence is not None else None,
        "confirmed": bool(asset.confirmed),
        "model_id": asset.model_id,
        "mask_version": int(asset.mask_version),
        "segmentation_warning": asset.segmentation_warning,
        "sam_selected_index": asset.sam_selected_index,
        "sam_points": [list(point) for point in asset.sam_points],
    }


def visual_asset_from_data(data: dict) -> VisualAsset:
    return VisualAsset(
        asset_id=str(data.get("asset_id") or "visual-asset"),
        bbox=_rect_from_data(data.get("bbox"), (0, 0, 1, 1)),
        enabled=bool(data.get("enabled", True)),
        source=str(data.get("source") or "opencv"),
        status=str(data.get("status") or "candidate"),
        layer=str(data.get("layer") or "below_text"),
        image_path=Path(data["image_path"]) if data.get("image_path") else None,
        mask_path=Path(data["mask_path"]) if data.get("mask_path") else None,
        segmentation_mode=str(data.get("segmentation_mode") or "opencv"),
        confidence=float(data["confidence"]) if data.get("confidence") is not None else None,
        confirmed=bool(data.get("confirmed", False)),
        model_id=str(data["model_id"]) if data.get("model_id") else None,
        mask_version=int(data.get("mask_version", 0)),
        segmentation_warning=str(data["segmentation_warning"]) if data.get("segmentation_warning") else None,
        sam_selected_index=_optional_int(data.get("sam_selected_index")),
        sam_points=tuple(
            (float(point[0]), float(point[1]), int(point[2]))
            for point in data.get("sam_points", [])
            if isinstance(point, (list, tuple)) and len(point) == 3
        ),
    )

def save_project_cache(project: PPTProject, cache_path: Path | None = None, progress: ProgressCB = None) -> Path:
    data = {
        "version": CACHE_VERSION,
        "source_name": project.source_pptx.name,
        "slide_width": int(project.slide_width),
        "slide_height": int(project.slide_height),
        "slides": [
            {
                "index": int(slide.index),
                "image_name": slide.image_name,
                "image_width": int(slide.image_width),
                "image_height": int(slide.image_height),
                "watermark_rect": _rect_to_list(slide.watermark_rect) if slide.watermark_rect else None,
                "remove_watermark": bool(slide.remove_watermark),
                "ocr_status": slide.ocr_status,
                "visual_assets": [visual_asset_to_data(asset) for asset in slide.visual_assets],
                "boxes": [
                    {
                        "text": box.text,
                        "score": float(box.score),
                        "bbox": _rect_to_list(box.bbox),
                        "erase_rect": _rect_to_list(box.erase_rect),
                        "enabled": bool(box.enabled),
                        "manual": bool(box.manual),
                        "edited": bool(box.edited),
                        "rotation": int(box.rotation),
                        "line_height": int(box.line_height) if box.line_height is not None else None,
                        "text_regions": _regions_to_data(box.text_regions),
                        "mask_mode": box.mask_mode,
                        "mask_reason": box.mask_reason,
                    }
                    for box in slide.boxes
                ],
            }
            for slide in project.slides
        ],
    }
    payload = json.dumps(data, ensure_ascii=False, indent=2)
    candidates = [cache_path] if cache_path else cache_path_candidates(project.source_pptx)
    last_error: Exception | None = None
    for candidate in candidates:
        try:
            candidate.parent.mkdir(parents=True, exist_ok=True)
            candidate.write_text(payload, encoding="utf-8")
        except OSError as exc:
            last_error = exc
            continue
        _log(progress, f"识别框缓存已保存：{candidate}")
        return candidate
    raise RuntimeError(f"识别框缓存保存失败：{last_error}")


def load_project_cache(project: PPTProject, cache_path: Path | None = None, progress: ProgressCB = None) -> bool:
    candidates = [cache_path] if cache_path else cache_path_candidates(project.source_pptx)
    cache_path = next((path for path in candidates if path.exists()), None)
    if cache_path is None:
        return False

    try:
        data = json.loads(cache_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        _log(progress, f"识别框缓存读取失败，已忽略：{cache_path} ({exc})")
        return False
    if _optional_int(data.get("version")) not in {2, 3, 4, CACHE_VERSION}:
        _log(progress, f"识别框缓存版本不匹配，已忽略：{cache_path}")
        return False
    cached_slides = data.get("slides")
    if not isinstance(cached_slides, list) or len(cached_slides) != len(project.slides):
        _log(progress, f"识别框缓存页数不匹配，已忽略：{cache_path}")
        return False

    for slide, cached in zip(project.slides, cached_slides):
        if (
            int(cached.get("index", -1)) != slide.index
            or int(cached.get("image_width", -1)) != slide.image_width
            or int(cached.get("image_height", -1)) != slide.image_height
        ):
            _log(progress, f"识别框缓存页面尺寸不匹配，已忽略：{cache_path}")
            return False

    for slide, cached in zip(project.slides, cached_slides):
        slide.remove_watermark = bool(cached.get("remove_watermark", slide.remove_watermark))
        slide.ocr_status = str(cached.get("ocr_status") or ("ok" if cached.get("boxes") else "pending"))
        cached_watermark_rect = (
            _rect_from_data(cached.get("watermark_rect"), slide.watermark_rect)
            if cached.get("watermark_rect")
            else slide.watermark_rect
        )
        default_rect = default_watermark_rect(slide.image_width, slide.image_height)
        if cached_watermark_rect is None:
            cached_watermark_rect = default_rect
        slide.watermark_rect = (
            min(cached_watermark_rect[0], default_rect[0]),
            min(cached_watermark_rect[1], default_rect[1]),
            max(cached_watermark_rect[2], default_rect[2]),
            max(cached_watermark_rect[3], default_rect[3]),
        )
        slide.boxes = []
        slide.visual_assets = [
            visual_asset_from_data(asset)
            for asset in cached.get("visual_assets", [])
            if isinstance(asset, dict)
        ]
        for item in cached.get("boxes", []):
            if not isinstance(item, dict):
                continue
            bbox = _rect_from_data(item.get("bbox"), (0, 0, 1, 1))
            erase_rect = _rect_from_data(item.get("erase_rect"), bbox)
            slide.boxes.append(
                OCRBox(
                    text=str(item.get("text") or ""),
                    score=float(item.get("score", 1.0)),
                    bbox=bbox,
                    erase_rect=erase_rect,
                    enabled=bool(item.get("enabled", True)),
                    manual=bool(item.get("manual", False)),
                    edited=bool(item.get("edited", False)),
                    rotation=int(item.get("rotation", 0)),
                    line_height=_optional_int(item.get("line_height")),
                    text_regions=_regions_from_data(item.get("text_regions")),
                    mask_mode=str(item.get("mask_mode") or "pending"),
                    mask_reason=str(item["mask_reason"]) if item.get("mask_reason") else None,
                )
            )

    _log(progress, f"已加载识别框缓存，跳过 OCR：{cache_path}")
    for slide in project.slides:
        _page_ready(progress, slide)
    return True


def box_to_rect(poly) -> tuple[int, int, int, int]:
    pts = np.array(poly, dtype=np.int32)
    x, y, w, h = cv2.boundingRect(pts)
    return x, y, w, h


def normalize_text_region(poly) -> tuple[tuple[int, int], ...]:
    points = np.asarray(poly, dtype=np.int32).reshape(-1, 2)
    return tuple((int(x), int(y)) for x, y in points)


def default_expand_rect(
    x: int, y: int, w: int, h: int, width: int, height: int
) -> tuple[int, int, int, int]:
    pad_x = max(8, int(w * 0.05))
    pad_y = max(6, int(h * 0.18))
    left = max(0, x - pad_x)
    top = max(0, y - pad_y)
    right = min(width - 1, x + w + pad_x)
    bottom = min(height - 1, y + h + pad_y)
    return left, top, right, bottom


def _log(progress: ProgressCB, message: str):
    if progress:
        progress(message)


def _progress(progress: ProgressCB, percent: int, message: str):
    percent = max(0, min(100, int(percent)))
    _log(progress, f"{PROGRESS_PREFIX}{percent}|{message}")


def _page_ready(progress: ProgressCB, slide: PPTSlide, status: str = "ok"):
    _log(progress, f"{PAGE_READY_PREFIX}{slide.index}|{len(slide.boxes)}|{status}")


def ocr_box_to_data(box: OCRBox) -> dict:
    return {
        "text": box.text,
        "score": float(box.score),
        "bbox": _rect_to_list(box.bbox),
        "erase_rect": _rect_to_list(box.erase_rect),
        "enabled": bool(box.enabled),
        "manual": bool(box.manual),
        "edited": bool(box.edited),
        "rotation": int(box.rotation),
        "line_height": int(box.line_height) if box.line_height is not None else None,
        "text_regions": _regions_to_data(box.text_regions),
        "mask_mode": box.mask_mode,
        "mask_reason": box.mask_reason,
    }


def ocr_box_from_data(data: dict) -> OCRBox:
    bbox = _rect_from_data(data.get("bbox"), (0, 0, 1, 1))
    erase_rect = _rect_from_data(data.get("erase_rect"), bbox)
    return OCRBox(
        text=str(data.get("text") or ""),
        score=float(data.get("score", 1.0)),
        bbox=bbox,
        erase_rect=erase_rect,
        enabled=bool(data.get("enabled", True)),
        manual=bool(data.get("manual", False)),
        edited=bool(data.get("edited", False)),
        rotation=int(data.get("rotation", 0)),
        line_height=_optional_int(data.get("line_height")),
        text_regions=_regions_from_data(data.get("text_regions")),
        mask_mode=str(data.get("mask_mode") or "pending"),
        mask_reason=str(data["mask_reason"]) if data.get("mask_reason") else None,
    )


def ppt_project_to_data(project: PPTProject) -> dict:
    return {
        "source_pptx": str(project.source_pptx),
        "work_dir": str(project.work_dir),
        "images_dir": str(project.images_dir),
        "masks_dir": str(project.masks_dir),
        "cleaned_dir": str(project.cleaned_dir),
        "assets_dir": str(project.assets_dir) if project.assets_dir else None,
        "slide_width": int(project.slide_width),
        "slide_height": int(project.slide_height),
        "slides": [
            {
                "index": int(slide.index),
                "image_name": slide.image_name,
                "image_path": str(slide.image_path),
                "image_width": int(slide.image_width),
                "image_height": int(slide.image_height),
                "watermark_rect": _rect_to_list(slide.watermark_rect) if slide.watermark_rect else None,
                "remove_watermark": bool(slide.remove_watermark),
                "ocr_status": slide.ocr_status,
                "visual_assets": [visual_asset_to_data(asset) for asset in slide.visual_assets],
                "boxes": [ocr_box_to_data(box) for box in slide.boxes],
            }
            for slide in project.slides
        ],
    }


def ppt_project_from_data(data: dict) -> PPTProject:
    slides = []
    for item in data.get("slides", []):
        slides.append(
            PPTSlide(
                index=int(item["index"]),
                image_name=str(item["image_name"]),
                image_path=Path(item["image_path"]),
                image_width=int(item["image_width"]),
                image_height=int(item["image_height"]),
                boxes=[ocr_box_from_data(box) for box in item.get("boxes", [])],
                watermark_rect=_rect_from_data(item.get("watermark_rect"), (0, 0, 0, 0))
                if item.get("watermark_rect")
                else None,
                remove_watermark=bool(item.get("remove_watermark", True)),
                ocr_status=str(item.get("ocr_status") or "pending"),
                visual_assets=[
                    visual_asset_from_data(asset)
                    for asset in item.get("visual_assets", [])
                    if isinstance(asset, dict)
                ],
            )
        )
    return PPTProject(
        source_pptx=Path(data["source_pptx"]),
        work_dir=Path(data["work_dir"]),
        images_dir=Path(data["images_dir"]),
        masks_dir=Path(data["masks_dir"]),
        cleaned_dir=Path(data["cleaned_dir"]),
        assets_dir=Path(data["assets_dir"]) if data.get("assets_dir") else None,
        slides=slides,
        slide_width=int(data["slide_width"]),
        slide_height=int(data["slide_height"]),
    )


def _is_horizontal_ocr_box(box: OCRBox) -> bool:
    if box.rotation in {90, 270, -90}:
        return False
    x, y, width, height = box.bbox
    return not (height > width * 1.45 and len(box.text.strip()) > 1)


def _horizontal_gap(first: OCRBox, second: OCRBox) -> int:
    first_left, _first_top, first_width, _first_height = first.bbox
    second_left, _second_top, second_width, _second_height = second.bbox
    first_right = first_left + first_width
    second_right = second_left + second_width
    if first_right < second_left:
        return second_left - first_right
    if second_right < first_left:
        return first_left - second_right
    return 0


def _can_join_ocr_lines(first: OCRBox, second: OCRBox) -> bool:
    if not _is_horizontal_ocr_box(first) or not _is_horizontal_ocr_box(second):
        return False

    first_x, first_y, first_width, first_height = first.bbox
    second_x, second_y, second_width, second_height = second.bbox
    first_height_ref = first.line_height or first_height
    second_height_ref = second.line_height or second_height
    min_height = min(first_height_ref, second_height_ref)
    max_height = max(first_height_ref, second_height_ref)
    if min_height <= 0 or max_height / min_height > 1.35:
        return False

    vertical_gap = second_y - (first_y + first_height)
    if vertical_gap < -0.25 * min_height or vertical_gap > max(8, 0.9 * max_height):
        return False

    first_center = first_x + first_width / 2
    second_center = second_x + second_width / 2
    first_right = first_x + first_width
    second_right = second_x + second_width
    alignment_delta = min(
        abs(first_x - second_x),
        abs(first_center - second_center),
        abs(first_right - second_right),
    )
    if alignment_delta > 0.75 * max_height:
        return False
    return _horizontal_gap(first, second) <= 1.5 * max_height


def _ocr_line_join_score(first: OCRBox, second: OCRBox) -> float:
    if not _can_join_ocr_lines(first, second):
        return float("inf")
    first_x, first_y, first_width, first_height = first.bbox
    second_x, second_y, second_width, _second_height = second.bbox
    first_height_ref = first.line_height or first_height
    second_height_ref = second.line_height or second.bbox[3]
    first_center = first_x + first_width / 2
    second_center = second_x + second_width / 2
    first_right = first_x + first_width
    second_right = second_x + second_width
    alignment_delta = min(
        abs(first_x - second_x),
        abs(first_center - second_center),
        abs(first_right - second_right),
    )
    vertical_gap = second_y - (first_y + first_height)
    return abs(vertical_gap) + alignment_delta + abs(first_height_ref - second_height_ref)


def _merge_ocr_line_group(group: list[OCRBox]) -> OCRBox:
    if len(group) == 1:
        return group[0]

    left = min(box.bbox[0] for box in group)
    top = min(box.bbox[1] for box in group)
    right = max(box.bbox[0] + box.bbox[2] for box in group)
    bottom = max(box.bbox[1] + box.bbox[3] for box in group)
    erase_left = min(box.erase_rect[0] for box in group)
    erase_top = min(box.erase_rect[1] for box in group)
    erase_right = max(box.erase_rect[2] for box in group)
    erase_bottom = max(box.erase_rect[3] for box in group)
    line_heights = [box.line_height or box.bbox[3] for box in group]
    return OCRBox(
        text="\n".join(box.text for box in group),
        score=sum(box.score for box in group) / len(group),
        bbox=(left, top, right - left, bottom - top),
        erase_rect=(erase_left, erase_top, erase_right, erase_bottom),
        enabled=all(box.enabled for box in group),
        manual=any(box.manual for box in group),
        edited=any(box.edited for box in group),
        line_height=int(round(median(line_heights))),
        text_regions=tuple(region for box in group for region in box.text_regions),
    )


def group_ocr_boxes(boxes: list[OCRBox]) -> list[OCRBox]:
    candidates = [
        box
        for box in boxes
        if box.text.strip() and box.text.strip() not in SKIP_TEXTS
    ]
    candidates.sort(key=lambda box: (box.bbox[1], box.bbox[0]))
    groups: list[list[OCRBox]] = []
    for box in candidates:
        joinable = [
            (index, _ocr_line_join_score(group[-1], box))
            for index, group in enumerate(groups)
            if _can_join_ocr_lines(group[-1], box)
        ]
        if joinable:
            group_index, _score = min(joinable, key=lambda item: item[1])
            groups[group_index].append(box)
        else:
            groups.append([box])
    return [_merge_ocr_line_group(group) for group in groups]


def build_ocr_boxes(page: dict, slide: PPTSlide) -> list[OCRBox]:
    boxes: list[OCRBox] = []
    for poly, text, score in zip(page["dt_polys"], page["rec_texts"], page["rec_scores"]):
        text = (text or "").strip()
        if not text:
            continue
        if text in SKIP_TEXTS:
            continue
        x, y, w, h = box_to_rect(poly)
        erase_rect = default_expand_rect(x, y, w, h, slide.image_width, slide.image_height)
        boxes.append(
            OCRBox(
                text=text,
                score=float(score),
                bbox=(x, y, w, h),
                erase_rect=erase_rect,
                line_height=h,
                text_regions=(normalize_text_region(poly),),
            )
        )
    return group_ocr_boxes(boxes)


def unique_output_path(path: Path) -> Path:
    if not path.exists():
        return path
    for index in range(2, 1000):
        candidate = path.with_name(f"{path.stem}-{index}{path.suffix}")
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"无法生成不覆盖已有文件的输出路径：{path}")


def default_pdf_pptx_path(source_pdf: Path) -> Path:
    return unique_output_path(source_pdf.with_name(f"{source_pdf.stem}-from-pdf.pptx"))


def fit_rect_to_slide(image_width: int, image_height: int, slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    scale = min(slide_width / image_width, slide_height / image_height)
    width = int(image_width * scale)
    height = int(image_height * scale)
    left = int((slide_width - width) / 2)
    top = int((slide_height - height) / 2)
    return left, top, width, height


def default_watermark_rect(width: int, height: int) -> tuple[int, int, int, int]:
    watermark_width = max(260, int(width * 0.16))
    watermark_height = max(120, int(height * 0.10))
    return (
        max(0, width - watermark_width),
        max(0, height - watermark_height),
        width - 1,
        height - 1,
    )


def convert_pdf_to_pptx(source_pdf: Path, output_pptx: Path | None = None, progress: ProgressCB = None) -> Path:
    import pypdfium2

    source_pdf = source_pdf.expanduser().resolve()
    if output_pptx is None:
        output_pptx = default_pdf_pptx_path(source_pdf)
    else:
        output_pptx = Path(output_pptx).expanduser().resolve()

    _log(progress, f"开始转换 PDF：{source_pdf}")
    pdf = pypdfium2.PdfDocument(str(source_pdf))
    page_count = len(pdf)
    if page_count <= 0:
        raise RuntimeError(f"PDF 没有可转换页面：{source_pdf}")

    first_page = pdf[0]
    first_width, first_height = first_page.get_size()
    first_page.close()
    if first_width <= 0 or first_height <= 0:
        raise RuntimeError(f"PDF 首页尺寸无效：{source_pdf}")

    out = Presentation()
    out.slide_width = int(DEFAULT_PPT_WIDTH)
    out.slide_height = int(DEFAULT_PPT_WIDTH * first_height / first_width)
    while out.slides:
        rel_id = out.slides._sldIdLst[0].rId
        out.part.drop_rel(rel_id)
        del out.slides._sldIdLst[0]

    blank = out.slide_layouts[6]
    _progress(progress, 0, f"PDF 转 PPT：0/{page_count}")
    with tempfile.TemporaryDirectory(prefix="ppttoedit_pdf_") as temp_dir:
        temp_root = Path(temp_dir)
        for page_index in range(page_count):
            page = pdf[page_index]
            image_path = temp_root / f"page_{page_index + 1:04d}.png"
            bitmap = page.render(scale=2)
            image = bitmap.to_pil().convert("RGB")
            image_width, image_height = image.size
            image.save(image_path)
            bitmap.close()
            page.close()

            slide = out.slides.add_slide(blank)
            left, top, width, height = fit_rect_to_slide(
                image_width,
                image_height,
                out.slide_width,
                out.slide_height,
            )
            slide.shapes.add_picture(
                str(image_path),
                left,
                top,
                width=width,
                height=height,
            )
            done = page_index + 1
            _progress(progress, int(done * 100 / page_count), f"PDF 转 PPT：{done}/{page_count}")
            _log(progress, f"第 {done} 页已加入 PPT")

        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        out.save(str(output_pptx))

    pdf.close()
    _log(progress, f"PDF 转 PPT 完成：{output_pptx}")
    return output_pptx


def extract_slide_images(source_pptx: Path, images_dir: Path, progress: ProgressCB = None):
    images_dir.mkdir(parents=True, exist_ok=True)
    src = Presentation(str(source_pptx))
    extracted: list[PPTSlide] = []

    for index, slide in enumerate(src.slides, start=1):
        pic = next((shape for shape in slide.shapes if getattr(shape, "image", None)), None)
        if pic is None:
            continue
        image_name = f"slide_{index:02d}.png"
        out_path = images_dir / image_name
        out_path.write_bytes(pic.image.blob)
        with Image.open(out_path) as image:
            width, height = image.size
        extracted.append(
            PPTSlide(
                index=index,
                image_name=image_name,
                image_path=out_path,
                image_width=width,
                image_height=height,
                watermark_rect=default_watermark_rect(width, height),
            )
        )
        _log(progress, f"提取第 {index} 页图片")
    return src, extracted


def prepare_pdf_project(
    source_pdf: Path,
    work_dir: Path | None = None,
    progress: ProgressCB = None,
    auto_ocr: bool = False,
    ocr_backend: str = OCR_BACKEND_LOCAL,
    ocr_token: str | None = None,
) -> PPTProject:
    import pypdfium2

    source_pdf = source_pdf.expanduser().resolve()
    if work_dir is None:
        work_dir = BASE / "_gui_workspace" / source_pdf.stem
    if work_dir.exists():
        shutil.rmtree(work_dir)
    images_dir = work_dir / "images"
    masks_dir = work_dir / "masks"
    cleaned_dir = work_dir / "cleaned"
    images_dir.mkdir(parents=True, exist_ok=True)
    work_dir.mkdir(parents=True, exist_ok=True)

    _log(progress, f"开始加载 PDF：{source_pdf}")
    pdf = pypdfium2.PdfDocument(str(source_pdf))
    try:
        page_count = len(pdf)
        if page_count <= 0:
            raise RuntimeError(f"PDF 没有可转换页面：{source_pdf}")

        first_page = pdf[0]
        first_width, first_height = first_page.get_size()
        first_page.close()
        if first_width <= 0 or first_height <= 0:
            raise RuntimeError(f"PDF 首页尺寸无效：{source_pdf}")

        slide_width = int(DEFAULT_PPT_WIDTH)
        slide_height = int(DEFAULT_PPT_WIDTH * first_height / first_width)
        slides: list[PPTSlide] = []
        _progress(progress, 0, f"PDF 页面载入：0/{page_count}")
        for page_index in range(page_count):
            page = pdf[page_index]
            image_name = f"slide_{page_index + 1:02d}.png"
            image_path = images_dir / image_name
            bitmap = page.render(scale=2)
            image = bitmap.to_pil().convert("RGB")
            image_width, image_height = image.size
            image.save(image_path)
            bitmap.close()
            page.close()
            slides.append(
                PPTSlide(
                    index=page_index + 1,
                    image_name=image_name,
                    image_path=image_path,
                    image_width=image_width,
                    image_height=image_height,
                    watermark_rect=default_watermark_rect(image_width, image_height),
                )
            )
            done = page_index + 1
            _progress(progress, int(done * 100 / page_count), f"PDF 页面载入：{done}/{page_count}")
            _log(progress, f"第 {done} 页 PDF 图片已载入")
    finally:
        pdf.close()

    project = PPTProject(
        source_pptx=source_pdf,
        work_dir=work_dir,
        images_dir=images_dir,
        masks_dir=masks_dir,
        cleaned_dir=cleaned_dir,
        slides=slides,
        slide_width=slide_width,
        slide_height=slide_height,
    )
    if not load_project_cache(project, progress=progress) and auto_ocr:
        run_ocr(slides, progress, ocr_backend=ocr_backend, ocr_token=ocr_token)
    _log(progress, "PDF 已载入工作区，未生成中间 PPT")
    return project

def _make_remote_ocr_options(OCROptions):
    if OCROptions is None:
        return None
    try:
        return OCROptions(
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            visualize=False,
        )
    except TypeError:
        try:
            return OCROptions(
                use_doc_orientation_classify=False,
                use_doc_unwarping=False,
                use_textline_orientation=False,
            )
        except TypeError:
            return None


def _result_attr(result, *names):
    for name in names:
        if isinstance(result, dict) and name in result:
            return result[name]
        if hasattr(result, name):
            return getattr(result, name)
    return None


def _result_as_page_dict(page) -> dict:
    pruned = _result_attr(page, "prunedResult", "pruned_result", "res")
    if pruned is None and hasattr(page, "to_dict"):
        try:
            pruned = _result_attr(page.to_dict(), "prunedResult", "pruned_result", "res")
        except Exception:
            pruned = None
    if pruned is None:
        pruned = page
    if hasattr(pruned, "to_dict"):
        pruned = pruned.to_dict()
    if not isinstance(pruned, dict):
        raise RuntimeError(f"远端 OCR 返回格式不支持：{type(pruned).__name__}")
    return {
        "dt_polys": pruned.get("dt_polys") or pruned.get("dtPolys") or [],
        "rec_texts": pruned.get("rec_texts") or pruned.get("recTexts") or [],
        "rec_scores": pruned.get("rec_scores") or pruned.get("recScores") or [],
    }


def _predict_remote_ocr_page(client, image_path: Path, Model=None, OCROptions=None) -> dict:
    model = getattr(Model, "PP_OCRV5", "PP-OCRv5") if Model is not None else "PP-OCRv5"
    kwargs = {
        "file_path": str(image_path),
        "model": model,
    }
    options = _make_remote_ocr_options(OCROptions)
    if options is not None:
        kwargs["options"] = options
    try:
        result = client.ocr(**kwargs)
    except TypeError:
        kwargs.pop("options", None)
        result = client.ocr(**kwargs)
    pages = _result_attr(result, "pages")
    if not pages:
        return {"dt_polys": [], "rec_texts": [], "rec_scores": []}
    return _result_as_page_dict(pages[0])


def predict_ocr_page(
    slide: PPTSlide,
    ocr_backend: str = OCR_BACKEND_LOCAL,
    ocr_token: str | None = None,
) -> list[OCRBox]:
    PaddleOCR, PaddleOCRClient, Model, OCROptions = load_paddleocr_symbols()
    det_dir = bundled_ocr_model_dir("PP-OCRv5_server_det")
    rec_dir = bundled_ocr_model_dir("PP-OCRv5_server_rec")
    use_remote = ocr_backend == OCR_BACKEND_REMOTE
    if use_remote:
        token = (ocr_token or os.environ.get("PADDLEOCR_ACCESS_TOKEN") or "").strip()
        if len(token) != REMOTE_OCR_TOKEN_LENGTH:
            raise RuntimeError("远端 OCR 需要 40 位访问令牌。")
        if PaddleOCRClient is None:
            raise RuntimeError("当前 PaddleOCR 版本不支持远端调用，请升级 paddleocr 后重试。")
        client = PaddleOCRClient(token=token)
        try:
            page = _predict_remote_ocr_page(client, slide.image_path, Model, OCROptions)
        finally:
            if hasattr(client, "close"):
                client.close()
        return build_ocr_boxes(page, slide)

    ocr = PaddleOCR(
        lang="ch",
        device=preferred_paddleocr_device(),
        use_doc_orientation_classify=False,
        use_doc_unwarping=False,
        use_textline_orientation=False,
        text_detection_model_dir=str(det_dir) if det_dir else None,
        text_recognition_model_dir=str(rec_dir) if rec_dir else None,
    )
    page = ocr.predict(str(slide.image_path))[0]
    return build_ocr_boxes(page, slide)


def run_ocr_page_subprocess(
    slide: PPTSlide,
    progress: ProgressCB = None,
    ocr_backend: str = OCR_BACKEND_LOCAL,
    ocr_token: str | None = None,
    runner=subprocess.run,
    timeout: int | None = None,
) -> list[OCRBox]:
    _log(progress, f"开始第 {slide.index} 页 OCR")
    with tempfile.TemporaryDirectory(prefix="ppttoedit_ocr_page_") as temp_dir:
        temp_root = Path(temp_dir)
        input_path = temp_root / "input.json"
        output_path = temp_root / "output.json"
        input_path.write_text(
            json.dumps(
                {
                    "slide": {
                        "index": slide.index,
                        "image_name": slide.image_name,
                        "image_path": str(slide.image_path),
                        "image_width": slide.image_width,
                        "image_height": slide.image_height,
                    },
                    "ocr_backend": ocr_backend,
                    "ocr_token": ocr_token,
                },
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        env = os.environ.copy()
        env.setdefault("PYTHONUTF8", "1")
        env.setdefault("PYTHONIOENCODING", "utf-8")
        env.setdefault("PYTHONFAULTHANDLER", "1")
        try:
            if getattr(sys, "frozen", False):
                command = [sys.executable, "--ocr-page-worker", str(input_path), str(output_path)]
            else:
                command = [
                    sys.executable,
                    "-X",
                    "faulthandler",
                    "-m",
                    "app.ocr_page_worker",
                    str(input_path),
                    str(output_path),
                ]
            result = runner(
                command,
                cwd=str(BASE),
                env=env,
                text=True,
                encoding="utf-8",
                errors="replace",
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=timeout,
            )
        except subprocess.TimeoutExpired as exc:
            output = "\n".join(part for part in [exc.stdout or "", exc.stderr or ""] if part)
            raise OCRPageProcessError(slide.index, f"第 {slide.index} 页 OCR 超时。", None, output) from exc

        output = "\n".join(part for part in [result.stdout, result.stderr] if part)
        if result.returncode != 0:
            raise OCRPageProcessError(
                slide.index,
                f"第 {slide.index} 页 OCR 子进程异常退出（代码 {result.returncode}）。",
                result.returncode,
                output,
            )
        try:
            payload = json.loads(output_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError) as exc:
            raise OCRPageProcessError(slide.index, f"第 {slide.index} 页 OCR 结果读取失败。", result.returncode, output) from exc
    return [ocr_box_from_data(item) for item in payload.get("boxes", [])]


def run_export_editable_ppt_subprocess(
    project: PPTProject,
    output_pptx: Path,
    progress: ProgressCB = None,
    enhance_images: bool = True,
    quality_mode: QualityMode = QualityMode.LOCAL_FAST,
    online_pages: set[int] | None = None,
    accepted_local_pages: set[int] | None = None,
    openai_api_key: str | None = None,
    popen=subprocess.Popen,
    timeout: int | None = None,
) -> Path:
    _log(progress, "启动独立导出进程")
    with tempfile.TemporaryDirectory(prefix="ppttoedit_export_") as temp_dir:
        temp_root = Path(temp_dir)
        input_path = temp_root / "input.json"
        payload = {
            "project": ppt_project_to_data(project),
            "output_pptx": str(output_pptx),
            "enhance_images": bool(enhance_images),
            "quality_mode": QualityMode(quality_mode).value,
            "online_pages": sorted(set(online_pages or ())),
            "accepted_local_pages": sorted(set(accepted_local_pages or ())),
        }
        input_path.write_text(json.dumps(payload, ensure_ascii=False), encoding="utf-8")

        env = os.environ.copy()
        env["PYTHONUTF8"] = "1"
        env["PYTHONIOENCODING"] = "utf-8"
        env["PYTHONFAULTHANDLER"] = "1"
        env["PYTHONUNBUFFERED"] = "1"
        if openai_api_key:
            env["PPTTOEDIT_OPENAI_API_KEY"] = openai_api_key
        if getattr(sys, "frozen", False):
            command = [sys.executable, "--export-worker", str(input_path)]
        else:
            command = [
                sys.executable,
                "-X",
                "faulthandler",
                "-m",
                "app.export_worker",
                str(input_path),
            ]

        process = popen(
            command,
            cwd=str(BASE),
            env=env,
            text=True,
            encoding="utf-8",
            errors="replace",
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
        )
        stdout_lines: list[str] = []
        if process.stdout is not None:
            for line in process.stdout:
                message = line.rstrip("\r\n")
                stdout_lines.append(message)
                if message:
                    _log(progress, message)
        try:
            returncode = process.wait(timeout=timeout)
        except subprocess.TimeoutExpired as exc:
            process.kill()
            raise RuntimeError("导出子进程超时。") from exc

        output = "\n".join(stdout_lines)
        if returncode != 0:
            raise RuntimeError(f"导出子进程异常退出（代码 {returncode}）。\n{output}")

    return output_pptx


def run_ocr(
    slides: list[PPTSlide],
    progress: ProgressCB = None,
    ocr_backend: str = OCR_BACKEND_LOCAL,
    ocr_token: str | None = None,
):
    PaddleOCR, PaddleOCRClient, Model, OCROptions = load_paddleocr_symbols()
    det_dir = bundled_ocr_model_dir("PP-OCRv5_server_det")
    rec_dir = bundled_ocr_model_dir("PP-OCRv5_server_rec")
    use_remote = ocr_backend == OCR_BACKEND_REMOTE

    def create_ocr():
        if use_remote:
            token = (ocr_token or os.environ.get("PADDLEOCR_ACCESS_TOKEN") or "").strip()
            if len(token) != REMOTE_OCR_TOKEN_LENGTH:
                raise RuntimeError("远端 OCR 需要 40 位访问令牌。")
            if PaddleOCRClient is None:
                raise RuntimeError("当前 PaddleOCR 版本不支持远端调用，请升级 paddleocr 后重试。")
            return PaddleOCRClient(token=token)
        ocr_device = preferred_paddleocr_device()
        return PaddleOCR(
            lang="ch",
            device=ocr_device,
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            text_detection_model_dir=str(det_dir) if det_dir else None,
            text_recognition_model_dir=str(rec_dir) if rec_dir else None,
        )

    if use_remote:
        _log(progress, "使用远端 PaddleOCR 识别")
    else:
        _log(progress, f"使用本地 PaddleOCR 识别（设备：{preferred_paddleocr_device()}）")
    ocr = create_ocr()
    for slide in slides:
        try:
            if use_remote:
                page = _predict_remote_ocr_page(ocr, slide.image_path, Model, OCROptions)
            else:
                page = ocr.predict(str(slide.image_path))[0]
        except Exception as exc:
            _log(progress, f"第 {slide.index} 页 OCR 失败，已跳过该页，可稍后手动新增框：{exc}")
            _log(progress, traceback.format_exc())
            slide.boxes = []
            slide.ocr_status = "failed"
            _page_ready(progress, slide, status="failed")
            try:
                if use_remote and hasattr(ocr, "close"):
                    ocr.close()
                ocr = create_ocr()
            except Exception as reset_exc:
                _log(progress, f"OCR 引擎重置失败，后续页面可能继续失败：{reset_exc}")
            continue
        boxes = build_ocr_boxes(page, slide)
        slide.boxes = boxes
        slide.ocr_status = "ok"
        _page_ready(progress, slide)
        _log(progress, f"第 {slide.index} 页 OCR 完成，共 {len(boxes)} 个框")
    if use_remote and hasattr(ocr, "close"):
        ocr.close()


def prepare_project(
    source_pptx: Path,
    work_dir: Path | None = None,
    progress: ProgressCB = None,
    ocr_backend: str = OCR_BACKEND_LOCAL,
    ocr_token: str | None = None,
    auto_ocr: bool = True,
) -> PPTProject:
    source_pptx = source_pptx.expanduser().resolve()
    if work_dir is None:
        work_dir = BASE / "_gui_workspace" / source_pptx.stem
    if work_dir.exists():
        shutil.rmtree(work_dir)
    images_dir = work_dir / "images"
    masks_dir = work_dir / "masks"
    cleaned_dir = work_dir / "cleaned"
    assets_dir = work_dir / "assets"
    work_dir.mkdir(parents=True, exist_ok=True)

    src, slides = extract_slide_images(source_pptx, images_dir, progress)
    project = PPTProject(
        source_pptx=source_pptx,
        work_dir=work_dir,
        images_dir=images_dir,
        masks_dir=masks_dir,
        cleaned_dir=cleaned_dir,
        assets_dir=assets_dir,
        slides=slides,
        slide_width=src.slide_width,
        slide_height=src.slide_height,
    )
    if not load_project_cache(project, progress=progress) and auto_ocr:
        run_ocr(slides, progress, ocr_backend=ocr_backend, ocr_token=ocr_token)
    for slide in project.slides:
        if slide.visual_assets or not slide.image_path.exists():
            continue
        with Image.open(slide.image_path) as source_image:
            slide.visual_assets = detect_visual_assets(np.asarray(source_image.convert("RGB")), slide)
        if slide.visual_assets:
            _log(progress, f"第 {slide.index} 页发现 {len(slide.visual_assets)} 个图片区候选")
    return project

def _rect_mask(image_width: int, image_height: int, rect: tuple[int, int, int, int]) -> np.ndarray:
    mask = np.zeros((image_height, image_width), dtype=np.uint8)
    left, top, right, bottom = rect
    cv2.rectangle(mask, (left, top), (right, bottom), 255, -1)
    return mask


def text_region_mask(image_shape: tuple[int, ...], boxes: list[OCRBox]) -> np.ndarray:
    """Return the real OCR glyph polygons, falling back only for legacy boxes."""
    height, width = image_shape[:2]
    mask = np.zeros((height, width), dtype=np.uint8)
    for box in boxes:
        if not box.enabled:
            continue
        if box.text_regions:
            for region in box.text_regions:
                cv2.fillPoly(mask, [np.asarray(region, dtype=np.int32)], 255)
        else:
            x, y, w, h = box.bbox
            cv2.rectangle(mask, (x, y), (x + w, y + h), 255, -1)
    return mask


def _visual_foreground_mask(image: np.ndarray) -> np.ndarray:
    """Conservative non-background mask used when an optional SAM model is unavailable."""
    hsv = cv2.cvtColor(image, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]
    value = hsv[:, :, 2]
    mask = np.zeros(image.shape[:2], dtype=np.uint8)
    # Colorful diagrams/photos are reliable candidates. Low-saturation page grids remain background.
    mask[(saturation >= 35) & (value <= 245)] = 255
    return cv2.morphologyEx(mask, cv2.MORPH_CLOSE, np.ones((5, 5), dtype=np.uint8))


def detect_visual_assets(image: np.ndarray, slide: PPTSlide) -> list[VisualAsset]:
    """Find conservative visual candidates after excluding actual text polygons."""
    candidates = _visual_foreground_mask(image)
    candidates[text_region_mask(image.shape, slide.boxes) > 0] = 0
    contours, _hierarchy = cv2.findContours(candidates, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    min_area = max(256, int(image.shape[0] * image.shape[1] * 0.01))
    max_area = int(image.shape[0] * image.shape[1] * 0.80)
    assets: list[VisualAsset] = []
    for contour in sorted(contours, key=cv2.contourArea, reverse=True):
        x, y, width, height = cv2.boundingRect(contour)
        area = width * height
        if area < min_area or area > max_area:
            continue
        assets.append(
            VisualAsset(
                asset_id=f"slide-{slide.index}-visual-{len(assets) + 1}",
                bbox=(int(x), int(y), int(width), int(height)),
                source="opencv",
                status="candidate",
                layer="below_text",
            )
        )
    return assets


def visual_asset_alpha_mask(image: np.ndarray, asset: VisualAsset, boxes: list[OCRBox]) -> np.ndarray:
    """Return a full-page alpha mask; overlapping text is repaired in the asset pixels later."""
    height, width = image.shape[:2]
    x, y, asset_width, asset_height = asset.bbox
    left, top = max(0, x), max(0, y)
    right, bottom = min(width, x + asset_width), min(height, y + asset_height)
    alpha = np.zeros((height, width), dtype=np.uint8)
    if right <= left or bottom <= top:
        return alpha
    if asset.segmentation_mode == "sam2":
        if asset.segmentation_warning or not asset.mask_path or not asset.mask_path.is_file():
            return alpha
        with Image.open(asset.mask_path) as stored_mask:
            cropped = np.asarray(stored_mask.convert("L"))
        expected_shape = (bottom - top, right - left)
        if cropped.shape != expected_shape:
            return alpha
        alpha[top:bottom, left:right] = cropped
    if not np.any(alpha[top:bottom, left:right]):
        if asset.segmentation_mode == "sam2":
            return alpha
        foreground = _visual_foreground_mask(image)
        alpha[top:bottom, left:right] = foreground[top:bottom, left:right]
    if not np.any(alpha[top:bottom, left:right]):
        alpha[top:bottom, left:right] = 255
    if asset.layer != "above_text":
        overlap_text = text_region_mask(image.shape, boxes)
        overlap_text[:top, :] = 0
        overlap_text[bottom:, :] = 0
        overlap_text[:, :left] = 0
        overlap_text[:, right:] = 0
        alpha[overlap_text > 0] = 255
    return alpha


def build_asset_text_mask(image: np.ndarray, asset: VisualAsset, boxes: list[OCRBox]) -> np.ndarray:
    """Build an expanded text-stroke mask clipped to one visual asset."""
    height, width = image.shape[:2]
    x, y, asset_width, asset_height = asset.bbox
    left, top = max(0, x), max(0, y)
    right, bottom = min(width, x + asset_width), min(height, y + asset_height)
    combined = np.zeros((height, width), dtype=np.uint8)
    if right <= left or bottom <= top or asset.layer == "above_text":
        return combined

    for box in boxes:
        if not box.enabled:
            continue
        box_x, box_y, box_width, box_height = box.bbox
        if box_x + box_width <= left or box_x >= right or box_y + box_height <= top or box_y >= bottom:
            continue
        strokes, _mode, _reason = build_text_stroke_mask(image, box)
        line_height = max(1, int(box.line_height or box_height))
        radius = max(1, int(round(line_height * 0.09)))
        kernel_size = radius * 2 + 1
        expanded = cv2.dilate(
            strokes,
            np.ones((kernel_size, kernel_size), dtype=np.uint8),
            iterations=1,
        )
        combined = cv2.bitwise_or(combined, expanded)

    allowed = np.zeros_like(combined)
    allowed[top:bottom, left:right] = 255
    return cv2.bitwise_and(combined, allowed)


def soft_blend_inpaint(
    original: np.ndarray,
    repaired: np.ndarray,
    mask: np.ndarray,
    feather_px: int = 24,
) -> np.ndarray:
    """Keep repaired pixels inside the mask and feather only into the surrounding edge."""
    if original.shape != repaired.shape or original.shape[:2] != mask.shape[:2]:
        raise ValueError("修复图、原图和蒙版尺寸必须一致。")
    binary = np.zeros(mask.shape[:2], dtype=np.uint8)
    binary[np.asarray(mask) > 0] = 255
    if not np.any(binary):
        return original.copy()
    feather_px = max(1, int(feather_px))
    kernel_size = feather_px * 2 + 1
    soft = cv2.GaussianBlur(binary, (kernel_size, kernel_size), sigmaX=max(1.0, feather_px / 3))
    # Never blend the source content back into the area that the model repaired.
    soft[binary > 0] = 255
    weight = soft.astype(np.float32)[:, :, None] / 255.0
    blended = repaired.astype(np.float32) * weight + original.astype(np.float32) * (1.0 - weight)
    return np.clip(np.rint(blended), 0, 255).astype(np.uint8)


def finish_page_inpaint(original: np.ndarray, generated: np.ndarray, mask: np.ndarray) -> np.ndarray:
    """Prefer deterministic texture continuation for large flat regions, then soften all seams."""
    binary = np.zeros(mask.shape[:2], dtype=np.uint8)
    binary[np.asarray(mask) > 0] = 255
    if not np.any(binary):
        return original.copy()
    candidate = generated.copy()
    count, labels, stats, _centroids = cv2.connectedComponentsWithStats(binary, connectivity=8)
    page_area = max(1, binary.shape[0] * binary.shape[1])
    ring_kernel = np.ones((15, 15), dtype=np.uint8)
    for component_index in range(1, count):
        area = int(stats[component_index, cv2.CC_STAT_AREA])
        if area < page_area * 0.02:
            continue
        component = np.zeros_like(binary)
        component[labels == component_index] = 255
        ring = cv2.subtract(cv2.dilate(component, ring_kernel, iterations=1), component)
        samples = original[ring > 0]
        if samples.size == 0 or float(np.mean(np.std(samples.astype(np.float32), axis=0))) >= 18.0:
            continue
        texture_fill = cv2.inpaint(original, component, 7, cv2.INPAINT_TELEA)
        candidate[component > 0] = texture_fill[component > 0]

    masked_width = int(stats[1:, cv2.CC_STAT_WIDTH].max()) if count > 1 else 1
    masked_height = int(stats[1:, cv2.CC_STAT_HEIGHT].max()) if count > 1 else 1
    feather = min(40, max(15, int(round(min(masked_width, masked_height) * 0.08))))
    return soft_blend_inpaint(original, candidate, binary, feather_px=feather)


def _clean_segmentation_mask(mask: np.ndarray, asset: VisualAsset, boxes: list[OCRBox]) -> np.ndarray:
    binary = np.zeros(mask.shape[:2], dtype=np.uint8)
    binary[np.asarray(mask) > 0] = 255
    height, width = binary.shape
    x, y, asset_width, asset_height = asset.bbox
    margin_x = max(2, int(asset_width * 0.08))
    margin_y = max(2, int(asset_height * 0.08))
    allowed = np.zeros_like(binary)
    left = max(0, x - margin_x)
    top = max(0, y - margin_y)
    right = min(width, x + asset_width + margin_x)
    bottom = min(height, y + asset_height + margin_y)
    allowed[top:bottom, left:right] = 255
    binary = cv2.bitwise_and(binary, allowed)

    count, labels, stats, _centroids = cv2.connectedComponentsWithStats(binary, connectivity=8)
    cleaned = np.zeros_like(binary)
    minimum_area = max(8, int(max(1, asset_width * asset_height) * 0.001))
    for component in range(1, count):
        if int(stats[component, cv2.CC_STAT_AREA]) >= minimum_area:
            cleaned[labels == component] = 255
    return cv2.GaussianBlur(cleaned, (3, 3), sigmaX=0.6)


def evaluate_segmentation_mask(mask: np.ndarray, asset: VisualAsset) -> tuple[dict[str, float | int], str | None]:
    binary = np.zeros(np.asarray(mask).shape[:2], dtype=np.uint8)
    binary[np.asarray(mask) > 0] = 255
    height, width = binary.shape
    x, y, asset_width, asset_height = asset.bbox
    left, top = max(0, x), max(0, y)
    right, bottom = min(width, x + asset_width), min(height, y + asset_height)
    box_area = max(1, (right - left) * (bottom - top))
    inside = binary[top:bottom, left:right]
    area = int(np.count_nonzero(inside))
    area_ratio = area / box_area

    ys, xs = np.nonzero(inside)
    if area:
        bbox_width = int(xs.max() - xs.min() + 1)
        bbox_height = int(ys.max() - ys.min() + 1)
        bbox_area = bbox_width * bbox_height
    else:
        bbox_width = bbox_height = bbox_area = 0
    bbox_coverage = bbox_area / box_area
    fill_ratio = area / max(1, bbox_area)

    component_count = 0
    largest_component_ratio = 0.0
    if area:
        count, _labels, stats, _centroids = cv2.connectedComponentsWithStats(inside, connectivity=8)
        component_areas = [int(stats[index, cv2.CC_STAT_AREA]) for index in range(1, count)]
        component_count = len(component_areas)
        largest_component_ratio = max(component_areas, default=0) / area

    border = max(1, int(round(min(max(1, inside.shape[0]), max(1, inside.shape[1])) * 0.03)))
    border_mask = np.zeros_like(inside)
    if inside.size:
        border_mask[:border, :] = 255
        border_mask[-border:, :] = 255
        border_mask[:, :border] = 255
        border_mask[:, -border:] = 255
    edge_contact_ratio = int(np.count_nonzero((inside > 0) & (border_mask > 0))) / max(1, area)

    source = np.asarray(mask)
    intermediate_ratio = int(np.count_nonzero((source > 0) & (source < 255))) / box_area
    metrics: dict[str, float | int] = {
        "area_pixels": area,
        "area_ratio": round(area_ratio, 6),
        "bbox_coverage": round(bbox_coverage, 6),
        "fill_ratio": round(fill_ratio, 6),
        "component_count": component_count,
        "largest_component_ratio": round(largest_component_ratio, 6),
        "edge_contact_ratio": round(edge_contact_ratio, 6),
        "intermediate_alpha_ratio": round(intermediate_ratio, 6),
    }

    warnings = []
    if not area:
        warnings.append("蒙版为空")
    elif area_ratio < 0.05 or bbox_coverage < 0.55:
        warnings.append("疑似只选择了局部结构")
    if component_count > 12 or (component_count > 1 and largest_component_ratio < 0.50):
        warnings.append("蒙版过于碎片化")
    if area_ratio > 0.40 and bbox_coverage > 0.45 and fill_ratio > 0.92:
        warnings.append("疑似选中了大面积矩形背景")
    if area_ratio > 0.70 and edge_contact_ratio > 0.35:
        warnings.append("蒙版主要贴着候选框边缘，疑似背景区域")
    if intermediate_ratio > 0.25 and area_ratio > 0.75 and fill_ratio > 0.85:
        warnings.append("蒙版包含大面积半透明矩形")
    return metrics, "；".join(warnings) if warnings else None


def save_segmentation_debug(
    project: PPTProject,
    slide: PPTSlide,
    asset: VisualAsset,
    result,
    points: list[tuple[float, float, int]] | tuple[tuple[float, float, int], ...] = (),
    run_id: str | None = None,
) -> Path:
    run_id = run_id or f"{time.strftime('%Y%m%d-%H%M%S')}-{time.time_ns() % 1_000_000_000:09d}"
    debug_dir = project.work_dir / "debug_masks" / f"slide_{slide.index:02d}" / asset.asset_id / run_id
    debug_dir.mkdir(parents=True, exist_ok=True)
    with Image.open(slide.image_path) as source_image:
        source = np.asarray(source_image.convert("RGB"))
    x, y, asset_width, asset_height = asset.bbox
    left, top = max(0, x), max(0, y)
    right, bottom = min(slide.image_width, x + asset_width), min(slide.image_height, y + asset_height)
    Image.fromarray(source[top:bottom, left:right]).save(debug_dir / "source_crop.png")
    (debug_dir / "box.json").write_text(
        json.dumps({"bbox": list(asset.bbox), "points": [list(point) for point in points]}, ensure_ascii=False, indent=2),
        encoding="utf-8",
    )
    candidates = []
    for index, (mask, score) in enumerate(zip(result.masks, result.scores), start=1):
        candidate = np.zeros(np.asarray(mask).shape[:2], dtype=np.uint8)
        candidate[np.asarray(mask) > 0] = 255
        Image.fromarray(candidate).save(debug_dir / f"candidate_{index}_raw.png")
        metrics, warning = evaluate_segmentation_mask(candidate, asset)
        candidates.append({"index": index - 1, "score": float(score), "metrics": metrics, "warning": warning})
    (debug_dir / "scores.json").write_text(
        json.dumps(
            {"selected_index": int(result.selected_index), "device": result.device, "model_id": result.model_id, "candidates": candidates},
            ensure_ascii=False,
            indent=2,
        ),
        encoding="utf-8",
    )
    selected_cleaned = _clean_segmentation_mask(result.mask, asset, slide.boxes)
    Image.fromarray(selected_cleaned).save(debug_dir / "selected_cleaned.png")
    rgba = np.dstack((source[top:bottom, left:right], selected_cleaned[top:bottom, left:right]))
    Image.fromarray(rgba, "RGBA").save(debug_dir / "selected_asset.png")
    return debug_dir


def store_visual_asset_mask(
    project: PPTProject,
    slide: PPTSlide,
    asset: VisualAsset,
    result,
    points: list[tuple[float, float, int]] | tuple[tuple[float, float, int], ...] = (),
):
    if result.mask.shape[:2] != (slide.image_height, slide.image_width):
        raise ValueError("SAM 2.1 蒙版尺寸与幻灯片图片不一致。")
    with Image.open(slide.image_path) as source_image:
        image = np.asarray(source_image.convert("RGB"))
    alpha = _clean_segmentation_mask(result.mask, asset, slide.boxes)
    if not np.any(alpha):
        raise RuntimeError("SAM 2.1 未在候选框内找到有效图片区域。")
    _metrics, warning = evaluate_segmentation_mask(alpha, asset)
    if warning:
        asset.segmentation_warning = warning
        asset.status = "segmentation_warning"
        asset.confirmed = False
        raise RuntimeError(f"SAM 2.1 蒙版未通过安全检查：{warning}")
    _write_visual_asset(project, slide, asset, image, alpha)
    asset.segmentation_mode = "sam2"
    asset.confidence = float(result.confidence)
    asset.model_id = str(result.model_id)
    asset.mask_version = 1
    asset.segmentation_warning = None
    asset.sam_selected_index = int(result.selected_index)
    asset.sam_points = tuple((float(px), float(py), int(label)) for px, py, label in points)
    asset.confirmed = True
    asset.status = "confirmed"


def _write_visual_asset(project: PPTProject, slide: PPTSlide, asset: VisualAsset, image: np.ndarray, alpha: np.ndarray):
    assets_dir = project.assets_dir or project.work_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    height, width = image.shape[:2]
    x, y, asset_width, asset_height = asset.bbox
    left, top = max(0, x), max(0, y)
    right, bottom = min(width, x + asset_width), min(height, y + asset_height)
    if right <= left or bottom <= top:
        return
    stem = f"{slide.image_path.stem}-{asset.asset_id}"
    image_path = assets_dir / f"{stem}.png"
    mask_path = assets_dir / f"{stem}.mask.png"
    local_alpha = refine_asset_alpha(image[top:bottom, left:right], alpha[top:bottom, left:right])
    expanded = build_background_erase_mask(local_alpha, 3)
    ring = (expanded > 0) & (local_alpha == 0)
    page_crop = image[top:bottom, left:right]
    samples = page_crop[ring]
    if samples.size == 0:
        samples = np.concatenate((page_crop[0], page_crop[-1], page_crop[:, 0], page_crop[:, -1]), axis=0)
    background_rgb = tuple(int(value) for value in np.median(samples, axis=0))
    rgba = decontaminate_asset_rgba(np.dstack((page_crop, local_alpha)), background_rgb)
    Image.fromarray(rgba, "RGBA").save(image_path)
    Image.fromarray(local_alpha).save(mask_path)
    asset.image_path = image_path
    asset.mask_path = mask_path


def build_text_stroke_mask(image: np.ndarray, box: OCRBox) -> tuple[np.ndarray, str, str | None]:
    """Return a narrow text-pixel mask, falling back to the legacy rectangle when needed."""
    image_height, image_width = image.shape[:2]
    if not box.text_regions:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "缺少 OCR 文字轮廓"

    allowed = _rect_mask(image_width, image_height, box.erase_rect)
    regions = np.zeros_like(allowed)
    for region in box.text_regions:
        points = np.asarray(region, dtype=np.int32)
        cv2.fillPoly(regions, [points], 255)
    allowed = cv2.bitwise_and(allowed, regions)
    region_pixels = int(np.count_nonzero(allowed))
    if region_pixels == 0:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "OCR 文字轮廓无效"

    gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
    # Text strokes have local contrast in either direction. This preserves the approach for
    # light text on dark fills without assuming a fixed text colour.
    blurred = cv2.GaussianBlur(gray, (0, 0), sigmaX=5, sigmaY=5)
    contrast = cv2.absdiff(gray, blurred)
    strokes = np.zeros_like(allowed)
    strokes[(contrast >= 18) & (allowed > 0)] = 255
    stroke_pixels = int(np.count_nonzero(strokes))
    density = stroke_pixels / region_pixels
    if density < 0.01:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "未找到足够的文字笔画"
    if density > 0.60:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "文字区域线条过密，可能与插图重叠"

    _count, _labels, stats, _centroids = cv2.connectedComponentsWithStats(strokes, connectivity=8)
    large_thin_component = any(
        int(component[cv2.CC_STAT_AREA]) > max(80, int(region_pixels * 0.01))
        and (
            int(component[cv2.CC_STAT_WIDTH]) / max(1, int(component[cv2.CC_STAT_HEIGHT])) <= 0.30
            or int(component[cv2.CC_STAT_HEIGHT]) / max(1, int(component[cv2.CC_STAT_WIDTH])) <= 0.30
        )
        for component in stats[1:]
    )
    if large_thin_component:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "检测到大面积连通线条，可能与插图重叠"

    line_height = max(1, int(box.line_height or box.bbox[3]))
    expansion_radius = min(4, max(1, int(round(line_height * 0.08))))
    kernel_size = expansion_radius * 2 + 1
    strokes = cv2.dilate(
        strokes,
        np.ones((kernel_size, kernel_size), dtype=np.uint8),
        iterations=1,
    )
    return cv2.bitwise_and(strokes, allowed), "text_stroke", None


def build_masks(project: PPTProject, progress: ProgressCB = None):
    project.masks_dir.mkdir(parents=True, exist_ok=True)
    (project.assets_dir or project.work_dir / "assets").mkdir(parents=True, exist_ok=True)
    total_slides = max(1, len(project.slides))
    for done, slide in enumerate(project.slides, start=1):
        with Image.open(slide.image_path) as source_image:
            image = np.asarray(source_image.convert("RGB"))
        mask = np.zeros((slide.image_height, slide.image_width), dtype=np.uint8)
        for box in slide.boxes:
            if not box.enabled:
                continue
            box_mask, box.mask_mode, box.mask_reason = build_text_stroke_mask(image, box)
            text_erase_mask = expand_text_erase_mask(box_mask, int(box.line_height or box.bbox[3]))
            mask = cv2.bitwise_or(mask, text_erase_mask)
            if box.mask_reason:
                _log(progress, f"第 {slide.index} 页文字框已回退矩形擦除：{box.mask_reason}")
        for asset in slide.visual_assets:
            if not asset.enabled or not asset.confirmed:
                continue
            if asset.segmentation_mode == "sam2" and (
                asset.segmentation_warning
                or asset.status == "segmentation_warning"
                or not asset.mask_path
                or not asset.mask_path.is_file()
            ):
                _log(progress, f"第 {slide.index} 页 SAM 图片未确认或存在警告，已保留原底图：{asset.asset_id}")
                continue
            alpha = visual_asset_alpha_mask(image, asset, slide.boxes)
            if not np.any(alpha):
                asset.enabled = False
                asset.status = "disabled_empty"
                _log(progress, f"第 {slide.index} 页图片区候选为空，已跳过：{asset.asset_id}")
                continue
            _write_visual_asset(project, slide, asset, image, alpha)
            expansion = min(20, max(3, int(round(min(asset.bbox[2], asset.bbox[3]) * 0.04))))
            mask = cv2.bitwise_or(mask, build_background_erase_mask(alpha, expansion))
        if slide.remove_watermark and slide.watermark_rect:
            left, top, right, bottom = slide.watermark_rect
            cv2.rectangle(mask, (left, top), (right, bottom), 255, -1)
        mask_path = project.masks_dir / slide.image_name
        Image.fromarray(mask).save(mask_path)
        if not mask_path.exists():
            raise RuntimeError(f"擦除蒙版写入失败：{mask_path}")
        _log(progress, f"第 {slide.index} 页擦除蒙版已生成")
        _progress(progress, int(done * 100 / total_slides), f"生成擦除蒙版：{done}/{total_slides}")


def os_environ_with_pythonpath():
    env = dict(os.environ)
    pythonpath_parts = [str(Path(IOPAINT_DEPS_DIR).resolve()), str(Path(DEPS_DIR).resolve())]
    if env.get("PYTHONPATH"):
        pythonpath_parts.append(env["PYTHONPATH"])
    env["PYTHONPATH"] = ";".join(pythonpath_parts)
    bundled_model_root = bundled_iopaint_model_root()
    if bundled_model_root:
        env["XDG_CACHE_HOME"] = str(bundled_model_root)
        env["U2NET_HOME"] = str(bundled_model_root)
    return env


def repair_visual_assets(
    project: PPTProject,
    model,
    inpaint_request,
    progress: ProgressCB = None,
) -> None:
    """Repair OCR text inside confirmed visual assets without punching alpha holes."""
    for slide in project.slides:
        if not slide.image_path.is_file():
            continue
        with Image.open(slide.image_path) as source_image:
            page = np.asarray(source_image.convert("RGB"))
        page_height, page_width = page.shape[:2]
        for asset in slide.visual_assets:
            if (
                not asset.enabled
                or not asset.confirmed
                or asset.segmentation_warning
                or asset.status == "segmentation_warning"
                or not asset.image_path
                or not asset.image_path.is_file()
            ):
                continue
            text_mask = build_asset_text_mask(page, asset, slide.boxes)
            if not np.any(text_mask):
                asset.mask_version = max(asset.mask_version, 2)
                continue

            x, y, asset_width, asset_height = asset.bbox
            padding = max(24, int(round(max(asset_width, asset_height) * 0.08)))
            crop_left = max(0, x - padding)
            crop_top = max(0, y - padding)
            crop_right = min(page_width, x + asset_width + padding)
            crop_bottom = min(page_height, y + asset_height + padding)
            crop = page[crop_top:crop_bottom, crop_left:crop_right].copy()
            local_text_mask = text_mask[crop_top:crop_bottom, crop_left:crop_right].copy()
            model_mask = local_text_mask.copy()

            available = model_mask == 0
            fill_color = np.median(crop[available], axis=0).astype(np.uint8) if np.any(available) else np.zeros(3, np.uint8)
            for neighbor in slide.visual_assets:
                if neighbor is asset or not neighbor.enabled or not neighbor.confirmed:
                    continue
                nx, ny, nw, nh = neighbor.bbox
                left = max(crop_left, nx)
                top = max(crop_top, ny)
                right = min(crop_right, nx + nw)
                bottom = min(crop_bottom, ny + nh)
                if right <= left or bottom <= top:
                    continue
                lx0, ly0 = left - crop_left, top - crop_top
                lx1, ly1 = right - crop_left, bottom - crop_top
                neighbor_context = np.zeros_like(model_mask)
                neighbor_context[ly0:ly1, lx0:lx1] = 255
                ax0, ay0 = max(0, x - crop_left), max(0, y - crop_top)
                ax1 = min(neighbor_context.shape[1], x + asset_width - crop_left)
                ay1 = min(neighbor_context.shape[0], y + asset_height - crop_top)
                neighbor_context[ay0:ay1, ax0:ax1] = 0
                crop[neighbor_context > 0] = fill_color
                model_mask[neighbor_context > 0] = 255

            try:
                generated = model(crop, model_mask, inpaint_request)
                generated = cv2.cvtColor(np.asarray(generated), cv2.COLOR_BGR2RGB)
                line_heights = [
                    int(box.line_height or box.bbox[3])
                    for box in slide.boxes
                    if box.enabled
                ]
                feather = min(40, max(15, int(round(median(line_heights))) if line_heights else 15))
                repaired_crop = soft_blend_inpaint(crop, generated, local_text_mask, feather_px=feather)

                with Image.open(asset.image_path) as stored_asset:
                    rgba = np.asarray(stored_asset.convert("RGBA")).copy()
                asset_left = max(0, x) - crop_left
                asset_top = max(0, y) - crop_top
                repaired_asset = repaired_crop[
                    asset_top:asset_top + rgba.shape[0],
                    asset_left:asset_left + rgba.shape[1],
                ]
                if repaired_asset.shape[:2] != rgba.shape[:2]:
                    raise RuntimeError("局部修复结果与图片资产尺寸不一致。")
                rgba[:, :, :3] = repaired_asset
                Image.fromarray(rgba, "RGBA").save(asset.image_path)
                asset.mask_version = 2
                if asset.status == "repair_warning":
                    asset.status = "confirmed"
                _log(progress, f"第 {slide.index} 页图片内文字已局部修复：{asset.asset_id}")
            except Exception as exc:
                asset.status = "repair_warning"
                _log(progress, f"第 {slide.index} 页图片内文字修复失败，已保留原图：{asset.asset_id} ({exc})")


def run_iopaint(
    images_dir: Path,
    masks_dir: Path,
    cleaned_dir: Path,
    progress: ProgressCB = None,
    project: PPTProject | None = None,
    skip_stems: set[str] | None = None,
):
    cleaned_dir.mkdir(parents=True, exist_ok=True)
    skip_stems = skip_stems or set()
    _log(progress, "开始用 IOPaint 擦除底图文字")

    env = os_environ_with_pythonpath()
    for key in ("XDG_CACHE_HOME", "U2NET_HOME", "PYTHONPATH"):
        if key in env:
            os.environ[key] = env[key]

    prepare_iopaint_torch(progress)

    LaMa = load_iopaint_lama_class()
    from iopaint.helper import pil_to_bytes
    from iopaint.model.utils import torch_gc
    from iopaint.schema import Device, InpaintRequest

    image_paths = {
        path.stem: path
        for path in images_dir.iterdir()
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg"}
    }
    mask_paths = {
        path.stem: path
        for path in masks_dir.iterdir()
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg"}
    }
    if not image_paths:
        raise RuntimeError(f"没有找到待处理图片：{images_dir}")
    if not mask_paths:
        raise RuntimeError(f"没有找到擦除蒙版：{masks_dir}")

    total_images = len(image_paths)
    _progress(progress, 0, f"IOPaint 擦除处理中：0/{total_images}")

    if not LaMa.is_downloaded():
        _log(progress, "本机未找到 IOPaint lama 模型，开始准备模型")
        LaMa.download()

    first_mask = next(iter(mask_paths.values()))

    def process_with_device(device) -> None:
        _log(progress, f"IOPaint 使用设备：{device}")
        model = LaMa(device=device)
        inpaint_request = InpaintRequest()
        for done, (stem, image_path) in enumerate(sorted(image_paths.items()), start=1):
            if stem in skip_stems:
                _log(progress, f"第 {done} 页复用已验证的本地修复结果：{stem}")
                _progress(progress, int(done * 100 / total_images), f"IOPaint 擦除处理中：{done}/{total_images}")
                continue
            mask_path = mask_paths.get(stem, first_mask)
            with Image.open(image_path) as source_image:
                infos = source_image.info
                image = np.array(source_image.convert("RGB"))
            with Image.open(mask_path) as source_mask:
                mask = np.array(source_mask.convert("L"))

            if mask.shape[:2] != image.shape[:2]:
                mask = cv2.resize(
                    mask,
                    (image.shape[1], image.shape[0]),
                    interpolation=cv2.INTER_NEAREST,
                )
            mask[mask >= 127] = 255
            mask[mask < 127] = 0

            result = model(image, mask, inpaint_request)
            result = cv2.cvtColor(result, cv2.COLOR_BGR2RGB)
            result = finish_page_inpaint(image, result, mask)
            output_path = cleaned_dir / f"{stem}.png"
            result_image = Image.fromarray(result)
            output_path.write_bytes(pil_to_bytes(result_image, "png", 100, infos))
            compressed_path = save_compressed_cleaned_image(result_image, output_path)
            torch_gc()
            _log(progress, f"第 {done} 页 IOPaint 擦除后已生成：{compressed_path.name}")
            _progress(progress, int(done * 100 / total_images), f"IOPaint 擦除处理中：{done}/{total_images}")
        if project is not None and len(skip_stems) < total_images:
            repair_visual_assets(project, model, inpaint_request, progress)

    device = preferred_iopaint_device(Device)
    try:
        process_with_device(device)
    except Exception as exc:
        if device == Device.cpu:
            raise
        _log(progress, f"IOPaint CUDA 处理失败，已回退 CPU：{exc}")
        cleaned_dir.mkdir(parents=True, exist_ok=True)
        process_with_device(Device.cpu)

    missing = [name for name in image_paths if not (cleaned_dir / f"{name}.png").exists()]
    if missing:
        raise RuntimeError(f"IOPaint 未生成以下页面：{', '.join(sorted(missing))}")

    _log(progress, "IOPaint 擦除完成")


def _write_openai_edit_mask(erase_mask: np.ndarray, output_path: Path) -> Path:
    binary = np.asarray(erase_mask) > 0
    rgba = np.full((*binary.shape, 4), 255, dtype=np.uint8)
    rgba[:, :, 3][binary] = 0
    output_path.parent.mkdir(parents=True, exist_ok=True)
    Image.fromarray(rgba, "RGBA").save(output_path)
    return output_path


def run_quality_iopaint(
    project: PPTProject,
    progress: ProgressCB = None,
    *,
    quality_mode: QualityMode = QualityMode.LOCAL_FAST,
    online_pages: set[int] | None = None,
    accepted_local_pages: set[int] | None = None,
    openai_api_key: str | None = None,
) -> list[PageQualityResult]:
    """Run the existing local repair only for invalidated pages and persist page-level QA."""
    online_pages = set(online_pages or ())
    accepted_local_pages = set(accepted_local_pages or ())
    pipeline = QualityPipeline(project.work_dir / "quality")
    project.cleaned_dir.mkdir(parents=True, exist_ok=True)
    sessions = {}
    pending_stems: set[str] = set()
    slide_by_stem = {slide.image_path.stem: slide for slide in project.slides}

    for stem, slide in slide_by_stem.items():
        mask_path = project.masks_dir / slide.image_name
        if not mask_path.is_file():
            raise FileNotFoundError(f"缺少第 {slide.index} 页擦除蒙版：{mask_path}")
        settings = {
            "quality_manifest_version": 1,
            "repair_backend": "local_lama",
            "mask_sha256": hashlib.sha256(mask_path.read_bytes()).hexdigest(),
            "asset_mask_versions": [asset.mask_version for asset in slide.visual_assets if asset.enabled],
        }
        session = pipeline.begin_page(slide.index, slide.image_path, settings)
        sessions[stem] = session
        output_path = project.cleaned_dir / f"{stem}.png"
        if session.reused and session.result.cleaned_path:
            shutil.copyfile(session.result.cleaned_path, output_path)
            with Image.open(output_path) as cached:
                save_compressed_cleaned_image(cached.convert("RGB"), output_path)
            _log(progress, f"第 {slide.index} 页质量缓存命中，跳过本地修复")
        else:
            pending_stems.add(stem)

    if pending_stems:
        run_iopaint(
            project.images_dir,
            project.masks_dir,
            project.cleaned_dir,
            progress,
            project=project,
            skip_stems=set(slide_by_stem) - pending_stems,
        )
    else:
        _log(progress, "所有页面均命中已验证质量缓存，已跳过 IOPaint")

    backend = OpenAIImageRepairBackend(openai_api_key) if online_pages else None
    results: list[PageQualityResult] = []
    for stem, slide in slide_by_stem.items():
        session = sessions[stem]
        if session.reused:
            results.append(session.result)
            continue
        output_path = project.cleaned_dir / f"{stem}.png"
        mask_path = project.masks_dir / slide.image_name
        with Image.open(slide.image_path) as source_image, Image.open(output_path) as repaired_image, Image.open(mask_path) as mask_image:
            source = np.asarray(source_image.convert("RGB"))
            repaired = np.asarray(repaired_image.convert("RGB"))
            erase_mask = np.asarray(mask_image.convert("L"))

        cache_path = pipeline.page_dir(slide.index) / "local_cleaned.png"
        cache_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copyfile(output_path, cache_path)
        result = pipeline.evaluate_local_quality(
            slide.index,
            source,
            repaired,
            erase_mask,
            source_path=slide.image_path,
            cleaned_path=cache_path,
        )

        if slide.index in online_pages:
            try:
                assert backend is not None
                edit_mask = _write_openai_edit_mask(erase_mask, cache_path.parent / "openai_edit_mask.png")
                online_path = cache_path.parent / "online_cleaned.png"
                backend.repair_background(
                    slide.image_path,
                    edit_mask,
                    online_path,
                    "Remove only the transparent masked text and foreground remnants. Preserve the original composition, colours, perspective, lighting, texture, and all unmasked content. Do not add text, logos, or new objects.",
                )
                with Image.open(online_path) as online_image:
                    online = np.asarray(online_image.convert("RGB"))
                if online.shape != source.shape:
                    raise RuntimeError("在线修复返回的页面尺寸与原图不一致。")
                online_result = pipeline.evaluate_local_quality(
                    slide.index,
                    source,
                    online,
                    erase_mask,
                    source_path=slide.image_path,
                    cleaned_path=online_path,
                )
                if online_result.status == QualityStatus.LOCAL_PROCESSED:
                    result = replace(online_result, status=QualityStatus.VALIDATED, mode=QualityMode.ONLINE_REPAIR)
                    shutil.copyfile(online_path, output_path)
                    with Image.open(output_path) as online_output:
                        save_compressed_cleaned_image(online_output.convert("RGB"), output_path)
                else:
                    result = replace(
                        result,
                        status=QualityStatus.REVIEW_REQUIRED,
                        issues=tuple((*result.issues, "online_quality_check_failed")),
                    )
            except Exception as exc:
                result = replace(
                    result,
                    status=QualityStatus.REVIEW_REQUIRED,
                    issues=tuple((*result.issues, "online_repair_failed")),
                )
                _log(progress, f"第 {slide.index} 页在线高质量修复失败，已保留本地结果：{exc}")

        if result.status == QualityStatus.LOCAL_PROCESSED:
            result = replace(result, status=QualityStatus.VALIDATED)
        elif slide.index in accepted_local_pages:
            result = replace(result, status=QualityStatus.VALIDATED, mode=QualityMode.LOCAL_REVIEWED)
        pipeline.complete_page(session, result)
        results.append(result)
        if result.status == QualityStatus.REVIEW_REQUIRED:
            _log(progress, f"第 {slide.index} 页需要人工检查：{', '.join(result.issues)}")
    return results


def save_compressed_cleaned_image(image: Image.Image, source_path: Path, quality: int = EXPORT_IMAGE_JPEG_QUALITY) -> Path:
    output_path = source_path.with_suffix(".jpg")
    save_kwargs = {
        "format": "JPEG",
        "quality": quality,
        "optimize": True,
        "progressive": True,
    }
    icc_profile = image.info.get("icc_profile")
    if icc_profile:
        save_kwargs["icc_profile"] = icc_profile
    rgb_image = image if image.mode == "RGB" else image.convert("RGB")
    rgb_image.save(output_path, **save_kwargs)
    return output_path


def cleaned_source_image_paths(cleaned_dir: Path) -> list[Path]:
    image_paths = [
        path
        for path in sorted(cleaned_dir.iterdir())
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg"}
    ]
    png_stems = {path.stem for path in image_paths if path.suffix.lower() == ".png"}
    return [
        path
        for path in image_paths
        if not (path.suffix.lower() in {".jpg", ".jpeg"} and path.stem in png_stems)
    ]


def upscale_cleaned_images(cleaned_dir: Path, scale: float = 2.0, progress: ProgressCB = None) -> int:
    if scale <= 1:
        _log(progress, "RealESRGAN 清晰化已跳过")
        return 0

    image_paths = cleaned_source_image_paths(cleaned_dir)
    if not image_paths:
        _log(progress, f"没有找到待清晰化图片：{cleaned_dir}")
        return 0

    env = os_environ_with_pythonpath()
    for key in ("XDG_CACHE_HOME", "U2NET_HOME", "PYTHONPATH"):
        if key in env:
            os.environ[key] = env[key]

    prepare_iopaint_torch(progress)

    from iopaint.model.utils import torch_gc
    from iopaint.plugins import RealESRGANUpscaler
    from iopaint.schema import Device, RunPluginRequest

    _log(progress, "开始用 RealESRGAN 提升导出底图清晰度")
    _progress(progress, 0, f"RealESRGAN 清晰化处理中：0/{len(image_paths)}")
    device = preferred_iopaint_device(Device)

    def create_upscaler(target_device):
        _log(progress, f"RealESRGAN 使用设备：{target_device}")
        return RealESRGANUpscaler("realesr-general-x4v3", target_device, no_half=(target_device == Device.cpu))

    try:
        upscaler = create_upscaler(device)
    except Exception as exc:
        if device == Device.cpu:
            raise
        _log(progress, f"RealESRGAN CUDA 处理失败，已回退 CPU：{exc}")
        device = Device.cpu
        upscaler = create_upscaler(device)

    for done, image_path in enumerate(image_paths, start=1):
        with Image.open(image_path) as source_image:
            infos = source_image.info
            rgb_image = np.array(source_image.convert("RGB"))

        try:
            bgr_result = upscaler.gen_image(
                rgb_image,
                RunPluginRequest(name=RealESRGANUpscaler.name, image="", scale=scale),
            )
        except Exception as exc:
            if device == Device.cpu:
                raise
            _log(progress, f"RealESRGAN CUDA 处理失败，已回退 CPU：{exc}")
            device = Device.cpu
            upscaler = create_upscaler(device)
            bgr_result = upscaler.gen_image(
                rgb_image,
                RunPluginRequest(name=RealESRGANUpscaler.name, image="", scale=scale),
            )
        rgb_result = cv2.cvtColor(bgr_result, cv2.COLOR_BGR2RGB)
        output_format = "JPEG" if image_path.suffix.lower() in {".jpg", ".jpeg"} else "PNG"
        result_image = Image.fromarray(rgb_result)
        if output_format == "JPEG":
            compressed_path = save_compressed_cleaned_image(result_image, image_path)
        else:
            result_image.save(image_path, format=output_format, **infos)
            compressed_path = save_compressed_cleaned_image(result_image, image_path)
        torch_gc()
        _log(progress, f"第 {done} 页已提升清晰度：{compressed_path.name}")
        _progress(progress, int(done * 100 / len(image_paths)), f"RealESRGAN 清晰化处理中：{done}/{len(image_paths)}")

    _log(progress, "RealESRGAN 清晰化完成")
    return len(image_paths)


def cleaned_image_path(cleaned_dir: Path, name: str):
    def preferred_existing(path: Path) -> Path | None:
        for candidate in [path.with_suffix(".jpg"), path.with_suffix(".jpeg"), path]:
            if candidate.exists():
                return candidate
        return None

    nested = cleaned_dir / name / name.replace("_clean", "")
    preferred = preferred_existing(nested)
    if preferred:
        return preferred
    direct = cleaned_dir / name
    preferred = preferred_existing(direct)
    if preferred:
        return preferred
    fallback = cleaned_dir / name.replace("_clean", "")
    preferred = preferred_existing(fallback)
    if preferred:
        return preferred
    raise FileNotFoundError(name)


def sample_text_color(image: Image.Image, rect, background_brightness: float | None = None):
    region = image.crop(rect).convert("RGB")
    arr = np.asarray(region)
    if arr.size == 0:
        return RGBColor(0, 0, 0)
    flat = arr.reshape(-1, 3)
    brightness = flat.astype(np.int32).sum(axis=1)
    if background_brightness is None:
        if arr.shape[0] > 2 and arr.shape[1] > 2:
            border = np.concatenate([arr[0, :, :], arr[-1, :, :], arr[:, 0, :], arr[:, -1, :]], axis=0)
        else:
            border = flat
        background_brightness = float(border.astype(np.int32).sum(axis=1).mean())
    sorted_indexes = np.argsort(brightness)
    sample_size = max(1, len(flat) // 8)
    if background_brightness < 382:
        sample = flat[sorted_indexes[-sample_size:]]
    else:
        sample = flat[sorted_indexes[:sample_size]]
    rgb = sample.mean(axis=0).astype(int)
    return RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def quantize_rgb_color(color: RGBColor, step: int = 24) -> RGBColor:
    channels = []
    for value in tuple(color):
        quantized = int(round(int(value) / step) * step)
        channels.append(max(0, min(255, quantized)))
    return RGBColor(*channels)


def character_color_runs(image: Image.Image, box: OCRBox, rotate_text: bool = False) -> list[tuple[str, RGBColor]]:
    text = box.text
    if not text:
        return []

    x, y, w, h = box.bbox
    runs: list[tuple[str, RGBColor]] = []
    current_text = ""
    current_color: RGBColor | None = None
    total = max(1, len(text))
    background_region = image.crop(box.erase_rect).convert("RGB")
    background_arr = np.asarray(background_region)
    background_brightness = None
    if background_arr.size:
        if background_arr.shape[0] > 2 and background_arr.shape[1] > 2:
            border = np.concatenate(
                [
                    background_arr[0, :, :],
                    background_arr[-1, :, :],
                    background_arr[:, 0, :],
                    background_arr[:, -1, :],
                ],
                axis=0,
            )
        else:
            border = background_arr.reshape(-1, 3)
        background_brightness = float(border.astype(np.int32).sum(axis=1).mean())

    for index, char in enumerate(text):
        if rotate_text:
            top = y + int(index * h / total)
            bottom = y + int((index + 1) * h / total)
            rect = (x, top, x + w, max(top + 1, bottom))
        else:
            left = x + int(index * w / total)
            right = x + int((index + 1) * w / total)
            rect = (left, y, max(left + 1, right), y + h)
        color = quantize_rgb_color(sample_text_color(image, rect, background_brightness=background_brightness))
        if current_color is not None and color == current_color:
            current_text += char
            continue
        if current_text and current_color is not None:
            runs.append((current_text, current_color))
        current_text = char
        current_color = color

    if current_text and current_color is not None:
        runs.append((current_text, current_color))
    return runs


def should_rotate_text(box: OCRBox) -> bool:
    if box.rotation in {90, 270, -90}:
        return True
    if "\n" in box.text:
        return False
    x, y, w, h = box.bbox
    return h > w * 1.45 and len(box.text.strip()) > 1


def add_textbox(slide, color_image, box: OCRBox, x_scale: float, y_scale: float):
    text = box.text
    if text in SKIP_TEXTS or (box.score < 0.45 and not box.manual and not box.edited) or not box.enabled:
        return False

    x, y, w, h = box.bbox
    left = int(x * x_scale)
    top = int(y * y_scale)
    width = int(w * x_scale)
    height = int(h * y_scale)
    rotate_text = should_rotate_text(box)
    font_axis = height
    if rotate_text:
        center_x = left + width / 2
        center_y = top + height / 2
        shape_left = int(center_x - height / 2)
        shape_top = int(center_y - width / 2)
        shape = slide.shapes.add_textbox(shape_left, shape_top, height, width)
        shape.rotation = box.rotation if box.rotation in {90, 270} else 270
        font_axis = width
    else:
        shape = slide.shapes.add_textbox(left, top, width, height)
        font_axis = box.line_height * y_scale if box.line_height is not None else height
    shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    font_size = Pt(max(8, font_axis * 0.72 / 12700))
    color_rect = (x, y, x + w, y + h)
    if "\n" in text:
        runs = [(text, sample_text_color(color_image, color_rect))]
    else:
        runs = character_color_runs(color_image, box, rotate_text=rotate_text) or [
            (text, sample_text_color(color_image, color_rect))
        ]
    for run_text, color in runs:
        run = p.add_run()
        run.text = run_text
        run.font.bold = True
        run.font.size = font_size
        run.font.color.rgb = color
    return True


def rebuild_ppt(project: PPTProject, output_pptx: Path, progress: ProgressCB = None):
    out = Presentation()
    out.slide_width = project.slide_width
    out.slide_height = project.slide_height
    while out.slides:
        rel_id = out.slides._sldIdLst[0].rId
        out.part.drop_rel(rel_id)
        del out.slides._sldIdLst[0]

    blank = out.slide_layouts[6]
    total_slides = max(1, len(project.slides))
    for done, slide_data in enumerate(project.slides, start=1):
        _progress(progress, int((done - 1) * 100 / total_slides), f"重建可编辑文本框：{done}/{total_slides}")
        cleaned_path = cleaned_image_path(project.cleaned_dir, slide_data.image_name)
        color_image = Image.open(slide_data.image_path).convert("RGB")
        dst = out.slides.add_slide(blank)
        dst.shapes.add_picture(
            str(cleaned_path),
            0,
            0,
            width=out.slide_width,
            height=out.slide_height,
        )
        x_scale = out.slide_width / slide_data.image_width
        y_scale = out.slide_height / slide_data.image_height

        def add_visual_asset(asset: VisualAsset):
            if not asset.enabled or not asset.confirmed or not asset.image_path or not asset.image_path.exists():
                return
            x, y, width, height = asset.bbox
            dst.shapes.add_picture(
                str(asset.image_path),
                int(x * x_scale),
                int(y * y_scale),
                width=int(width * x_scale),
                height=int(height * y_scale),
            )

        for asset in slide_data.visual_assets:
            if asset.layer != "above_text":
                add_visual_asset(asset)
        count = 0
        for box in slide_data.boxes:
            if slide_data.remove_watermark and slide_data.watermark_rect:
                x, y, width, height = box.bbox
                center_x = x + width / 2
                center_y = y + height / 2
                left, top, right, bottom = slide_data.watermark_rect
                if left <= center_x <= right and top <= center_y <= bottom:
                    continue
            if add_textbox(dst, color_image, box, x_scale, y_scale):
                count += 1
        for asset in slide_data.visual_assets:
            if asset.layer == "above_text":
                add_visual_asset(asset)
        _log(progress, f"第 {slide_data.index} 页已重建 {count} 个可编辑文本框")
        _progress(progress, int(done * 100 / total_slides), f"重建可编辑文本框：{done}/{total_slides}")
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    _progress(progress, 100, "正在保存 PPT 文件...")
    out.save(str(output_pptx))
    _log(progress, f"已导出：{output_pptx}")


def export_editable_ppt(
    project: PPTProject,
    output_pptx: Path,
    progress: ProgressCB = None,
    enhance_images: bool = True,
    quality_mode: QualityMode = QualityMode.LOCAL_FAST,
    online_pages: set[int] | None = None,
    accepted_local_pages: set[int] | None = None,
    openai_api_key: str | None = None,
):
    build_masks(project, progress)
    run_quality_iopaint(
        project,
        progress,
        quality_mode=quality_mode,
        online_pages=online_pages,
        accepted_local_pages=accepted_local_pages,
        openai_api_key=openai_api_key,
    )
    if enhance_images:
        upscale_cleaned_images(project.cleaned_dir, progress=progress)
    else:
        _log(progress, "已跳过 RealESRGAN 底图清晰化")
    rebuild_ppt(project, output_pptx, progress)
