import json
import subprocess
import sys
import types
import unittest
import unittest.mock
import zipfile
from pathlib import Path
from tempfile import TemporaryDirectory

import numpy as np
from PIL import Image

from app.core import (
    OCRBox,
    PPTProject,
    PPTSlide,
    VisualAsset,
    cleaned_image_path,
    export_editable_ppt,
    prefer_system_cuda_torch,
    prepare_iopaint_torch,
    preferred_iopaint_device,
    preferred_paddleocr_device,
    ppt_project_from_data,
    ppt_project_to_data,
    rebuild_ppt,
    run_export_editable_ppt_subprocess,
    run_iopaint,
    sample_text_color,
    default_watermark_rect,
    build_ocr_boxes,
    build_masks,
    build_asset_text_mask,
    build_text_stroke_mask,
    detect_visual_assets,
    finish_page_inpaint,
    visual_asset_alpha_mask,
    soft_blend_inpaint,
    load_project_cache,
    save_compressed_cleaned_image,
    upscale_cleaned_images,
)


class UpscaleCleanedImagesTest(unittest.TestCase):
    def test_detect_visual_assets_excludes_text_regions(self):
        image = np.full((100, 160, 3), 250, dtype=np.uint8)
        image[20:80, 100:145] = (0, 180, 220)
        image[25:45, 15:85] = (0, 0, 120)
        slide = PPTSlide(
            index=1,
            image_name="slide.png",
            image_path=Path("slide.png"),
            image_width=160,
            image_height=100,
            boxes=[
                OCRBox(
                    text="标题",
                    score=1.0,
                    bbox=(15, 25, 70, 20),
                    erase_rect=(10, 20, 90, 50),
                    text_regions=(((15, 25), (85, 25), (85, 45), (15, 45)),),
                )
            ],
        )

        assets = detect_visual_assets(image, slide)

        self.assertEqual(len(assets), 1)
        self.assertGreaterEqual(assets[0].bbox[0], 95)
        self.assertEqual(assets[0].layer, "below_text")

    def test_visual_asset_alpha_keeps_text_area_opaque_for_local_repair(self):
        image = np.full((80, 100, 3), 255, dtype=np.uint8)
        image[10:70, 20:90] = (10, 160, 220)
        image[30:45, 40:60] = (255, 255, 255)
        asset = VisualAsset(asset_id="asset-1", bbox=(20, 10, 70, 60))
        box = OCRBox(
            text="标题",
            score=1.0,
            bbox=(40, 30, 20, 15),
            erase_rect=(40, 30, 60, 45),
            text_regions=(((40, 30), (60, 30), (60, 45), (40, 45)),),
        )

        alpha = visual_asset_alpha_mask(image, asset, [box])

        self.assertEqual(alpha[25, 35], 255)
        self.assertEqual(alpha[5, 30], 0)
        self.assertEqual(alpha[35, 50], 255)

    def test_asset_text_mask_is_expanded_and_clipped_to_asset(self):
        image = np.full((80, 100, 3), 255, dtype=np.uint8)
        cv2 = __import__("cv2")
        cv2.putText(image, "A", (35, 45), cv2.FONT_HERSHEY_SIMPLEX, 0.8, (0, 0, 0), 2, cv2.LINE_AA)
        asset = VisualAsset(asset_id="asset-1", bbox=(20, 10, 50, 50))
        box = OCRBox(
            text="A",
            score=1.0,
            bbox=(30, 20, 30, 30),
            erase_rect=(25, 15, 65, 55),
            text_regions=(((30, 20), (60, 20), (60, 50), (30, 50)),),
            line_height=30,
        )

        mask = build_asset_text_mask(image, asset, [box])

        self.assertGreater(np.count_nonzero(mask), 30)
        self.assertEqual(int(np.count_nonzero(mask[:10])), 0)
        self.assertEqual(int(np.count_nonzero(mask[:, :20])), 0)
        self.assertEqual(int(np.count_nonzero(mask[:, 70:])), 0)

    def test_soft_blend_inpaint_preserves_unmasked_pixels_and_softens_boundary(self):
        original = np.zeros((31, 31, 3), dtype=np.uint8)
        repaired = np.full((31, 31, 3), 200, dtype=np.uint8)
        mask = np.zeros((31, 31), dtype=np.uint8)
        mask[8:23, 8:23] = 255

        blended = soft_blend_inpaint(original, repaired, mask, feather_px=6)

        self.assertEqual(tuple(blended[0, 0]), (0, 0, 0))
        self.assertGreater(int(blended[8, 15, 0]), 0)
        self.assertLess(int(blended[8, 15, 0]), 200)
        self.assertEqual(tuple(blended[15, 15]), (200, 200, 200))

    def test_finish_page_inpaint_prefers_surrounding_texture_for_large_flat_region(self):
        original = np.full((80, 100, 3), (30, 50, 70), dtype=np.uint8)
        original[:, :, 0] += np.arange(100, dtype=np.uint8)[None, :] // 10
        generated = np.full_like(original, (220, 20, 20))
        mask = np.zeros((80, 100), dtype=np.uint8)
        mask[15:65, 20:80] = 255

        finished = finish_page_inpaint(original, generated, mask)

        self.assertEqual(tuple(finished[0, 0]), tuple(original[0, 0]))
        self.assertLess(int(finished[40, 50, 0]), 100)
        self.assertGreater(int(finished[40, 50, 2]), 50)
        self.assertLess(abs(int(finished[15, 50, 0]) - int(original[15, 50, 0])), 30)

    def test_rebuild_places_visual_asset_between_background_and_text(self):
        from pptx import Presentation

        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir, masks_dir, cleaned_dir, assets_dir = (root / "images", root / "masks", root / "cleaned", root / "assets")
            for directory in (images_dir, masks_dir, cleaned_dir, assets_dir):
                directory.mkdir()
            original_path = images_dir / "slide.png"
            cleaned_path = cleaned_dir / "slide.png"
            asset_path = assets_dir / "device.png"
            Image.new("RGB", (100, 80), "white").save(original_path)
            Image.new("RGB", (100, 80), "white").save(cleaned_path)
            Image.new("RGBA", (30, 50), (0, 150, 220, 255)).save(asset_path)
            project = PPTProject(
                source_pptx=root / "source.pptx", work_dir=root, images_dir=images_dir,
                masks_dir=masks_dir, cleaned_dir=cleaned_dir, assets_dir=assets_dir,
                slide_width=914400, slide_height=914400,
                slides=[PPTSlide(
                    index=1, image_name="slide.png", image_path=original_path,
                    image_width=100, image_height=80,
                    visual_assets=[VisualAsset("device", (60, 10, 30, 50), image_path=asset_path, confirmed=True)],
                    boxes=[OCRBox("标题", 1.0, (10, 10, 30, 20), (10, 10, 40, 30))],
                )],
            )

            output = root / "out.pptx"
            rebuild_ppt(project, output)

            self.assertEqual(len(Presentation(output).slides[0].shapes), 3)
    def test_text_stroke_mask_is_narrow_and_preserves_nearby_illustration(self):
        image = np.full((100, 160, 3), 255, dtype=np.uint8)
        cv2 = __import__("cv2")
        cv2.putText(image, "TEXT", (10, 50), cv2.FONT_HERSHEY_SIMPLEX, 0.7, (0, 0, 0), 2, cv2.LINE_AA)
        cv2.line(image, (125, 10), (125, 90), (0, 0, 0), 2)
        box = OCRBox(
            text="TEXT",
            score=1.0,
            bbox=(5, 20, 95, 40),
            erase_rect=(0, 0, 150, 99),
            text_regions=(((5, 20), (105, 20), (105, 60), (5, 60)),),
        )

        mask, mode, reason = build_text_stroke_mask(image, box)

        self.assertEqual(mode, "text_stroke")
        self.assertIsNone(reason)
        self.assertGreater(np.count_nonzero(mask), 100)
        self.assertLess(np.count_nonzero(mask), 100 * 40)
        self.assertEqual(mask[50, 125], 0)

    def test_text_stroke_mask_falls_back_for_legacy_box_without_regions(self):
        image = np.full((40, 40, 3), 255, dtype=np.uint8)
        box = OCRBox(text="旧框", score=1.0, bbox=(10, 10, 10, 10), erase_rect=(8, 8, 22, 22))

        mask, mode, reason = build_text_stroke_mask(image, box)

        self.assertEqual(mode, "rectangle_fallback")
        self.assertEqual(reason, "缺少 OCR 文字轮廓")
        self.assertEqual(mask[8, 8], 255)
        self.assertEqual(mask[0, 0], 0)

    def test_text_stroke_mask_marks_large_thin_line_as_overlap_risk(self):
        image = np.full((100, 120, 3), 255, dtype=np.uint8)
        cv2 = __import__("cv2")
        cv2.line(image, (80, 10), (80, 90), (0, 0, 0), 3)
        box = OCRBox(
            text="疑似文本",
            score=1.0,
            bbox=(10, 10, 90, 80),
            erase_rect=(5, 5, 110, 95),
            text_regions=(((10, 10), (100, 10), (100, 90), (10, 90)),),
        )

        _mask, mode, reason = build_text_stroke_mask(image, box)

        self.assertEqual(mode, "rectangle_fallback")
        self.assertIn("连通线条", reason)

    def test_build_ocr_boxes_omits_notebooklm_watermark_text(self):
        slide = PPTSlide(
            index=1,
            image_path=Path("slide.png"),
            image_name="slide.png",
            image_width=1000,
            image_height=600,
        )
        page = {
            "dt_polys": [
                [[10, 10], [60, 10], [60, 30], [10, 30]],
                [[900, 560], [990, 560], [990, 590], [900, 590]],
            ],
            "rec_texts": ["Title", "NotebookLM"],
            "rec_scores": [0.99, 0.99],
        }

        boxes = build_ocr_boxes(page, slide)

        self.assertEqual([box.text for box in boxes], ["Title"])
        self.assertTrue(boxes[0].enabled)
        self.assertEqual(boxes[0].text_regions, (((10, 10), (60, 10), (60, 30), (10, 30)),))

    def test_prefer_system_cuda_torch_keeps_cuda_module_loaded(self):
        loaded_modules = {}
        cuda_torch = types.SimpleNamespace(
            __name__="torch",
            cuda=types.SimpleNamespace(is_available=lambda: True),
        )

        def fake_import(name):
            if name == "torch":
                loaded_modules[name] = cuda_torch
                sys.modules[name] = cuda_torch
                return cuda_torch
            if name == "torchvision":
                torchvision_module = types.SimpleNamespace(__name__="torchvision")
                loaded_modules[name] = torchvision_module
                sys.modules[name] = torchvision_module
                return torchvision_module
            raise AssertionError(name)

        with unittest.mock.patch.dict(sys.modules, {}, clear=True):
            selected = prefer_system_cuda_torch(
                original_sys_path=["system-site-packages"],
                bundled_paths={Path("bundled")},
                import_module=fake_import,
            )

            self.assertTrue(selected)
            self.assertIs(sys.modules["torch"], cuda_torch)

    def test_prefer_system_cuda_torch_requires_system_torchvision(self):
        cuda_torch = types.SimpleNamespace(
            __name__="torch",
            cuda=types.SimpleNamespace(is_available=lambda: True),
        )

        def fake_import(name):
            if name == "torch":
                sys.modules[name] = cuda_torch
                return cuda_torch
            if name == "torchvision":
                raise ImportError("system torchvision missing")
            raise AssertionError(name)

        with unittest.mock.patch.dict(sys.modules, {}, clear=True):
            selected = prefer_system_cuda_torch(
                original_sys_path=["system-site-packages"],
                bundled_paths={Path("bundled")},
                import_module=fake_import,
            )

            self.assertFalse(selected)
            self.assertNotIn("torch", sys.modules)

    def test_prefer_system_cuda_torch_removes_cpu_only_module(self):
        cpu_torch = types.SimpleNamespace(
            __name__="torch",
            cuda=types.SimpleNamespace(is_available=lambda: False),
        )

        def fake_import(name):
            sys.modules[name] = cpu_torch
            sys.modules["torch.cuda"] = types.SimpleNamespace()
            return cpu_torch

        with unittest.mock.patch.dict(sys.modules, {}, clear=True):
            selected = prefer_system_cuda_torch(
                original_sys_path=["system-site-packages"],
                bundled_paths={Path("bundled")},
                import_module=fake_import,
            )

            self.assertFalse(selected)
            self.assertNotIn("torch", sys.modules)
            self.assertNotIn("torch.cuda", sys.modules)

    def test_prefer_system_cuda_torch_keeps_frozen_cpu_module_loaded(self):
        cpu_torch = types.SimpleNamespace(
            __name__="torch",
            cuda=types.SimpleNamespace(is_available=lambda: False),
        )
        torch_cuda = types.SimpleNamespace()

        def fake_import(name):
            if name == "torch":
                sys.modules[name] = cpu_torch
                sys.modules["torch.cuda"] = torch_cuda
                return cpu_torch
            raise AssertionError(name)

        with unittest.mock.patch.dict(sys.modules, {}, clear=True):
            selected = prefer_system_cuda_torch(
                original_sys_path=["system-site-packages"],
                bundled_paths={Path("bundled")},
                import_module=fake_import,
                frozen_app=True,
            )

            self.assertFalse(selected)
            self.assertIs(sys.modules["torch"], cpu_torch)
            self.assertIs(sys.modules["torch.cuda"], torch_cuda)

    def test_prepare_iopaint_torch_logs_when_system_cuda_torch_selected(self):
        messages = []
        with unittest.mock.patch("app.core.prefer_system_cuda_torch", return_value=True):
            selected = prepare_iopaint_torch(progress=messages.append)

        self.assertTrue(selected)
        self.assertIn("已优先使用系统 CUDA Torch", messages)

    def test_prefers_cuda_for_iopaint_when_torch_cuda_is_available(self):
        device = types.SimpleNamespace(cpu="cpu", cuda="cuda")
        torch_module = types.SimpleNamespace(cuda=types.SimpleNamespace(is_available=lambda: True))

        self.assertEqual(preferred_iopaint_device(device, torch_module=torch_module), "cuda")

    def test_uses_cpu_for_iopaint_when_torch_cuda_is_unavailable(self):
        device = types.SimpleNamespace(cpu="cpu", cuda="cuda")
        torch_module = types.SimpleNamespace(cuda=types.SimpleNamespace(is_available=lambda: False))

        self.assertEqual(preferred_iopaint_device(device, torch_module=torch_module), "cpu")

    def test_prefers_gpu_for_paddleocr_when_paddle_is_cuda_build(self):
        paddle_module = types.SimpleNamespace(
            device=types.SimpleNamespace(is_compiled_with_cuda=lambda: True)
        )

        self.assertEqual(preferred_paddleocr_device(paddle_module=paddle_module), "gpu")

    def test_uses_cpu_for_paddleocr_when_paddle_is_cpu_build(self):
        paddle_module = types.SimpleNamespace(
            device=types.SimpleNamespace(is_compiled_with_cuda=lambda: False)
        )

        self.assertEqual(preferred_paddleocr_device(paddle_module=paddle_module), "cpu")

    def test_upscales_cleaned_images_with_realesrgan_and_preserves_rgb_colors(self):
        with TemporaryDirectory() as temp_dir:
            cleaned_dir = Path(temp_dir)
            image_path = cleaned_dir / "slide_01.png"
            Image.new("RGB", (2, 1), (10, 20, 30)).save(image_path)

            calls = []

            class FakeUpscaler:
                name = "RealESRGAN"

                def __init__(self, model_name, device, no_half=False):
                    calls.append(("init", model_name, device, no_half))

                def gen_image(self, rgb_np_img, request):
                    calls.append(("gen_image", tuple(rgb_np_img[0, 0]), request.scale))
                    return np.array([[[90, 80, 70], [60, 50, 40]]], dtype=np.uint8)

            class FakeRequest:
                def __init__(self, name, image, scale):
                    self.name = name
                    self.image = image
                    self.scale = scale

            plugins_module = types.ModuleType("iopaint.plugins")
            plugins_module.RealESRGANUpscaler = FakeUpscaler
            schema_module = types.ModuleType("iopaint.schema")
            schema_module.Device = types.SimpleNamespace(cpu="cpu", cuda="cuda")
            schema_module.RunPluginRequest = FakeRequest
            model_utils_module = types.ModuleType("iopaint.model.utils")
            model_utils_module.torch_gc = lambda: calls.append(("torch_gc",))
            torch_module = types.SimpleNamespace(cuda=types.SimpleNamespace(is_available=lambda: True))
            messages = []

            with unittest.mock.patch.dict(
                sys.modules,
                {
                    "iopaint.plugins": plugins_module,
                    "iopaint.schema": schema_module,
                    "iopaint.model.utils": model_utils_module,
                    "torch": torch_module,
                },
            ):
                processed = upscale_cleaned_images(cleaned_dir, scale=2, progress=messages.append)

            self.assertEqual(processed, 1)
            self.assertIn(("init", "realesr-general-x4v3", "cuda", False), calls)
            self.assertIn(("gen_image", (10, 20, 30), 2), calls)
            self.assertIn(("torch_gc",), calls)
            self.assertIn("第 1 页已提升清晰度：slide_01.jpg", messages)
            with Image.open(image_path) as output:
                self.assertEqual(output.size, (2, 1))
                self.assertEqual(output.getpixel((0, 0)), (70, 80, 90))
            jpeg_path = cleaned_dir / "slide_01.jpg"
            self.assertTrue(jpeg_path.exists())
            with Image.open(jpeg_path) as output:
                self.assertEqual(output.format, "JPEG")

    def test_realesrgan_falls_back_to_cpu_when_cuda_processing_fails(self):
        with TemporaryDirectory() as temp_dir:
            cleaned_dir = Path(temp_dir)
            image_path = cleaned_dir / "slide_01.png"
            Image.new("RGB", (2, 1), (10, 20, 30)).save(image_path)

            calls = []

            class FakeUpscaler:
                name = "RealESRGAN"

                def __init__(self, model_name, device, no_half=False):
                    self.device = device
                    calls.append(("init", model_name, device, no_half))

                def gen_image(self, rgb_np_img, request):
                    calls.append(("gen_image", self.device, request.scale))
                    if self.device == "cuda":
                        raise RuntimeError("CUDA out of memory")
                    return np.array([[[90, 80, 70], [60, 50, 40]]], dtype=np.uint8)

            class FakeRequest:
                def __init__(self, name, image, scale):
                    self.name = name
                    self.image = image
                    self.scale = scale

            plugins_module = types.ModuleType("iopaint.plugins")
            plugins_module.RealESRGANUpscaler = FakeUpscaler
            schema_module = types.ModuleType("iopaint.schema")
            schema_module.Device = types.SimpleNamespace(cpu="cpu", cuda="cuda")
            schema_module.RunPluginRequest = FakeRequest
            model_utils_module = types.ModuleType("iopaint.model.utils")
            model_utils_module.torch_gc = lambda: calls.append(("torch_gc",))
            torch_module = types.SimpleNamespace(cuda=types.SimpleNamespace(is_available=lambda: True))
            messages = []

            with unittest.mock.patch.dict(
                sys.modules,
                {
                    "iopaint.plugins": plugins_module,
                    "iopaint.schema": schema_module,
                    "iopaint.model.utils": model_utils_module,
                    "torch": torch_module,
                },
            ):
                processed = upscale_cleaned_images(cleaned_dir, scale=2, progress=messages.append)

            self.assertEqual(processed, 1)
            self.assertIn(("init", "realesr-general-x4v3", "cuda", False), calls)
            self.assertIn(("init", "realesr-general-x4v3", "cpu", True), calls)
            self.assertIn("RealESRGAN CUDA 处理失败，已回退 CPU：CUDA out of memory", messages)
            self.assertTrue((cleaned_dir / "slide_01.jpg").exists())

    def test_cleaned_image_path_prefers_compressed_jpeg_for_ppt_embedding(self):
        with TemporaryDirectory() as temp_dir:
            cleaned_dir = Path(temp_dir)
            png_path = cleaned_dir / "slide_01.png"
            jpeg_path = cleaned_dir / "slide_01.jpg"
            Image.new("RGB", (10, 10), (255, 255, 255)).save(png_path)
            Image.new("RGB", (10, 10), (255, 255, 255)).save(jpeg_path, format="JPEG")

            self.assertEqual(cleaned_image_path(cleaned_dir, "slide_01.png"), jpeg_path)

    def test_save_compressed_cleaned_image_does_not_convert_existing_rgb_images(self):
        with TemporaryDirectory() as temp_dir:
            source_path = Path(temp_dir) / "slide_01.png"
            image = Image.new("RGB", (4, 4), (255, 255, 255))

            with unittest.mock.patch.object(Image.Image, "convert", wraps=image.convert) as convert:
                output_path = save_compressed_cleaned_image(image, source_path)

            self.assertTrue(output_path.exists())
            convert.assert_not_called()

    def test_default_watermark_rect_covers_notebooklm_icon_area(self):
        rect = default_watermark_rect(1000, 600)

        self.assertLessEqual(rect[0], 740)
        self.assertLessEqual(rect[1], 480)
        self.assertEqual(rect[2:], (999, 599))

    def test_build_masks_erases_expanded_watermark_area(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            masks_dir = root / "masks"
            cleaned_dir = root / "cleaned"
            images_dir.mkdir()
            image_path = images_dir / "slide_01.png"
            Image.new("RGB", (1000, 600), (255, 255, 255)).save(image_path)
            slide = PPTSlide(
                index=1,
                image_name="slide_01.png",
                image_path=image_path,
                image_width=1000,
                image_height=600,
                watermark_rect=default_watermark_rect(1000, 600),
            )
            project = PPTProject(
                source_pptx=root / "source.pptx",
                work_dir=root,
                images_dir=images_dir,
                masks_dir=masks_dir,
                cleaned_dir=cleaned_dir,
                slides=[slide],
                slide_width=914400,
                slide_height=914400,
            )

            build_masks(project)

            with Image.open(masks_dir / "slide_01.png") as mask:
                self.assertEqual(mask.getpixel((745, 500)), 255)
                self.assertEqual(mask.getpixel((700, 500)), 0)

    def test_load_project_cache_expands_old_watermark_rect(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            image_path = root / "slide_01.png"
            Image.new("RGB", (1000, 600), (255, 255, 255)).save(image_path)
            slide = PPTSlide(
                index=1,
                image_name="slide_01.png",
                image_path=image_path,
                image_width=1000,
                image_height=600,
                watermark_rect=(810, 510, 999, 599),
            )
            project = PPTProject(
                source_pptx=root / "source.pptx",
                work_dir=root,
                images_dir=root,
                masks_dir=root / "masks",
                cleaned_dir=root / "cleaned",
                slides=[slide],
                slide_width=914400,
                slide_height=914400,
            )
            cache_path = root / "cache.json"
            cache_path.write_text(
                json.dumps(
                    {
                        "version": 2,
                        "slides": [
                            {
                                "index": 1,
                                "image_width": 1000,
                                "image_height": 600,
                                "watermark_rect": [810, 510, 999, 599],
                                "remove_watermark": True,
                                "boxes": [],
                            }
                        ]
                    }
                ),
                encoding="utf-8",
            )

            self.assertTrue(load_project_cache(project, cache_path=cache_path))

            self.assertEqual(project.slides[0].watermark_rect, default_watermark_rect(1000, 600))

    def test_iopaint_logs_each_generated_cleaned_page(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            masks_dir = root / "masks"
            cleaned_dir = root / "cleaned"
            images_dir.mkdir()
            masks_dir.mkdir()
            Image.new("RGB", (2, 1), (10, 20, 30)).save(images_dir / "slide_01.png")
            Image.new("L", (2, 1), 255).save(masks_dir / "slide_01.png")

            calls = []

            class FakeLaMa:
                @staticmethod
                def is_downloaded():
                    return True

                @staticmethod
                def download():
                    calls.append(("download", "lama"))

                def __init__(self, device):
                    calls.append(("init", "lama", device))

                def __call__(self, image, mask, request):
                    calls.append(("call", tuple(image[0, 0]), int(mask[0, 0])))
                    return np.array([[[90, 80, 70], [60, 50, 40]]], dtype=np.uint8)

            class FakeRequest:
                pass

            helper_module = types.ModuleType("iopaint.helper")

            def fake_pil_to_bytes(image, output_format, quality, infos):
                import io

                buffer = io.BytesIO()
                image.save(buffer, format=output_format.upper())
                return buffer.getvalue()

            helper_module.pil_to_bytes = fake_pil_to_bytes
            model_package = types.ModuleType("iopaint.model")
            model_package.__path__ = []
            lama_module = types.ModuleType("iopaint.model.lama")
            lama_module.LaMa = FakeLaMa
            model_utils_module = types.ModuleType("iopaint.model.utils")
            model_utils_module.torch_gc = lambda: calls.append(("torch_gc",))
            schema_module = types.ModuleType("iopaint.schema")
            schema_module.Device = types.SimpleNamespace(cpu="cpu", cuda="cuda")
            schema_module.InpaintRequest = FakeRequest
            torch_module = types.SimpleNamespace(cuda=types.SimpleNamespace(is_available=lambda: True))
            messages = []

            with unittest.mock.patch.dict(
                sys.modules,
                {
                    "iopaint.helper": helper_module,
                    "iopaint.model": model_package,
                    "iopaint.model.lama": lama_module,
                    "iopaint.model.utils": model_utils_module,
                    "iopaint.schema": schema_module,
                    "torch": torch_module,
                },
            ):
                run_iopaint(images_dir, masks_dir, cleaned_dir, progress=messages.append)

            self.assertIn(("init", "lama", "cuda"), calls)
            self.assertTrue((cleaned_dir / "slide_01.png").exists())
            self.assertTrue((cleaned_dir / "slide_01.jpg").exists())
            self.assertIn("第 1 页 IOPaint 擦除后已生成：slide_01.jpg", messages)

    def test_iopaint_falls_back_to_cpu_when_cuda_processing_fails(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            masks_dir = root / "masks"
            cleaned_dir = root / "cleaned"
            images_dir.mkdir()
            masks_dir.mkdir()
            Image.new("RGB", (2, 1), (10, 20, 30)).save(images_dir / "slide_01.png")
            Image.new("L", (2, 1), 255).save(masks_dir / "slide_01.png")

            calls = []

            class FakeLaMa:
                @staticmethod
                def is_downloaded():
                    return True

                @staticmethod
                def download():
                    calls.append(("download", "lama"))

                def __init__(self, device):
                    self.device = device
                    calls.append(("init", "lama", device))

                def __call__(self, image, mask, request):
                    calls.append(("call", self.device))
                    if self.device == "cuda":
                        raise RuntimeError("CUDA out of memory")
                    return np.array([[[90, 80, 70], [60, 50, 40]]], dtype=np.uint8)

            class FakeRequest:
                pass

            helper_module = types.ModuleType("iopaint.helper")

            def fake_pil_to_bytes(image, output_format, quality, infos):
                import io

                buffer = io.BytesIO()
                image.save(buffer, format=output_format.upper())
                return buffer.getvalue()

            helper_module.pil_to_bytes = fake_pil_to_bytes
            model_package = types.ModuleType("iopaint.model")
            model_package.__path__ = []
            lama_module = types.ModuleType("iopaint.model.lama")
            lama_module.LaMa = FakeLaMa
            model_utils_module = types.ModuleType("iopaint.model.utils")
            model_utils_module.torch_gc = lambda: calls.append(("torch_gc",))
            schema_module = types.ModuleType("iopaint.schema")
            schema_module.Device = types.SimpleNamespace(cpu="cpu", cuda="cuda")
            schema_module.InpaintRequest = FakeRequest
            torch_module = types.SimpleNamespace(cuda=types.SimpleNamespace(is_available=lambda: True))
            messages = []

            with unittest.mock.patch.dict(
                sys.modules,
                {
                    "iopaint.helper": helper_module,
                    "iopaint.model": model_package,
                    "iopaint.model.lama": lama_module,
                    "iopaint.model.utils": model_utils_module,
                    "iopaint.schema": schema_module,
                    "torch": torch_module,
                },
            ):
                run_iopaint(images_dir, masks_dir, cleaned_dir, progress=messages.append)

            self.assertIn(("init", "lama", "cuda"), calls)
            self.assertIn(("init", "lama", "cpu"), calls)
            self.assertIn("IOPaint CUDA 处理失败，已回退 CPU：CUDA out of memory", messages)
            self.assertTrue((cleaned_dir / "slide_01.jpg").exists())

    def test_rebuild_keeps_text_coordinates_based_on_original_slide_image_size(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            cleaned_dir = root / "cleaned"
            masks_dir = root / "masks"
            images_dir.mkdir()
            cleaned_dir.mkdir()
            masks_dir.mkdir()

            original_path = images_dir / "slide_01.png"
            cleaned_path = cleaned_dir / "slide_01.png"
            Image.new("RGB", (100, 100), (255, 255, 255)).save(original_path)
            Image.new("RGB", (200, 200), (255, 255, 255)).save(cleaned_path)

            project = PPTProject(
                source_pptx=root / "source.pptx",
                work_dir=root,
                images_dir=images_dir,
                masks_dir=masks_dir,
                cleaned_dir=cleaned_dir,
                slides=[
                    PPTSlide(
                        index=1,
                        image_name="slide_01.png",
                        image_path=original_path,
                        image_width=100,
                        image_height=100,
                        boxes=[
                            OCRBox(
                                text="Hello",
                                score=1.0,
                                bbox=(50, 20, 10, 10),
                                erase_rect=(50, 20, 60, 30),
                            )
                        ],
                    )
                ],
                slide_width=914400,
                slide_height=914400,
            )

            output_path = root / "out.pptx"
            rebuild_ppt(project, output_path)

            from pptx import Presentation

            deck = Presentation(str(output_path))
            text_shape = next(shape for shape in deck.slides[0].shapes if getattr(shape, "text", "") == "Hello")
            self.assertEqual(text_shape.left, 457200)

    def test_rebuild_omits_ocr_text_inside_enabled_watermark_region(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            cleaned_dir = root / "cleaned"
            masks_dir = root / "masks"
            images_dir.mkdir()
            cleaned_dir.mkdir()
            masks_dir.mkdir()

            original_path = images_dir / "slide_01.png"
            cleaned_path = cleaned_dir / "slide_01.png"
            Image.new("RGB", (100, 100), (255, 255, 255)).save(original_path)
            Image.new("RGB", (100, 100), (255, 255, 255)).save(cleaned_path)

            project = PPTProject(
                source_pptx=root / "source.pptx",
                work_dir=root,
                images_dir=images_dir,
                masks_dir=masks_dir,
                cleaned_dir=cleaned_dir,
                slides=[
                    PPTSlide(
                        index=1,
                        image_name="slide_01.png",
                        image_path=original_path,
                        image_width=100,
                        image_height=100,
                        boxes=[
                            OCRBox(
                                text="Heading",
                                score=1.0,
                                bbox=(10, 10, 30, 10),
                                erase_rect=(10, 10, 40, 20),
                            ),
                            OCRBox(
                                text="Gemini Notebook",
                                score=0.96,
                                bbox=(75, 85, 20, 10),
                                erase_rect=(73, 83, 97, 97),
                            ),
                        ],
                        watermark_rect=(70, 80, 99, 99),
                        remove_watermark=True,
                    )
                ],
                slide_width=914400,
                slide_height=914400,
            )

            output_path = root / "out.pptx"
            rebuild_ppt(project, output_path)

            from pptx import Presentation

            deck = Presentation(str(output_path))
            texts = [shape.text for shape in deck.slides[0].shapes if getattr(shape, "text", "")]
            self.assertEqual(texts, ["Heading"])

    def test_rebuild_preserves_light_text_color_on_dark_background(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            cleaned_dir = root / "cleaned"
            masks_dir = root / "masks"
            images_dir.mkdir()
            cleaned_dir.mkdir()
            masks_dir.mkdir()

            original_path = images_dir / "slide_01.png"
            cleaned_path = cleaned_dir / "slide_01.png"
            original = Image.new("RGB", (100, 40), (8, 18, 30))
            pixels = original.load()
            for x in range(20, 80):
                for y in range(10, 24):
                    pixels[x, y] = (235, 240, 245)
            original.save(original_path)
            Image.new("RGB", (100, 40), (8, 18, 30)).save(cleaned_path)

            project = PPTProject(
                source_pptx=root / "source.pptx",
                work_dir=root,
                images_dir=images_dir,
                masks_dir=masks_dir,
                cleaned_dir=cleaned_dir,
                slides=[
                    PPTSlide(
                        index=1,
                        image_name="slide_01.png",
                        image_path=original_path,
                        image_width=100,
                        image_height=40,
                        boxes=[
                            OCRBox(
                                text="Hello",
                                score=1.0,
                                bbox=(20, 10, 60, 14),
                                erase_rect=(16, 6, 84, 28),
                            )
                        ],
                    )
                ],
                slide_width=914400,
                slide_height=914400,
            )

            output_path = root / "out.pptx"
            rebuild_ppt(project, output_path)

            from pptx import Presentation

            deck = Presentation(str(output_path))
            text_shape = next(shape for shape in deck.slides[0].shapes if getattr(shape, "text", "") == "Hello")
            color = text_shape.text_frame.paragraphs[0].runs[0].font.color.rgb
            self.assertGreater(sum(tuple(color)), 600)

    def test_rebuild_preserves_mixed_text_colors_with_same_font_size(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            cleaned_dir = root / "cleaned"
            masks_dir = root / "masks"
            images_dir.mkdir()
            cleaned_dir.mkdir()
            masks_dir.mkdir()

            original_path = images_dir / "slide_01.png"
            cleaned_path = cleaned_dir / "slide_01.png"
            original = Image.new("RGB", (120, 30), (4, 18, 32))
            pixels = original.load()
            colors = [(245, 245, 245), (235, 130, 30), (245, 245, 245)]
            for index, color in enumerate(colors):
                for x in range(10 + index * 30, 38 + index * 30):
                    for y in range(6, 22):
                        pixels[x, y] = color
            original.save(original_path)
            Image.new("RGB", (120, 30), (4, 18, 32)).save(cleaned_path)

            project = PPTProject(
                source_pptx=root / "source.pptx",
                work_dir=root,
                images_dir=images_dir,
                masks_dir=masks_dir,
                cleaned_dir=cleaned_dir,
                slides=[
                    PPTSlide(
                        index=1,
                        image_name="slide_01.png",
                        image_path=original_path,
                        image_width=120,
                        image_height=30,
                        boxes=[
                            OCRBox(
                                text="ABC",
                                score=1.0,
                                bbox=(10, 6, 90, 16),
                                erase_rect=(8, 4, 102, 24),
                            )
                        ],
                    )
                ],
                slide_width=914400,
                slide_height=914400,
            )

            output_path = root / "out.pptx"
            rebuild_ppt(project, output_path)

            from pptx import Presentation

            deck = Presentation(str(output_path))
            text_shape = next(shape for shape in deck.slides[0].shapes if getattr(shape, "text", "") == "ABC")
            runs = text_shape.text_frame.paragraphs[0].runs
            run_colors = [tuple(run.font.color.rgb) for run in runs]
            font_sizes = {run.font.size for run in runs}

            self.assertEqual(len(font_sizes), 1)
            self.assertTrue(any(color[0] > 220 and color[1] > 220 and color[2] > 220 for color in run_colors))
            self.assertTrue(any(color[0] > 180 and 70 < color[1] < 180 and color[2] < 90 for color in run_colors))

    def test_sample_text_color_uses_light_text_on_dark_background(self):
        image = Image.new("RGB", (40, 20), (8, 18, 30))
        pixels = image.load()
        for x in range(12, 28):
            for y in range(6, 14):
                pixels[x, y] = (235, 240, 245)

        color = sample_text_color(image, (8, 4, 32, 16))

        self.assertGreater(sum(tuple(color)), 600)

    def test_sample_text_color_uses_dark_text_on_light_background(self):
        image = Image.new("RGB", (40, 20), (235, 240, 245))
        pixels = image.load()
        for x in range(12, 28):
            for y in range(6, 14):
                pixels[x, y] = (8, 18, 30)

        color = sample_text_color(image, (8, 4, 32, 16))

        self.assertLess(sum(tuple(color)), 100)

    def test_rebuild_embeds_compressed_jpeg_when_available(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            images_dir = root / "images"
            cleaned_dir = root / "cleaned"
            masks_dir = root / "masks"
            images_dir.mkdir()
            cleaned_dir.mkdir()
            masks_dir.mkdir()

            original_path = images_dir / "slide_01.png"
            cleaned_png_path = cleaned_dir / "slide_01.png"
            cleaned_jpeg_path = cleaned_dir / "slide_01.jpg"
            Image.new("RGB", (100, 100), (255, 255, 255)).save(original_path)
            Image.new("RGB", (100, 100), (255, 255, 255)).save(cleaned_png_path)
            Image.new("RGB", (100, 100), (255, 255, 255)).save(cleaned_jpeg_path, format="JPEG")

            project = PPTProject(
                source_pptx=root / "source.pptx",
                work_dir=root,
                images_dir=images_dir,
                masks_dir=masks_dir,
                cleaned_dir=cleaned_dir,
                slides=[
                    PPTSlide(
                        index=1,
                        image_name="slide_01.png",
                        image_path=original_path,
                        image_width=100,
                        image_height=100,
                        boxes=[],
                    )
                ],
                slide_width=914400,
                slide_height=914400,
            )

            output_path = root / "out.pptx"
            rebuild_ppt(project, output_path)

            with zipfile.ZipFile(output_path) as deck:
                media_names = [name for name in deck.namelist() if name.startswith("ppt/media/")]
            self.assertTrue(any(name.lower().endswith((".jpg", ".jpeg")) for name in media_names))

    def test_export_can_skip_realesrgan_enhancement(self):
        project = PPTProject(
            source_pptx=Path("source.pptx"),
            work_dir=Path("work"),
            images_dir=Path("images"),
            masks_dir=Path("masks"),
            cleaned_dir=Path("cleaned"),
            slides=[],
            slide_width=914400,
            slide_height=914400,
        )
        messages = []
        with (
            unittest.mock.patch("app.core.build_masks"),
            unittest.mock.patch("app.core.run_iopaint"),
            unittest.mock.patch("app.core.upscale_cleaned_images") as upscale,
            unittest.mock.patch("app.core.rebuild_ppt"),
        ):
            export_editable_ppt(project, Path("out.pptx"), progress=messages.append, enhance_images=False)

        upscale.assert_not_called()
        self.assertIn("已跳过 RealESRGAN 底图清晰化", messages)


    def test_project_serialization_preserves_export_state(self):
        project = PPTProject(
            source_pptx=Path("source.pptx"),
            work_dir=Path("work"),
            images_dir=Path("images"),
            masks_dir=Path("masks"),
            cleaned_dir=Path("cleaned"),
            slides=[
                PPTSlide(
                    index=1,
                    image_name="slide_01.png",
                    image_path=Path("images/slide_01.png"),
                    image_width=100,
                    image_height=80,
                    boxes=[
                        OCRBox(
                            text="Hello",
                            score=0.9,
                            bbox=(1, 2, 3, 4),
                            erase_rect=(1, 2, 10, 12),
                            enabled=False,
                            manual=True,
                            edited=True,
                            rotation=270,
                        )
                    ],
                    watermark_rect=(70, 60, 99, 79),
                    remove_watermark=False,
                    ocr_status="ok",
                    visual_assets=[
                        VisualAsset(
                            asset_id="device", bbox=(50, 5, 40, 65),
                            status="rule_candidate", layer="below_text",
                            image_path=Path("assets/device.png"),
                        )
                    ],
                )
            ],
            slide_width=914400,
            slide_height=685800,
        )

        restored = ppt_project_from_data(ppt_project_to_data(project))

        self.assertEqual(restored.source_pptx, Path("source.pptx"))
        self.assertEqual(restored.slides[0].watermark_rect, (70, 60, 99, 79))
        self.assertFalse(restored.slides[0].remove_watermark)
        self.assertEqual(restored.slides[0].boxes[0].text, "Hello")
        self.assertEqual(restored.slides[0].boxes[0].rotation, 270)
        self.assertEqual(restored.slides[0].visual_assets[0].asset_id, "device")
        self.assertEqual(restored.slides[0].visual_assets[0].image_path, Path("assets/device.png"))

    def test_export_subprocess_relays_progress_and_returns_output_path(self):
        class FakeStdout:
            def __iter__(self):
                return iter(["step 1\n", "step 2\n"])

        class FakeProcess:
            stdout = FakeStdout()

            def wait(self, timeout=None):
                return 0

            def kill(self):
                raise AssertionError("unexpected kill")

        captured = {}

        def fake_popen(command, **kwargs):
            captured["command"] = command
            captured["kwargs"] = kwargs
            return FakeProcess()

        project = PPTProject(
            source_pptx=Path("source.pptx"),
            work_dir=Path("work"),
            images_dir=Path("images"),
            masks_dir=Path("masks"),
            cleaned_dir=Path("cleaned"),
            slides=[],
            slide_width=914400,
            slide_height=685800,
        )
        messages = []

        output = run_export_editable_ppt_subprocess(
            project,
            Path("out.pptx"),
            progress=messages.append,
            enhance_images=False,
            popen=fake_popen,
        )

        self.assertEqual(output, Path("out.pptx"))
        self.assertIn("step 1", messages)
        self.assertIn("step 2", messages)
        self.assertIn("app.export_worker", captured["command"])
        self.assertEqual(captured["kwargs"]["stderr"], subprocess.STDOUT)
        self.assertEqual(captured["kwargs"]["env"]["PYTHONUTF8"], "1")
        self.assertEqual(captured["kwargs"]["env"]["PYTHONIOENCODING"], "utf-8")


if __name__ == "__main__":
    unittest.main()
