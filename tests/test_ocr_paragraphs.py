import json
import unittest
from pathlib import Path
from tempfile import TemporaryDirectory
from unittest.mock import patch

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from app.core import (
    OCRBox,
    PPTProject,
    PPTSlide,
    add_textbox,
    build_ocr_boxes,
    group_ocr_boxes,
    load_project_cache,
    save_project_cache,
    should_rotate_text,
)


class OcrParagraphGroupingTest(unittest.TestCase):
    @staticmethod
    def box(
        text: str,
        x: int,
        y: int,
        width: int = 100,
        height: int = 20,
        score: float = 1.0,
        erase_rect: tuple[int, int, int, int] | None = None,
        rotation: int = 0,
    ) -> OCRBox:
        return OCRBox(
            text=text,
            score=score,
            bbox=(x, y, width, height),
            erase_rect=erase_rect or (x - 2, y - 2, x + width + 2, y + height + 2),
            rotation=rotation,
            line_height=height,
        )

    def test_adjacent_same_format_lines_merge_and_preserve_line_height(self):
        boxes = group_ocr_boxes(
            [
                self.box("第一行", 20, 100, score=0.9),
                self.box("第二行", 20, 130, score=0.8),
                self.box("第三行", 20, 160, score=1.0),
            ]
        )

        self.assertEqual(len(boxes), 1)
        self.assertEqual(boxes[0].text, "第一行\n第二行\n第三行")
        self.assertEqual(boxes[0].bbox, (20, 100, 100, 80))
        self.assertEqual(boxes[0].erase_rect, (18, 98, 122, 182))
        self.assertAlmostEqual(boxes[0].score, 0.9)
        self.assertEqual(boxes[0].line_height, 20)

    def test_large_gap_starts_a_new_paragraph(self):
        boxes = group_ocr_boxes(
            [
                self.box("上段", 20, 100),
                self.box("上段续行", 20, 130),
                self.box("下段", 20, 210),
            ]
        )

        self.assertEqual([box.text for box in boxes], ["上段\n上段续行", "下段"])

    def test_different_height_columns_and_indent_do_not_merge(self):
        boxes = group_ocr_boxes(
            [
                self.box("标题", 20, 20, width=140, height=40),
                self.box("正文左栏", 20, 80),
                self.box("正文右栏", 240, 82),
                self.box("项目符号", 42, 110),
            ]
        )

        self.assertEqual(
            [box.text for box in boxes],
            ["标题", "正文左栏", "正文右栏", "项目符号"],
        )

    def test_two_columns_keep_their_own_multiline_paragraphs(self):
        boxes = group_ocr_boxes(
            [
                self.box("左栏第一行", 20, 100, width=100),
                self.box("右栏第一行", 240, 102, width=100),
                self.box("左栏第二行", 20, 130, width=100),
                self.box("右栏第二行", 240, 132, width=100),
            ]
        )

        self.assertEqual(
            [box.text for box in boxes],
            ["左栏第一行\n左栏第二行", "右栏第一行\n右栏第二行"],
        )

    def test_center_aligned_lines_merge_by_center(self):
        boxes = group_ocr_boxes(
            [
                self.box("短标题", 80, 100, width=80),
                self.box("较长的第二行", 50, 130, width=140),
            ]
        )

        self.assertEqual(len(boxes), 1)
        self.assertEqual(boxes[0].text, "短标题\n较长的第二行")

    def test_vertical_text_remains_single_box(self):
        boxes = group_ocr_boxes(
            [
                self.box("竖一", 20, 100, width=18, height=60, rotation=90),
                self.box("竖二", 20, 170, width=18, height=60, rotation=90),
            ]
        )

        self.assertEqual([box.text for box in boxes], ["竖一", "竖二"])

    def test_multiline_export_uses_source_line_height_for_font_size(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = self.box(
            "第一行\n第二行\n第三行",
            20,
            100,
            width=80,
            height=60,
        )
        box.line_height = 20

        with patch("app.core.sample_text_color", return_value=RGBColor(0, 0, 0)):
            self.assertTrue(
                add_textbox(
                    slide,
                    Image.new("RGB", (1, 1), (255, 255, 255)),
                    box,
                    10000.0,
                    10000.0,
                )
            )

        shape = slide.shapes[-1]
        run = shape.text_frame.paragraphs[0].runs[0]
        self.assertEqual(shape.text, "第一行\n第二行\n第三行")
        self.assertAlmostEqual(
            run.font.size,
            Pt(20 * 10000 * 0.72 / 12700),
            delta=256,
        )

    def test_export_does_not_force_font_family(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = self.box("文字", 20, 100, width=100, height=30)

        with patch("app.core.sample_text_color", return_value=RGBColor(0, 0, 0)):
            self.assertTrue(
                add_textbox(
                    slide,
                    Image.new("RGB", (500, 500), (255, 255, 255)),
                    box,
                    10000.0,
                    10000.0,
                )
            )

        run = slide.shapes[-1].text_frame.paragraphs[0].runs[0]
        self.assertIsNone(run.font.name)

    def test_multiline_export_converts_bbox_for_color_sampling(self):
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        box = self.box("第一行\n第二行", 20, 100, width=100, height=50)

        self.assertTrue(
            add_textbox(
                slide,
                Image.new("RGB", (500, 500), (255, 255, 255)),
                box,
                1.0,
                1.0,
            )
        )

    def test_multiline_horizontal_box_is_not_treated_as_vertical(self):
        box = self.box("第一行\n第二行", 20, 100, width=40, height=80)

        self.assertFalse(should_rotate_text(box))

    def test_build_ocr_boxes_returns_paragraph_boxes(self):
        slide = PPTSlide(
            index=1,
            image_name="slide.png",
            image_path=Path("slide.png"),
            image_width=500,
            image_height=400,
        )
        page = {
            "dt_polys": [
                [[20, 100], [120, 100], [120, 120], [20, 120]],
                [[20, 130], [120, 130], [120, 150], [20, 150]],
            ],
            "rec_texts": ["第一行", "第二行"],
            "rec_scores": [0.9, 0.8],
        }

        boxes = build_ocr_boxes(page, slide)

        self.assertEqual(len(boxes), 1)
        self.assertEqual(boxes[0].text, "第一行\n第二行")


class OcrParagraphCacheTest(unittest.TestCase):
    @staticmethod
    def project(root: Path) -> PPTProject:
        slide = PPTSlide(
            index=1,
            image_name="slide.png",
            image_path=root / "slide.png",
            image_width=500,
            image_height=400,
            boxes=[
                OCRBox(
                    text="第一行\n第二行",
                    score=0.9,
                    bbox=(20, 100, 100, 50),
                    erase_rect=(18, 98, 122, 152),
                    line_height=20,
                    text_regions=(((20, 100), (120, 100), (120, 120), (20, 120)),),
                    mask_mode="text_stroke",
                )
            ],
            ocr_status="ok",
        )
        return PPTProject(
            source_pptx=root / "source.pptx",
            work_dir=root,
            images_dir=root,
            masks_dir=root,
            cleaned_dir=root,
            slides=[slide],
            slide_width=1000,
            slide_height=800,
        )

    def test_old_cache_version_is_ignored(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            project = self.project(root)
            cache = root / "cache.json"
            save_project_cache(project, cache_path=cache)
            payload = json.loads(cache.read_text(encoding="utf-8"))
            payload["version"] = 1
            cache.write_text(json.dumps(payload), encoding="utf-8")

            self.assertFalse(load_project_cache(project, cache_path=cache))

    def test_new_cache_round_trips_line_height(self):
        with TemporaryDirectory() as temp_dir:
            root = Path(temp_dir)
            project = self.project(root)
            cache = root / "cache.json"
            save_project_cache(project, cache_path=cache)

            restored = self.project(root)
            restored.slides[0].boxes = []
            self.assertTrue(load_project_cache(restored, cache_path=cache))
            self.assertEqual(restored.slides[0].boxes[0].line_height, 20)
            self.assertEqual(restored.slides[0].boxes[0].mask_mode, "text_stroke")
            self.assertEqual(len(restored.slides[0].boxes[0].text_regions), 1)


if __name__ == "__main__":
    unittest.main()
