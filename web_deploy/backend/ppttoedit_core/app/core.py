from __future__ import annotations

import os
import re
import shutil
import subprocess
import sys
import json
import hashlib
import tempfile
import traceback
import importlib
import importlib.util
import types
from dataclasses import dataclass, field
from pathlib import Path
from statistics import median
from typing import Callable

BASE = Path(__file__).resolve().parent.parent
DEPS_DIR = BASE / ".py310deps"
IOPAINT_DEPS_DIR = BASE / ".py310iopaint"
ORIGINAL_SYS_PATH = list(sys.path)
for deps in [IOPAINT_DEPS_DIR, DEPS_DIR]:
    if deps.exists() and str(deps) not in sys.path:
        sys.path.insert(0, str(deps))


def _resolved_sys_path(path: str) -> Path:
    return Path(path or ".").resolve()


def prefer_system_cuda_torch(
    original_sys_path: list[str] | None = None,
    bundled_paths: set[Path] | None = None,
    import_module=importlib.import_module,
) -> bool:
    if "torch" in sys.modules:
        return bool(getattr(getattr(sys.modules["torch"], "cuda", None), "is_available", lambda: False)())

    original_sys_path = list(original_sys_path or ORIGINAL_SYS_PATH)
    bundled_paths = {path.resolve() for path in (bundled_paths or {IOPAINT_DEPS_DIR, DEPS_DIR})}
    current_sys_path = list(sys.path)
    before_modules = set(sys.modules)
    try:
        sys.path[:] = [
            path
            for path in original_sys_path
            if _resolved_sys_path(path) not in bundled_paths
        ]
        torch_module = import_module("torch")
        cuda = getattr(torch_module, "cuda", None)
        is_available = getattr(cuda, "is_available", None)
        if callable(is_available) and is_available():
            import_module("torchvision")
            return True
    except Exception:
        pass
    finally:
        sys.path[:] = current_sys_path

    for name in list(sys.modules):
        if name == "torch" or name.startswith("torch."):
            if name not in before_modules:
                sys.modules.pop(name, None)
    return False


prefer_system_cuda_torch()

import cv2
import numpy as np
from PIL import Image
from paddleocr import PaddleOCR
try:
    from paddleocr import PaddleOCRClient
except ImportError:
    PaddleOCRClient = None
try:
    from paddleocr import Model
except ImportError:
    Model = None
try:
    from paddleocr import OCROptions
except ImportError:
    OCROptions = None
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SKIP_TEXTS = {"NotebookLM"}
ProgressCB = Callable[[str], None] | None
PROGRESS_PREFIX = "__PPTTOEDIT_PROGRESS__|"
PAGE_READY_PREFIX = "__PPTTOEDIT_PAGE_READY__|"
CACHE_VERSION = 4
DEFAULT_PPT_WIDTH = Inches(13.333333)
OCR_BACKEND_LOCAL = "local"
OCR_BACKEND_REMOTE = "remote"
REMOTE_OCR_TOKEN_LENGTH = 40
EXPORT_IMAGE_JPEG_QUALITY = 85


def preferred_iopaint_device(device_enum, torch_module=None):
    if torch_module is None:
        try:
            import torch as torch_module
        except Exception:
            torch_module = None
    cuda_available = False
    if torch_module is not None:
        cuda = getattr(torch_module, "cuda", None)
        is_available = getattr(cuda, "is_available", None)
        if callable(is_available):
            try:
                cuda_available = bool(is_available())
            except Exception:
                cuda_available = False
    if cuda_available and hasattr(device_enum, "cuda"):
        return device_enum.cuda
    return device_enum.cpu


def preferred_paddleocr_device(paddle_module=None) -> str:
    if paddle_module is None:
        try:
            import paddle as paddle_module
        except Exception:
            paddle_module = None
    if paddle_module is not None:
        device = getattr(paddle_module, "device", None)
        is_compiled_with_cuda = getattr(device, "is_compiled_with_cuda", None)
        if callable(is_compiled_with_cuda):
            try:
                if is_compiled_with_cuda():
                    return "gpu"
            except Exception:
                pass
    return "cpu"


def bundle_root() -> Path:
    if getattr(sys, "frozen", False) and hasattr(sys, "_MEIPASS"):
        return Path(sys._MEIPASS)
    return BASE


def bundled_models_root() -> Path | None:
    candidate = bundle_root() / "models"
    return candidate if candidate.exists() else None


def bundled_ocr_model_dir(model_name: str) -> Path | None:
    root = bundled_models_root()
    if not root:
        return None
    candidate = root / "paddlex" / "official_models" / model_name
    return candidate if candidate.exists() else None


def bundled_iopaint_model_root() -> Path | None:
    root = bundled_models_root()
    return root if root and (root / "torch" / "hub" / "checkpoints" / "big-lama.pt").exists() else None


def load_iopaint_lama_class():
    module = sys.modules.get("iopaint.model.lama")
    if module is not None and hasattr(module, "LaMa"):
        return module.LaMa

    iopaint_spec = importlib.util.find_spec("iopaint")
    package_locations = list(iopaint_spec.submodule_search_locations or []) if iopaint_spec else []
    if not package_locations:
        raise RuntimeError("未找到 IOPaint 运行依赖。")

    model_dir = Path(package_locations[0]) / "model"
    lama_path = model_dir / "lama.py"
    if not lama_path.exists():
        raise RuntimeError(f"未找到 IOPaint LaMa 模块：{lama_path}")

    if "iopaint.model" not in sys.modules:
        package = types.ModuleType("iopaint.model")
        package.__path__ = [str(model_dir)]  # type: ignore[attr-defined]
        package.__package__ = "iopaint"
        sys.modules["iopaint.model"] = package

    spec = importlib.util.spec_from_file_location("iopaint.model.lama", lama_path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"无法加载 IOPaint LaMa 模块：{lama_path}")
    module = importlib.util.module_from_spec(spec)
    sys.modules["iopaint.model.lama"] = module
    spec.loader.exec_module(module)
    return module.LaMa


@dataclass
class OCRBox:
    text: str
    score: float
    bbox: tuple[int, int, int, int]
    erase_rect: tuple[int, int, int, int]
    enabled: bool = True
    manual: bool = False
    edited: bool = False
    rotation: int = 0
    line_height: int | None = None
    text_regions: tuple[tuple[tuple[int, int], ...], ...] = ()
    mask_mode: str = "pending"
    mask_reason: str | None = None

    def set_erase_rect(self, rect: tuple[int, int, int, int]):
        left, top, right, bottom = rect
        self.erase_rect = (
            int(left),
            int(top),
            int(max(left + 1, right)),
            int(max(top + 1, bottom)),
        )

    def set_bbox_from_rect(self, rect: tuple[int, int, int, int]):
        left, top, right, bottom = rect
        self.bbox = (
            int(left),
            int(top),
            int(max(1, right - left)),
            int(max(1, bottom - top)),
        )

    def reset_from_bbox(self, pad_x: int, pad_y: int, image_width: int, image_height: int):
        x, y, w, h = self.bbox
        left = max(0, x - pad_x)
        top = max(0, y - pad_y)
        right = min(image_width - 1, x + w + pad_x)
        bottom = min(image_height - 1, y + h + pad_y)
        self.erase_rect = (left, top, right, bottom)


@dataclass
class VisualAsset:
    asset_id: str
    bbox: tuple[int, int, int, int]
    enabled: bool = True
    source: str = "opencv"
    status: str = "rule_candidate"
    layer: str = "below_text"
    image_path: Path | None = None
    mask_path: Path | None = None
    mask_version: int = 0


@dataclass
class PPTSlide:
    index: int
    image_name: str
    image_path: Path
    image_width: int
    image_height: int
    boxes: list[OCRBox] = field(default_factory=list)
    watermark_rect: tuple[int, int, int, int] | None = None
    remove_watermark: bool = True
    visual_assets: list[VisualAsset] = field(default_factory=list)

    def reset_boxes(self, pad_x: int, pad_y: int):
        for box in self.boxes:
            box.reset_from_bbox(pad_x, pad_y, self.image_width, self.image_height)


@dataclass
class PPTProject:
    source_pptx: Path
    work_dir: Path
    images_dir: Path
    masks_dir: Path
    cleaned_dir: Path
    slides: list[PPTSlide]
    slide_width: int
    slide_height: int
    assets_dir: Path | None = None


def default_cache_path(source_pptx: Path) -> Path:
    return source_pptx.with_name(f"{source_pptx.stem}.ppttoedit.json")


def app_cache_path(source_pptx: Path) -> Path:
    cache_root = Path(os.environ.get("LOCALAPPDATA") or Path.home()) / "PPTEditableOCR" / "ocr_caches"
    digest = hashlib.sha1(str(source_pptx).encode("utf-8", errors="ignore")).hexdigest()[:16]
    return cache_root / f"{source_pptx.stem}-{digest}.ppttoedit.json"


def cache_path_candidates(source_pptx: Path) -> list[Path]:
    sidecar = default_cache_path(source_pptx)
    fallback = app_cache_path(source_pptx)
    return [sidecar] if sidecar == fallback else [sidecar, fallback]


def _rect_to_list(rect: tuple[int, int, int, int]) -> list[int]:
    return [int(value) for value in rect]


def _rect_from_data(value, fallback: tuple[int, int, int, int]) -> tuple[int, int, int, int]:
    if not isinstance(value, (list, tuple)) or len(value) != 4:
        return fallback
    try:
        return tuple(int(item) for item in value)
    except (TypeError, ValueError):
        return fallback


def _optional_int(value) -> int | None:
    if value is None:
        return None
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _regions_to_data(regions: tuple[tuple[tuple[int, int], ...], ...]) -> list[list[list[int]]]:
    return [[[int(x), int(y)] for x, y in region] for region in regions]


def _regions_from_data(value) -> tuple[tuple[tuple[int, int], ...], ...]:
    if not isinstance(value, (list, tuple)):
        return ()
    regions = []
    for region in value:
        if not isinstance(region, (list, tuple)) or len(region) < 3:
            continue
        points = []
        for point in region:
            if not isinstance(point, (list, tuple)) or len(point) < 2:
                points = []
                break
            try:
                points.append((int(point[0]), int(point[1])))
            except (TypeError, ValueError):
                points = []
                break
        if len(points) >= 3:
            regions.append(tuple(points))
    return tuple(regions)


def visual_asset_to_data(asset: VisualAsset) -> dict:
    return {
        "asset_id": asset.asset_id, "bbox": _rect_to_list(asset.bbox),
        "enabled": bool(asset.enabled), "source": asset.source, "status": asset.status,
        "layer": asset.layer, "image_path": str(asset.image_path) if asset.image_path else None,
        "mask_path": str(asset.mask_path) if asset.mask_path else None,
        "mask_version": int(asset.mask_version),
    }


def visual_asset_from_data(data: dict) -> VisualAsset:
    return VisualAsset(
        asset_id=str(data.get("asset_id") or "visual-asset"),
        bbox=_rect_from_data(data.get("bbox"), (0, 0, 1, 1)),
        enabled=bool(data.get("enabled", True)), source=str(data.get("source") or "opencv"),
        status=str(data.get("status") or "rule_candidate"), layer=str(data.get("layer") or "below_text"),
        image_path=Path(data["image_path"]) if data.get("image_path") else None,
        mask_path=Path(data["mask_path"]) if data.get("mask_path") else None,
        mask_version=int(data.get("mask_version", 0)),
    )


def save_project_cache(project: PPTProject, cache_path: Path | None = None, progress: ProgressCB = None) -> Path:
    data = {
        "version": CACHE_VERSION,
        "source_name": project.source_pptx.name,
        "slide_width": int(project.slide_width),
        "slide_height": int(project.slide_height),
        "slides": [
            {
                "index": int(slide.index),
                "image_name": slide.image_name,
                "image_width": int(slide.image_width),
                "image_height": int(slide.image_height),
                "watermark_rect": _rect_to_list(slide.watermark_rect) if slide.watermark_rect else None,
                "remove_watermark": bool(slide.remove_watermark),
                "visual_assets": [visual_asset_to_data(asset) for asset in slide.visual_assets],
                "boxes": [
                    {
                        "text": box.text,
                        "score": float(box.score),
                        "bbox": _rect_to_list(box.bbox),
                        "erase_rect": _rect_to_list(box.erase_rect),
                        "enabled": bool(box.enabled),
                        "manual": bool(box.manual),
                        "edited": bool(box.edited),
                        "rotation": int(box.rotation),
                        "line_height": int(box.line_height) if box.line_height is not None else None,
                        "text_regions": _regions_to_data(box.text_regions),
                        "mask_mode": box.mask_mode,
                        "mask_reason": box.mask_reason,
                    }
                    for box in slide.boxes
                ],
            }
            for slide in project.slides
        ],
    }
    payload = json.dumps(data, ensure_ascii=False, indent=2)
    candidates = [cache_path] if cache_path else cache_path_candidates(project.source_pptx)
    last_error: Exception | None = None
    for candidate in candidates:
        try:
            candidate.parent.mkdir(parents=True, exist_ok=True)
            candidate.write_text(payload, encoding="utf-8")
        except OSError as exc:
            last_error = exc
            continue
        _log(progress, f"识别框缓存已保存：{candidate}")
        return candidate
    raise RuntimeError(f"识别框缓存保存失败：{last_error}")


def load_project_cache(project: PPTProject, cache_path: Path | None = None, progress: ProgressCB = None) -> bool:
    candidates = [cache_path] if cache_path else cache_path_candidates(project.source_pptx)
    cache_path = next((path for path in candidates if path.exists()), None)
    if cache_path is None:
        return False

    try:
        data = json.loads(cache_path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        _log(progress, f"识别框缓存读取失败，已忽略：{cache_path} ({exc})")
        return False
    if _optional_int(data.get("version")) not in {2, 3, CACHE_VERSION}:
        _log(progress, f"识别框缓存版本不匹配，已忽略：{cache_path}")
        return False
    cached_slides = data.get("slides")
    if not isinstance(cached_slides, list) or len(cached_slides) != len(project.slides):
        _log(progress, f"识别框缓存页数不匹配，已忽略：{cache_path}")
        return False

    for slide, cached in zip(project.slides, cached_slides):
        if (
            int(cached.get("index", -1)) != slide.index
            or int(cached.get("image_width", -1)) != slide.image_width
            or int(cached.get("image_height", -1)) != slide.image_height
        ):
            _log(progress, f"识别框缓存页面尺寸不匹配，已忽略：{cache_path}")
            return False

    for slide, cached in zip(project.slides, cached_slides):
        slide.remove_watermark = bool(cached.get("remove_watermark", slide.remove_watermark))
        slide.watermark_rect = _rect_from_data(cached.get("watermark_rect"), slide.watermark_rect) if cached.get("watermark_rect") else slide.watermark_rect
        slide.boxes = []
        slide.visual_assets = [visual_asset_from_data(asset) for asset in cached.get("visual_assets", []) if isinstance(asset, dict)]
        for item in cached.get("boxes", []):
            if not isinstance(item, dict):
                continue
            bbox = _rect_from_data(item.get("bbox"), (0, 0, 1, 1))
            erase_rect = _rect_from_data(item.get("erase_rect"), bbox)
            slide.boxes.append(
                OCRBox(
                    text=str(item.get("text") or ""),
                    score=float(item.get("score", 1.0)),
                    bbox=bbox,
                    erase_rect=erase_rect,
                    enabled=bool(item.get("enabled", True)),
                    manual=bool(item.get("manual", False)),
                    edited=bool(item.get("edited", False)),
                    rotation=int(item.get("rotation", 0)),
                    line_height=_optional_int(item.get("line_height")),
                    text_regions=_regions_from_data(item.get("text_regions")),
                    mask_mode=str(item.get("mask_mode") or "pending"),
                    mask_reason=str(item["mask_reason"]) if item.get("mask_reason") else None,
                )
            )

    _log(progress, f"已加载识别框缓存，跳过 OCR：{cache_path}")
    for slide in project.slides:
        _page_ready(progress, slide)
    return True


def box_to_rect(poly) -> tuple[int, int, int, int]:
    pts = np.array(poly, dtype=np.int32)
    x, y, w, h = cv2.boundingRect(pts)
    return x, y, w, h


def normalize_text_region(poly) -> tuple[tuple[int, int], ...]:
    points = np.asarray(poly, dtype=np.int32).reshape(-1, 2)
    return tuple((int(x), int(y)) for x, y in points)


def default_expand_rect(
    x: int, y: int, w: int, h: int, width: int, height: int
) -> tuple[int, int, int, int]:
    pad_x = max(8, int(w * 0.05))
    pad_y = max(6, int(h * 0.18))
    left = max(0, x - pad_x)
    top = max(0, y - pad_y)
    right = min(width - 1, x + w + pad_x)
    bottom = min(height - 1, y + h + pad_y)
    return left, top, right, bottom


def _log(progress: ProgressCB, message: str):
    if progress:
        progress(message)


def _progress(progress: ProgressCB, percent: int, message: str):
    percent = max(0, min(100, int(percent)))
    _log(progress, f"{PROGRESS_PREFIX}{percent}|{message}")


def _page_ready(progress: ProgressCB, slide: PPTSlide, status: str = "ok"):
    _log(progress, f"{PAGE_READY_PREFIX}{slide.index}|{len(slide.boxes)}|{status}")


def unique_output_path(path: Path) -> Path:
    if not path.exists():
        return path
    for index in range(2, 1000):
        candidate = path.with_name(f"{path.stem}-{index}{path.suffix}")
        if not candidate.exists():
            return candidate
    raise RuntimeError(f"无法生成不覆盖已有文件的输出路径：{path}")


def default_pdf_pptx_path(source_pdf: Path) -> Path:
    return unique_output_path(source_pdf.with_name(f"{source_pdf.stem}-from-pdf.pptx"))


def fit_rect_to_slide(image_width: int, image_height: int, slide_width: int, slide_height: int) -> tuple[int, int, int, int]:
    scale = min(slide_width / image_width, slide_height / image_height)
    width = int(image_width * scale)
    height = int(image_height * scale)
    left = int((slide_width - width) / 2)
    top = int((slide_height - height) / 2)
    return left, top, width, height


def convert_pdf_to_pptx(source_pdf: Path, output_pptx: Path | None = None, progress: ProgressCB = None) -> Path:
    import pypdfium2

    source_pdf = source_pdf.expanduser().resolve()
    if output_pptx is None:
        output_pptx = default_pdf_pptx_path(source_pdf)
    else:
        output_pptx = Path(output_pptx).expanduser().resolve()

    _log(progress, f"开始转换 PDF：{source_pdf}")
    pdf = pypdfium2.PdfDocument(str(source_pdf))
    page_count = len(pdf)
    if page_count <= 0:
        raise RuntimeError(f"PDF 没有可转换页面：{source_pdf}")

    first_page = pdf[0]
    first_width, first_height = first_page.get_size()
    first_page.close()
    if first_width <= 0 or first_height <= 0:
        raise RuntimeError(f"PDF 首页尺寸无效：{source_pdf}")

    out = Presentation()
    out.slide_width = int(DEFAULT_PPT_WIDTH)
    out.slide_height = int(DEFAULT_PPT_WIDTH * first_height / first_width)
    while out.slides:
        rel_id = out.slides._sldIdLst[0].rId
        out.part.drop_rel(rel_id)
        del out.slides._sldIdLst[0]

    blank = out.slide_layouts[6]
    _progress(progress, 0, f"PDF 转 PPT：0/{page_count}")
    with tempfile.TemporaryDirectory(prefix="ppttoedit_pdf_") as temp_dir:
        temp_root = Path(temp_dir)
        for page_index in range(page_count):
            page = pdf[page_index]
            image_path = temp_root / f"page_{page_index + 1:04d}.png"
            bitmap = page.render(scale=2)
            image = bitmap.to_pil().convert("RGB")
            image_width, image_height = image.size
            image.save(image_path)
            bitmap.close()
            page.close()

            slide = out.slides.add_slide(blank)
            left, top, width, height = fit_rect_to_slide(
                image_width,
                image_height,
                out.slide_width,
                out.slide_height,
            )
            slide.shapes.add_picture(
                str(image_path),
                left,
                top,
                width=width,
                height=height,
            )
            done = page_index + 1
            _progress(progress, int(done * 100 / page_count), f"PDF 转 PPT：{done}/{page_count}")
            _log(progress, f"第 {done} 页已加入 PPT")

        output_pptx.parent.mkdir(parents=True, exist_ok=True)
        out.save(str(output_pptx))

    pdf.close()
    _log(progress, f"PDF 转 PPT 完成：{output_pptx}")
    return output_pptx


def extract_slide_images(source_pptx: Path, images_dir: Path, progress: ProgressCB = None):
    images_dir.mkdir(parents=True, exist_ok=True)
    src = Presentation(str(source_pptx))
    extracted: list[PPTSlide] = []

    for index, slide in enumerate(src.slides, start=1):
        pic = next((shape for shape in slide.shapes if getattr(shape, "image", None)), None)
        if pic is None:
            continue
        image_name = f"slide_{index:02d}.png"
        out_path = images_dir / image_name
        out_path.write_bytes(pic.image.blob)
        with Image.open(out_path) as image:
            width, height = image.size
        watermark_rect = (
            max(0, width - 190),
            max(0, height - 90),
            width - 1,
            height - 1,
        )
        extracted.append(
            PPTSlide(
                index=index,
                image_name=image_name,
                image_path=out_path,
                image_width=width,
                image_height=height,
                watermark_rect=watermark_rect,
            )
        )
        _log(progress, f"提取第 {index} 页图片")
    return src, extracted


def _make_remote_ocr_options():
    if OCROptions is None:
        return None
    try:
        return OCROptions(
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            visualize=False,
        )
    except TypeError:
        try:
            return OCROptions(
                use_doc_orientation_classify=False,
                use_doc_unwarping=False,
                use_textline_orientation=False,
            )
        except TypeError:
            return None


def _result_attr(result, *names):
    for name in names:
        if isinstance(result, dict) and name in result:
            return result[name]
        if hasattr(result, name):
            return getattr(result, name)
    return None


def _result_as_page_dict(page) -> dict:
    pruned = _result_attr(page, "prunedResult", "pruned_result", "res")
    if pruned is None and hasattr(page, "to_dict"):
        try:
            pruned = _result_attr(page.to_dict(), "prunedResult", "pruned_result", "res")
        except Exception:
            pruned = None
    if pruned is None:
        pruned = page
    if hasattr(pruned, "to_dict"):
        pruned = pruned.to_dict()
    if not isinstance(pruned, dict):
        raise RuntimeError(f"远端 OCR 返回格式不支持：{type(pruned).__name__}")
    return {
        "dt_polys": pruned.get("dt_polys") or pruned.get("dtPolys") or [],
        "rec_texts": pruned.get("rec_texts") or pruned.get("recTexts") or [],
        "rec_scores": pruned.get("rec_scores") or pruned.get("recScores") or [],
    }


def _predict_remote_ocr_page(client, image_path: Path) -> dict:
    model = getattr(Model, "PP_OCRV5", "PP-OCRv5") if Model is not None else "PP-OCRv5"
    kwargs = {
        "file_path": str(image_path),
        "model": model,
    }
    options = _make_remote_ocr_options()
    if options is not None:
        kwargs["options"] = options
    try:
        result = client.ocr(**kwargs)
    except TypeError:
        kwargs.pop("options", None)
        result = client.ocr(**kwargs)
    pages = _result_attr(result, "pages")
    if not pages:
        return {"dt_polys": [], "rec_texts": [], "rec_scores": []}
    return _result_as_page_dict(pages[0])


def _is_horizontal_ocr_box(box: OCRBox) -> bool:
    if box.rotation in {90, 270, -90}:
        return False
    x, y, width, height = box.bbox
    return not (height > width * 1.45 and len(box.text.strip()) > 1)


def _horizontal_gap(first: OCRBox, second: OCRBox) -> int:
    first_left, _first_top, first_width, _first_height = first.bbox
    second_left, _second_top, second_width, _second_height = second.bbox
    first_right = first_left + first_width
    second_right = second_left + second_width
    if first_right < second_left:
        return second_left - first_right
    if second_right < first_left:
        return first_left - second_right
    return 0


def _can_join_ocr_lines(first: OCRBox, second: OCRBox) -> bool:
    if not _is_horizontal_ocr_box(first) or not _is_horizontal_ocr_box(second):
        return False

    first_x, first_y, first_width, first_height = first.bbox
    second_x, second_y, second_width, second_height = second.bbox
    first_height_ref = first.line_height or first_height
    second_height_ref = second.line_height or second_height
    min_height = min(first_height_ref, second_height_ref)
    max_height = max(first_height_ref, second_height_ref)
    if min_height <= 0 or max_height / min_height > 1.35:
        return False

    vertical_gap = second_y - (first_y + first_height)
    if vertical_gap < -0.25 * min_height or vertical_gap > max(8, 0.9 * max_height):
        return False

    first_center = first_x + first_width / 2
    second_center = second_x + second_width / 2
    first_right = first_x + first_width
    second_right = second_x + second_width
    alignment_delta = min(
        abs(first_x - second_x),
        abs(first_center - second_center),
        abs(first_right - second_right),
    )
    if alignment_delta > 0.75 * max_height:
        return False
    return _horizontal_gap(first, second) <= 1.5 * max_height


def _ocr_line_join_score(first: OCRBox, second: OCRBox) -> float:
    if not _can_join_ocr_lines(first, second):
        return float("inf")
    first_x, first_y, first_width, first_height = first.bbox
    second_x, second_y, second_width, _second_height = second.bbox
    first_height_ref = first.line_height or first_height
    second_height_ref = second.line_height or second.bbox[3]
    first_center = first_x + first_width / 2
    second_center = second_x + second_width / 2
    first_right = first_x + first_width
    second_right = second_x + second_width
    alignment_delta = min(
        abs(first_x - second_x),
        abs(first_center - second_center),
        abs(first_right - second_right),
    )
    vertical_gap = second_y - (first_y + first_height)
    return abs(vertical_gap) + alignment_delta + abs(first_height_ref - second_height_ref)


def _merge_ocr_line_group(group: list[OCRBox]) -> OCRBox:
    if len(group) == 1:
        return group[0]

    left = min(box.bbox[0] for box in group)
    top = min(box.bbox[1] for box in group)
    right = max(box.bbox[0] + box.bbox[2] for box in group)
    bottom = max(box.bbox[1] + box.bbox[3] for box in group)
    erase_left = min(box.erase_rect[0] for box in group)
    erase_top = min(box.erase_rect[1] for box in group)
    erase_right = max(box.erase_rect[2] for box in group)
    erase_bottom = max(box.erase_rect[3] for box in group)
    line_heights = [box.line_height or box.bbox[3] for box in group]
    return OCRBox(
        text="\n".join(box.text for box in group),
        score=sum(box.score for box in group) / len(group),
        bbox=(left, top, right - left, bottom - top),
        erase_rect=(erase_left, erase_top, erase_right, erase_bottom),
        enabled=all(box.enabled for box in group),
        manual=any(box.manual for box in group),
        edited=any(box.edited for box in group),
        line_height=int(round(median(line_heights))),
        text_regions=tuple(region for box in group for region in box.text_regions),
    )


def group_ocr_boxes(boxes: list[OCRBox]) -> list[OCRBox]:
    candidates = [
        box
        for box in boxes
        if box.text.strip() and box.text.strip() not in SKIP_TEXTS
    ]
    candidates.sort(key=lambda box: (box.bbox[1], box.bbox[0]))
    groups: list[list[OCRBox]] = []
    for box in candidates:
        joinable = [
            (index, _ocr_line_join_score(group[-1], box))
            for index, group in enumerate(groups)
            if _can_join_ocr_lines(group[-1], box)
        ]
        if joinable:
            group_index, _score = min(joinable, key=lambda item: item[1])
            groups[group_index].append(box)
        else:
            groups.append([box])
    return [_merge_ocr_line_group(group) for group in groups]


def run_ocr(
    slides: list[PPTSlide],
    progress: ProgressCB = None,
    ocr_backend: str = OCR_BACKEND_LOCAL,
    ocr_token: str | None = None,
):
    det_dir = bundled_ocr_model_dir("PP-OCRv5_server_det")
    rec_dir = bundled_ocr_model_dir("PP-OCRv5_server_rec")
    use_remote = ocr_backend == OCR_BACKEND_REMOTE

    def create_ocr():
        if use_remote:
            token = (ocr_token or os.environ.get("PADDLEOCR_ACCESS_TOKEN") or "").strip()
            if len(token) != REMOTE_OCR_TOKEN_LENGTH:
                raise RuntimeError("远端 OCR 需要 40 位访问令牌。")
            if PaddleOCRClient is None:
                raise RuntimeError("当前 PaddleOCR 版本不支持远端调用，请升级 paddleocr 后重试。")
            return PaddleOCRClient(token=token)
        ocr_device = preferred_paddleocr_device()
        return PaddleOCR(
            lang="ch",
            device=ocr_device,
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            text_detection_model_dir=str(det_dir) if det_dir else None,
            text_recognition_model_dir=str(rec_dir) if rec_dir else None,
        )

    if use_remote:
        _log(progress, "使用远端 PaddleOCR 识别")
    else:
        _log(progress, f"使用本地 PaddleOCR 识别（设备：{preferred_paddleocr_device()}）")
    ocr = create_ocr()
    for slide in slides:
        try:
            if use_remote:
                page = _predict_remote_ocr_page(ocr, slide.image_path)
            else:
                page = ocr.predict(str(slide.image_path))[0]
        except Exception as exc:
            _log(progress, f"第 {slide.index} 页 OCR 失败，已跳过该页，可稍后手动新增框：{exc}")
            _log(progress, traceback.format_exc())
            slide.boxes = []
            _page_ready(progress, slide, status="failed")
            try:
                if use_remote and hasattr(ocr, "close"):
                    ocr.close()
                ocr = create_ocr()
            except Exception as reset_exc:
                _log(progress, f"OCR 引擎重置失败，后续页面可能继续失败：{reset_exc}")
            continue
        boxes: list[OCRBox] = []
        for poly, text, score in zip(page["dt_polys"], page["rec_texts"], page["rec_scores"]):
            text = (text or "").strip()
            if not text:
                continue
            if text in SKIP_TEXTS:
                continue
            x, y, w, h = box_to_rect(poly)
            erase_rect = default_expand_rect(x, y, w, h, slide.image_width, slide.image_height)
            boxes.append(
                OCRBox(
                    text=text,
                    score=float(score),
                    bbox=(x, y, w, h),
                erase_rect=erase_rect,
                line_height=h,
                text_regions=(normalize_text_region(poly),),
                )
            )
        slide.boxes = group_ocr_boxes(boxes)
        _page_ready(progress, slide)
        _log(progress, f"第 {slide.index} 页 OCR 完成，共 {len(slide.boxes)} 个框")
    if use_remote and hasattr(ocr, "close"):
        ocr.close()


def prepare_project(
    source_pptx: Path,
    work_dir: Path | None = None,
    progress: ProgressCB = None,
    ocr_backend: str = OCR_BACKEND_LOCAL,
    ocr_token: str | None = None,
    auto_ocr: bool = True,
) -> PPTProject:
    source_pptx = source_pptx.expanduser().resolve()
    if work_dir is None:
        work_dir = BASE / "_gui_workspace" / source_pptx.stem
    if work_dir.exists():
        shutil.rmtree(work_dir)
    images_dir = work_dir / "images"
    masks_dir = work_dir / "masks"
    cleaned_dir = work_dir / "cleaned"
    assets_dir = work_dir / "assets"
    work_dir.mkdir(parents=True, exist_ok=True)

    src, slides = extract_slide_images(source_pptx, images_dir, progress)
    project = PPTProject(
        source_pptx=source_pptx,
        work_dir=work_dir,
        images_dir=images_dir,
        masks_dir=masks_dir,
        cleaned_dir=cleaned_dir,
        assets_dir=assets_dir,
        slides=slides,
        slide_width=src.slide_width,
        slide_height=src.slide_height,
    )
    if not load_project_cache(project, progress=progress) and auto_ocr:
        run_ocr(slides, progress, ocr_backend=ocr_backend, ocr_token=ocr_token)
    for slide in project.slides:
        if slide.visual_assets or not slide.image_path.exists():
            continue
        with Image.open(slide.image_path) as source_image:
            slide.visual_assets = detect_visual_assets(np.asarray(source_image.convert("RGB")), slide)
    return project

def _rect_mask(image_width: int, image_height: int, rect: tuple[int, int, int, int]) -> np.ndarray:
    mask = np.zeros((image_height, image_width), dtype=np.uint8)
    left, top, right, bottom = rect
    cv2.rectangle(mask, (left, top), (right, bottom), 255, -1)
    return mask


def text_region_mask(image_shape: tuple[int, ...], boxes: list[OCRBox]) -> np.ndarray:
    height, width = image_shape[:2]
    mask = np.zeros((height, width), dtype=np.uint8)
    for box in boxes:
        if not box.enabled:
            continue
        if box.text_regions:
            for region in box.text_regions:
                cv2.fillPoly(mask, [np.asarray(region, dtype=np.int32)], 255)
        else:
            x, y, box_width, box_height = box.bbox
            cv2.rectangle(mask, (x, y), (x + box_width, y + box_height), 255, -1)
    return mask


def _visual_foreground_mask(image: np.ndarray) -> np.ndarray:
    hsv = cv2.cvtColor(image, cv2.COLOR_RGB2HSV)
    mask = np.zeros(image.shape[:2], dtype=np.uint8)
    mask[(hsv[:, :, 1] >= 35) & (hsv[:, :, 2] <= 245)] = 255
    return cv2.morphologyEx(mask, cv2.MORPH_CLOSE, np.ones((5, 5), dtype=np.uint8))


def detect_visual_assets(image: np.ndarray, slide: PPTSlide) -> list[VisualAsset]:
    candidates = _visual_foreground_mask(image)
    candidates[text_region_mask(image.shape, slide.boxes) > 0] = 0
    contours, _hierarchy = cv2.findContours(candidates, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
    min_area = max(256, int(image.shape[0] * image.shape[1] * 0.01))
    max_area = int(image.shape[0] * image.shape[1] * 0.80)
    assets = []
    for contour in sorted(contours, key=cv2.contourArea, reverse=True):
        x, y, width, height = cv2.boundingRect(contour)
        if not min_area <= width * height <= max_area:
            continue
        assets.append(VisualAsset(f"slide-{slide.index}-visual-{len(assets) + 1}", (x, y, width, height)))
    return assets


def visual_asset_alpha_mask(image: np.ndarray, asset: VisualAsset, boxes: list[OCRBox]) -> np.ndarray:
    height, width = image.shape[:2]
    x, y, asset_width, asset_height = asset.bbox
    left, top = max(0, x), max(0, y)
    right, bottom = min(width, x + asset_width), min(height, y + asset_height)
    alpha = np.zeros((height, width), dtype=np.uint8)
    if right <= left or bottom <= top:
        return alpha
    foreground = _visual_foreground_mask(image)
    alpha[top:bottom, left:right] = foreground[top:bottom, left:right]
    if not np.any(alpha[top:bottom, left:right]):
        alpha[top:bottom, left:right] = 255
    if asset.layer != "above_text":
        overlap_text = text_region_mask(image.shape, boxes)
        overlap_text[:top, :] = 0
        overlap_text[bottom:, :] = 0
        overlap_text[:, :left] = 0
        overlap_text[:, right:] = 0
        alpha[overlap_text > 0] = 255
    return alpha


def build_asset_text_mask(image: np.ndarray, asset: VisualAsset, boxes: list[OCRBox]) -> np.ndarray:
    height, width = image.shape[:2]
    x, y, asset_width, asset_height = asset.bbox
    left, top = max(0, x), max(0, y)
    right, bottom = min(width, x + asset_width), min(height, y + asset_height)
    combined = np.zeros((height, width), dtype=np.uint8)
    if right <= left or bottom <= top or asset.layer == "above_text":
        return combined
    for box in boxes:
        if not box.enabled:
            continue
        bx, by, bw, bh = box.bbox
        if bx + bw <= left or bx >= right or by + bh <= top or by >= bottom:
            continue
        strokes, _mode, _reason = build_text_stroke_mask(image, box)
        radius = max(1, int(round(max(1, int(box.line_height or bh)) * 0.09)))
        expanded = cv2.dilate(strokes, np.ones((radius * 2 + 1, radius * 2 + 1), np.uint8), iterations=1)
        combined = cv2.bitwise_or(combined, expanded)
    allowed = np.zeros_like(combined)
    allowed[top:bottom, left:right] = 255
    return cv2.bitwise_and(combined, allowed)


def soft_blend_inpaint(original: np.ndarray, repaired: np.ndarray, mask: np.ndarray, feather_px: int = 24) -> np.ndarray:
    if original.shape != repaired.shape or original.shape[:2] != mask.shape[:2]:
        raise ValueError("修复图、原图和蒙版尺寸必须一致。")
    binary = np.zeros(mask.shape[:2], dtype=np.uint8)
    binary[np.asarray(mask) > 0] = 255
    if not np.any(binary):
        return original.copy()
    feather_px = max(1, int(feather_px))
    soft = cv2.GaussianBlur(binary, (feather_px * 2 + 1, feather_px * 2 + 1), sigmaX=max(1.0, feather_px / 3))
    distance = cv2.distanceTransform(binary, cv2.DIST_L2, 3)
    soft[distance >= feather_px] = 255
    weight = soft.astype(np.float32)[:, :, None] / 255.0
    return np.clip(np.rint(repaired.astype(np.float32) * weight + original.astype(np.float32) * (1.0 - weight)), 0, 255).astype(np.uint8)


def finish_page_inpaint(original: np.ndarray, generated: np.ndarray, mask: np.ndarray) -> np.ndarray:
    binary = np.zeros(mask.shape[:2], dtype=np.uint8)
    binary[np.asarray(mask) > 0] = 255
    if not np.any(binary):
        return original.copy()
    candidate = generated.copy()
    count, labels, stats, _centroids = cv2.connectedComponentsWithStats(binary, connectivity=8)
    page_area = max(1, binary.shape[0] * binary.shape[1])
    for component_index in range(1, count):
        if int(stats[component_index, cv2.CC_STAT_AREA]) < page_area * 0.02:
            continue
        component = np.zeros_like(binary)
        component[labels == component_index] = 255
        ring = cv2.subtract(cv2.dilate(component, np.ones((15, 15), np.uint8), iterations=1), component)
        samples = original[ring > 0]
        if samples.size == 0 or float(np.mean(np.std(samples.astype(np.float32), axis=0))) >= 18.0:
            continue
        texture_fill = cv2.inpaint(original, component, 7, cv2.INPAINT_TELEA)
        candidate[component > 0] = texture_fill[component > 0]
    masked_width = int(stats[1:, cv2.CC_STAT_WIDTH].max()) if count > 1 else 1
    masked_height = int(stats[1:, cv2.CC_STAT_HEIGHT].max()) if count > 1 else 1
    feather = min(40, max(15, int(round(min(masked_width, masked_height) * 0.08))))
    return soft_blend_inpaint(original, candidate, binary, feather_px=feather)


def _write_visual_asset(project: PPTProject, slide: PPTSlide, asset: VisualAsset, image: np.ndarray, alpha: np.ndarray):
    assets_dir = project.assets_dir or project.work_dir / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)
    x, y, asset_width, asset_height = asset.bbox
    height, width = image.shape[:2]
    left, top = max(0, x), max(0, y)
    right, bottom = min(width, x + asset_width), min(height, y + asset_height)
    if right <= left or bottom <= top:
        return
    stem = f"{slide.image_path.stem}-{asset.asset_id}"
    asset.image_path = assets_dir / f"{stem}.png"
    asset.mask_path = assets_dir / f"{stem}.mask.png"
    Image.fromarray(np.dstack((image[top:bottom, left:right], alpha[top:bottom, left:right])), "RGBA").save(asset.image_path)
    Image.fromarray(alpha[top:bottom, left:right]).save(asset.mask_path)


def build_text_stroke_mask(image: np.ndarray, box: OCRBox) -> tuple[np.ndarray, str, str | None]:
    """Return a narrow text-pixel mask, falling back to the legacy rectangle when needed."""
    image_height, image_width = image.shape[:2]
    if not box.text_regions:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "缺少 OCR 文字轮廓"

    allowed = _rect_mask(image_width, image_height, box.erase_rect)
    regions = np.zeros_like(allowed)
    for region in box.text_regions:
        cv2.fillPoly(regions, [np.asarray(region, dtype=np.int32)], 255)
    allowed = cv2.bitwise_and(allowed, regions)
    region_pixels = int(np.count_nonzero(allowed))
    if region_pixels == 0:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "OCR 文字轮廓无效"

    gray = cv2.cvtColor(image, cv2.COLOR_RGB2GRAY)
    blurred = cv2.GaussianBlur(gray, (0, 0), sigmaX=5, sigmaY=5)
    contrast = cv2.absdiff(gray, blurred)
    strokes = np.zeros_like(allowed)
    strokes[(contrast >= 18) & (allowed > 0)] = 255
    density = int(np.count_nonzero(strokes)) / region_pixels
    if density < 0.01:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "未找到足够的文字笔画"
    if density > 0.60:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "文字区域线条过密，可能与插图重叠"
    _count, _labels, stats, _centroids = cv2.connectedComponentsWithStats(strokes, connectivity=8)
    large_thin_component = any(
        int(component[cv2.CC_STAT_AREA]) > max(80, int(region_pixels * 0.01))
        and (
            int(component[cv2.CC_STAT_WIDTH]) / max(1, int(component[cv2.CC_STAT_HEIGHT])) <= 0.30
            or int(component[cv2.CC_STAT_HEIGHT]) / max(1, int(component[cv2.CC_STAT_WIDTH])) <= 0.30
        )
        for component in stats[1:]
    )
    if large_thin_component:
        return _rect_mask(image_width, image_height, box.erase_rect), "rectangle_fallback", "检测到大面积连通线条，可能与插图重叠"
    strokes = cv2.dilate(strokes, np.ones((3, 3), dtype=np.uint8), iterations=1)
    return strokes, "text_stroke", None


def build_masks(project: PPTProject, progress: ProgressCB = None):
    project.masks_dir.mkdir(parents=True, exist_ok=True)
    (project.assets_dir or project.work_dir / "assets").mkdir(parents=True, exist_ok=True)
    total_slides = max(1, len(project.slides))
    for done, slide in enumerate(project.slides, start=1):
        with Image.open(slide.image_path) as source_image:
            image = np.asarray(source_image.convert("RGB"))
        mask = np.zeros((slide.image_height, slide.image_width), dtype=np.uint8)
        for box in slide.boxes:
            if not box.enabled:
                continue
            box_mask, box.mask_mode, box.mask_reason = build_text_stroke_mask(image, box)
            mask = cv2.bitwise_or(mask, box_mask)
            if box.mask_reason:
                _log(progress, f"第 {slide.index} 页文字框已回退矩形擦除：{box.mask_reason}")
        for asset in slide.visual_assets:
            if not asset.enabled:
                continue
            alpha = visual_asset_alpha_mask(image, asset, slide.boxes)
            if not np.any(alpha):
                asset.enabled = False
                asset.status = "disabled_empty"
                continue
            _write_visual_asset(project, slide, asset, image, alpha)
            mask = cv2.bitwise_or(mask, alpha)
        if slide.remove_watermark and slide.watermark_rect:
            left, top, right, bottom = slide.watermark_rect
            cv2.rectangle(mask, (left, top), (right, bottom), 255, -1)
        mask_path = project.masks_dir / slide.image_name
        Image.fromarray(mask).save(mask_path)
        if not mask_path.exists():
            raise RuntimeError(f"擦除蒙版写入失败：{mask_path}")
        _log(progress, f"第 {slide.index} 页擦除蒙版已生成")
        _progress(progress, int(done * 100 / total_slides), f"生成擦除蒙版：{done}/{total_slides}")


def os_environ_with_pythonpath():
    env = dict(os.environ)
    pythonpath_parts = [str(Path(IOPAINT_DEPS_DIR).resolve()), str(Path(DEPS_DIR).resolve())]
    if env.get("PYTHONPATH"):
        pythonpath_parts.append(env["PYTHONPATH"])
    env["PYTHONPATH"] = ";".join(pythonpath_parts)
    bundled_model_root = bundled_iopaint_model_root()
    if bundled_model_root:
        env["XDG_CACHE_HOME"] = str(bundled_model_root)
        env["U2NET_HOME"] = str(bundled_model_root)
    return env


def repair_visual_assets(project: PPTProject, model, inpaint_request, progress: ProgressCB = None) -> None:
    for slide in project.slides:
        if not slide.image_path.is_file():
            continue
        with Image.open(slide.image_path) as source_image:
            page = np.asarray(source_image.convert("RGB"))
        page_height, page_width = page.shape[:2]
        for asset in slide.visual_assets:
            if not asset.enabled or not asset.image_path or not asset.image_path.is_file():
                continue
            text_mask = build_asset_text_mask(page, asset, slide.boxes)
            if not np.any(text_mask):
                asset.mask_version = max(asset.mask_version, 2)
                continue
            x, y, asset_width, asset_height = asset.bbox
            padding = max(24, int(round(max(asset_width, asset_height) * 0.08)))
            crop_left, crop_top = max(0, x - padding), max(0, y - padding)
            crop_right = min(page_width, x + asset_width + padding)
            crop_bottom = min(page_height, y + asset_height + padding)
            crop = page[crop_top:crop_bottom, crop_left:crop_right].copy()
            local_text_mask = text_mask[crop_top:crop_bottom, crop_left:crop_right].copy()
            model_mask = local_text_mask.copy()
            available = model_mask == 0
            fill_color = np.median(crop[available], axis=0).astype(np.uint8) if np.any(available) else np.zeros(3, np.uint8)
            for neighbor in slide.visual_assets:
                if neighbor is asset or not neighbor.enabled:
                    continue
                nx, ny, nw, nh = neighbor.bbox
                left, top = max(crop_left, nx), max(crop_top, ny)
                right, bottom = min(crop_right, nx + nw), min(crop_bottom, ny + nh)
                if right <= left or bottom <= top:
                    continue
                lx0, ly0, lx1, ly1 = left - crop_left, top - crop_top, right - crop_left, bottom - crop_top
                neighbor_context = np.zeros_like(model_mask)
                neighbor_context[ly0:ly1, lx0:lx1] = 255
                ax0, ay0 = max(0, x - crop_left), max(0, y - crop_top)
                ax1 = min(neighbor_context.shape[1], x + asset_width - crop_left)
                ay1 = min(neighbor_context.shape[0], y + asset_height - crop_top)
                neighbor_context[ay0:ay1, ax0:ax1] = 0
                crop[neighbor_context > 0] = fill_color
                model_mask[neighbor_context > 0] = 255
            try:
                generated = cv2.cvtColor(np.asarray(model(crop, model_mask, inpaint_request)), cv2.COLOR_BGR2RGB)
                heights = [int(box.line_height or box.bbox[3]) for box in slide.boxes if box.enabled]
                feather = min(40, max(15, int(round(median(heights))) if heights else 15))
                repaired_crop = soft_blend_inpaint(crop, generated, local_text_mask, feather_px=feather)
                with Image.open(asset.image_path) as stored_asset:
                    rgba = np.asarray(stored_asset.convert("RGBA")).copy()
                asset_left, asset_top = max(0, x) - crop_left, max(0, y) - crop_top
                repaired_asset = repaired_crop[asset_top:asset_top + rgba.shape[0], asset_left:asset_left + rgba.shape[1]]
                if repaired_asset.shape[:2] != rgba.shape[:2]:
                    raise RuntimeError("局部修复结果与图片资产尺寸不一致。")
                rgba[:, :, :3] = repaired_asset
                Image.fromarray(rgba, "RGBA").save(asset.image_path)
                asset.mask_version = 2
                _log(progress, f"第 {slide.index} 页图片内文字已局部修复：{asset.asset_id}")
            except Exception as exc:
                asset.status = "repair_warning"
                _log(progress, f"第 {slide.index} 页图片内文字修复失败，已保留原图：{asset.asset_id} ({exc})")


def run_iopaint(
    images_dir: Path,
    masks_dir: Path,
    cleaned_dir: Path,
    progress: ProgressCB = None,
    project: PPTProject | None = None,
):
    if cleaned_dir.exists():
        shutil.rmtree(cleaned_dir)
    cleaned_dir.mkdir(parents=True, exist_ok=True)
    _log(progress, "开始用 IOPaint 擦除底图文字")

    env = os_environ_with_pythonpath()
    for key in ("XDG_CACHE_HOME", "U2NET_HOME", "PYTHONPATH"):
        if key in env:
            os.environ[key] = env[key]

    LaMa = load_iopaint_lama_class()
    from iopaint.helper import pil_to_bytes
    from iopaint.model.utils import torch_gc
    from iopaint.schema import Device, InpaintRequest

    image_paths = {
        path.stem: path
        for path in images_dir.iterdir()
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg"}
    }
    mask_paths = {
        path.stem: path
        for path in masks_dir.iterdir()
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg"}
    }
    if not image_paths:
        raise RuntimeError(f"没有找到待处理图片：{images_dir}")
    if not mask_paths:
        raise RuntimeError(f"没有找到擦除蒙版：{masks_dir}")

    total_images = len(image_paths)
    _progress(progress, 0, f"IOPaint 擦除处理中：0/{total_images}")

    if not LaMa.is_downloaded():
        _log(progress, "本机未找到 IOPaint lama 模型，开始准备模型")
        LaMa.download()

    first_mask = next(iter(mask_paths.values()))

    def process_with_device(device) -> None:
        _log(progress, f"IOPaint 使用设备：{device}")
        model = LaMa(device=device)
        inpaint_request = InpaintRequest()
        for done, (stem, image_path) in enumerate(sorted(image_paths.items()), start=1):
            mask_path = mask_paths.get(stem, first_mask)
            with Image.open(image_path) as source_image:
                infos = source_image.info
                image = np.array(source_image.convert("RGB"))
            with Image.open(mask_path) as source_mask:
                mask = np.array(source_mask.convert("L"))

            if mask.shape[:2] != image.shape[:2]:
                mask = cv2.resize(
                    mask,
                    (image.shape[1], image.shape[0]),
                    interpolation=cv2.INTER_NEAREST,
                )
            mask[mask >= 127] = 255
            mask[mask < 127] = 0

            result = model(image, mask, inpaint_request)
            result = cv2.cvtColor(result, cv2.COLOR_BGR2RGB)
            result = finish_page_inpaint(image, result, mask)
            output_path = cleaned_dir / f"{stem}.png"
            result_image = Image.fromarray(result)
            output_path.write_bytes(pil_to_bytes(result_image, "png", 100, infos))
            compressed_path = save_compressed_cleaned_image(result_image, output_path)
            torch_gc()
            _log(progress, f"第 {done} 页 IOPaint 擦除后已生成：{compressed_path.name}")
            _progress(progress, int(done * 100 / total_images), f"IOPaint 擦除处理中：{done}/{total_images}")
        if project is not None:
            repair_visual_assets(project, model, inpaint_request, progress)

    device = preferred_iopaint_device(Device)
    try:
        process_with_device(device)
    except Exception as exc:
        if device == Device.cpu:
            raise
        _log(progress, f"IOPaint CUDA 处理失败，已回退 CPU：{exc}")
        if cleaned_dir.exists():
            shutil.rmtree(cleaned_dir)
        cleaned_dir.mkdir(parents=True, exist_ok=True)
        process_with_device(Device.cpu)

    missing = [name for name in image_paths if not (cleaned_dir / f"{name}.png").exists()]
    if missing:
        raise RuntimeError(f"IOPaint 未生成以下页面：{', '.join(sorted(missing))}")

    _log(progress, "IOPaint 擦除完成")


def save_compressed_cleaned_image(image: Image.Image, source_path: Path, quality: int = EXPORT_IMAGE_JPEG_QUALITY) -> Path:
    output_path = source_path.with_suffix(".jpg")
    save_kwargs = {
        "format": "JPEG",
        "quality": quality,
        "optimize": True,
        "progressive": True,
    }
    icc_profile = image.info.get("icc_profile")
    if icc_profile:
        save_kwargs["icc_profile"] = icc_profile
    rgb_image = image if image.mode == "RGB" else image.convert("RGB")
    rgb_image.save(output_path, **save_kwargs)
    return output_path


def cleaned_source_image_paths(cleaned_dir: Path) -> list[Path]:
    image_paths = [
        path
        for path in sorted(cleaned_dir.iterdir())
        if path.is_file() and path.suffix.lower() in {".png", ".jpg", ".jpeg"}
    ]
    png_stems = {path.stem for path in image_paths if path.suffix.lower() == ".png"}
    return [
        path
        for path in image_paths
        if not (path.suffix.lower() in {".jpg", ".jpeg"} and path.stem in png_stems)
    ]


def upscale_cleaned_images(cleaned_dir: Path, scale: float = 2.0, progress: ProgressCB = None) -> int:
    if scale <= 1:
        _log(progress, "RealESRGAN 清晰化已跳过")
        return 0

    image_paths = cleaned_source_image_paths(cleaned_dir)
    if not image_paths:
        _log(progress, f"没有找到待清晰化图片：{cleaned_dir}")
        return 0

    env = os_environ_with_pythonpath()
    for key in ("XDG_CACHE_HOME", "U2NET_HOME", "PYTHONPATH"):
        if key in env:
            os.environ[key] = env[key]

    from iopaint.model.utils import torch_gc
    from iopaint.plugins import RealESRGANUpscaler
    from iopaint.schema import Device, RunPluginRequest

    _log(progress, "开始用 RealESRGAN 提升导出底图清晰度")
    _progress(progress, 0, f"RealESRGAN 清晰化处理中：0/{len(image_paths)}")
    device = preferred_iopaint_device(Device)

    def create_upscaler(target_device):
        _log(progress, f"RealESRGAN 使用设备：{target_device}")
        return RealESRGANUpscaler("realesr-general-x4v3", target_device, no_half=(target_device == Device.cpu))

    try:
        upscaler = create_upscaler(device)
    except Exception as exc:
        if device == Device.cpu:
            raise
        _log(progress, f"RealESRGAN CUDA 处理失败，已回退 CPU：{exc}")
        device = Device.cpu
        upscaler = create_upscaler(device)

    for done, image_path in enumerate(image_paths, start=1):
        with Image.open(image_path) as source_image:
            infos = source_image.info
            rgb_image = np.array(source_image.convert("RGB"))

        try:
            bgr_result = upscaler.gen_image(
                rgb_image,
                RunPluginRequest(name=RealESRGANUpscaler.name, image="", scale=scale),
            )
        except Exception as exc:
            if device == Device.cpu:
                raise
            _log(progress, f"RealESRGAN CUDA 处理失败，已回退 CPU：{exc}")
            device = Device.cpu
            upscaler = create_upscaler(device)
            bgr_result = upscaler.gen_image(
                rgb_image,
                RunPluginRequest(name=RealESRGANUpscaler.name, image="", scale=scale),
            )
        rgb_result = cv2.cvtColor(bgr_result, cv2.COLOR_BGR2RGB)
        output_format = "JPEG" if image_path.suffix.lower() in {".jpg", ".jpeg"} else "PNG"
        result_image = Image.fromarray(rgb_result)
        if output_format == "JPEG":
            compressed_path = save_compressed_cleaned_image(result_image, image_path)
        else:
            result_image.save(image_path, format=output_format, **infos)
            compressed_path = save_compressed_cleaned_image(result_image, image_path)
        torch_gc()
        _log(progress, f"第 {done} 页已提升清晰度：{compressed_path.name}")
        _progress(progress, int(done * 100 / len(image_paths)), f"RealESRGAN 清晰化处理中：{done}/{len(image_paths)}")

    _log(progress, "RealESRGAN 清晰化完成")
    return len(image_paths)


def cleaned_image_path(cleaned_dir: Path, name: str):
    def preferred_existing(path: Path) -> Path | None:
        for candidate in [path.with_suffix(".jpg"), path.with_suffix(".jpeg"), path]:
            if candidate.exists():
                return candidate
        return None

    nested = cleaned_dir / name / name.replace("_clean", "")
    preferred = preferred_existing(nested)
    if preferred:
        return preferred
    direct = cleaned_dir / name
    preferred = preferred_existing(direct)
    if preferred:
        return preferred
    fallback = cleaned_dir / name.replace("_clean", "")
    preferred = preferred_existing(fallback)
    if preferred:
        return preferred
    raise FileNotFoundError(name)


def sample_text_color(image: Image.Image, rect):
    region = image.crop(rect).convert("RGB")
    arr = np.asarray(region)
    if arr.size == 0:
        return RGBColor(0, 0, 0)
    flat = arr.reshape(-1, 3)
    brightness = flat.sum(axis=1)
    sample = flat[np.argsort(brightness)[: max(1, len(flat) // 8)]]
    rgb = sample.mean(axis=0).astype(int)
    return RGBColor(int(rgb[0]), int(rgb[1]), int(rgb[2]))


def should_rotate_text(box: OCRBox) -> bool:
    if box.rotation in {90, 270, -90}:
        return True
    if "\n" in box.text:
        return False
    x, y, w, h = box.bbox
    return h > w * 1.45 and len(box.text.strip()) > 1


def add_textbox(slide, color_image, box: OCRBox, x_scale: float, y_scale: float):
    text = box.text
    if text in SKIP_TEXTS or (box.score < 0.45 and not box.manual and not box.edited) or not box.enabled:
        return False

    x, y, w, h = box.bbox
    left = int(x * x_scale)
    top = int(y * y_scale)
    width = int(w * x_scale)
    height = int(h * y_scale)
    rotate_text = should_rotate_text(box)
    font_axis = height
    if rotate_text:
        center_x = left + width / 2
        center_y = top + height / 2
        shape_left = int(center_x - height / 2)
        shape_top = int(center_y - width / 2)
        shape = slide.shapes.add_textbox(shape_left, shape_top, height, width)
        shape.rotation = box.rotation if box.rotation in {90, 270} else 270
        font_axis = width
    else:
        shape = slide.shapes.add_textbox(left, top, width, height)
        font_axis = box.line_height * y_scale if box.line_height is not None else height
    shape.fill.background()
    shape.line.fill.background()

    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = True
    run.font.size = Pt(max(8, font_axis * 0.72 / 12700))
    run.font.color.rgb = sample_text_color(color_image, box.erase_rect)
    return True


def rebuild_ppt(project: PPTProject, output_pptx: Path, progress: ProgressCB = None):
    out = Presentation()
    out.slide_width = project.slide_width
    out.slide_height = project.slide_height
    while out.slides:
        rel_id = out.slides._sldIdLst[0].rId
        out.part.drop_rel(rel_id)
        del out.slides._sldIdLst[0]

    blank = out.slide_layouts[6]
    total_slides = max(1, len(project.slides))
    for done, slide_data in enumerate(project.slides, start=1):
        _progress(progress, int((done - 1) * 100 / total_slides), f"重建可编辑文本框：{done}/{total_slides}")
        cleaned_path = cleaned_image_path(project.cleaned_dir, slide_data.image_name)
        color_image = Image.open(slide_data.image_path).convert("RGB")
        dst = out.slides.add_slide(blank)
        dst.shapes.add_picture(
            str(cleaned_path),
            0,
            0,
            width=out.slide_width,
            height=out.slide_height,
        )
        x_scale = out.slide_width / slide_data.image_width
        y_scale = out.slide_height / slide_data.image_height

        def add_visual_asset(asset: VisualAsset):
            if not asset.enabled or not asset.image_path or not asset.image_path.exists():
                return
            x, y, width, height = asset.bbox
            dst.shapes.add_picture(
                str(asset.image_path), int(x * x_scale), int(y * y_scale),
                width=int(width * x_scale), height=int(height * y_scale),
            )

        for asset in slide_data.visual_assets:
            if asset.layer != "above_text":
                add_visual_asset(asset)
        count = 0
        for box in slide_data.boxes:
            if add_textbox(dst, color_image, box, x_scale, y_scale):
                count += 1
        for asset in slide_data.visual_assets:
            if asset.layer == "above_text":
                add_visual_asset(asset)
        _log(progress, f"第 {slide_data.index} 页已重建 {count} 个可编辑文本框")
        _progress(progress, int(done * 100 / total_slides), f"重建可编辑文本框：{done}/{total_slides}")
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    _progress(progress, 100, "正在保存 PPT 文件...")
    out.save(str(output_pptx))
    _log(progress, f"已导出：{output_pptx}")


def export_editable_ppt(
    project: PPTProject,
    output_pptx: Path,
    progress: ProgressCB = None,
    enhance_images: bool = True,
):
    build_masks(project, progress)
    run_iopaint(project.images_dir, project.masks_dir, project.cleaned_dir, progress, project=project)
    if enhance_images:
        upscale_cleaned_images(project.cleaned_dir, progress=progress)
    else:
        _log(progress, "已跳过 RealESRGAN 底图清晰化")
    rebuild_ppt(project, output_pptx, progress)
